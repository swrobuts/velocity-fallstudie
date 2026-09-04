#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Erzeugt das Foliendeck „Datenbankentwurf am Fallbeispiel VeloCity".

Aufruf:
    bash tools/render_diagrams.sh      # Diagramme zuerst
    python3 slides/build_deck.py

Als Vorlage dient ein bestehendes BINT-Deck: daraus kommen Thema,
Layouts und Fußmasken. Alle Folien werden neu gebaut.

Didaktisches Gerüst
-------------------
Ein durchgehender Fall trägt das Deck: **Anna fährt 61 Minuten mit einem
E-Bike**. Jedes Kapitel öffnet mit einer Leitfrage, die dieser Fall
aufwirft, und schließt mit ihrer Beantwortung. Auf den Inhaltsfolien
stellt ein roter Streifen am Fuß den Bezug zu Annas Fahrt her.

Wo es trägt, kommt erst der naive Versuch, der scheitert, und dann die
Lösung — nicht die Lösung allein.
"""
from __future__ import annotations

import pathlib
import sys

from pptx import Presentation
from pptx.util import Pt

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent))
from thws import (  # noqa: E402
    BLAU, BREITE, FLUCHT_L, GRUEN_D, ORANGE, ROT_A, TEXT_SEK,
    ZONE_OBEN, ZONE_UNTEN,
    ampel_matrix, code_kacheln, diagramm, faden, kachelreihe, kopf, leitfrage,
    notizen, prozesskette, regel_streifen, sandband, sandkarte,
    schichtenstapel, tabelle, vorher_nachher,
)

WURZEL  = pathlib.Path(__file__).resolve().parent.parent

# Frueher diente ein konkretes BINT-Vorlesungsdeck aus dem Werkzeug
# thws-deck-batch als Vorlage (nur fuer Layouts und Fussmasken, keine
# Inhalte). Dieses Werkzeug ist abgeloest, der Ordner existiert nicht
# mehr - ein fest verdrahteter absoluter Pfad auf ein fremdes Deck war
# der Grund, warum das Bauen kaputtging. Die Nachfolge ist der Skill
# thws-slides: sein Master traegt dieselben drei benoetigten Layouts
# (Frontpage_Digital, Chapter, Slide) und enthaelt bereits null Folien,
# sodass leere_praesentation() nichts zu entfernen hat.
VORLAGE = pathlib.Path(
    "/Users/robert/.claude/skills/thws-slides/assets/template.pptx"
)
_VORLAGE_ALT = pathlib.Path(
    "/Users/robert/Library/CloudStorage/OneDrive-Persönlich/Vorlesungen/"
    "thws-deck-batch/decks/BINT_E4_Datenmodellierung_WS2627_v3.pptx"
)
if not VORLAGE.exists():
    if _VORLAGE_ALT.exists():
        VORLAGE = _VORLAGE_ALT
    else:
        raise SystemExit(
            "Vorlage fuer das Foliendeck fehlt. Gesucht wurden:\n"
            f"  1) {VORLAGE}  (Master des Skills thws-slides)\n"
            f"  2) {_VORLAGE_ALT}  (alte, abgeloeste Vorlage)\n"
            "Keine der beiden Dateien existiert. Ohne Vorlage fehlen die "
            "Layouts Frontpage_Digital/Chapter/Slide, das Deck kann nicht "
            "gebaut werden - Pfad korrigieren oder Vorlage bereitstellen."
        )

ZIEL   = WURZEL / "slides" / "velocity-datenbankentwurf.pptx"
ASSETS = WURZEL / "slides" / "assets"

Q = ("Fallstudie VeloCity, Schema velocity auf supabase.butscher.cloud. "
     "Quellen und Diagramme im Repository unter doku/datenmodell/.")


def bild(name: str) -> str:
    pfad = ASSETS / f"{name}.png"
    if not pfad.exists():
        raise SystemExit(f"Diagramm fehlt: {pfad}\nZuerst: bash tools/render_diagrams.sh")
    return str(pfad)


def leere_praesentation() -> Presentation:
    prs = Presentation(str(VORLAGE))
    liste = prs.slides._sldIdLst
    for sld in list(liste):
        prs.part.drop_rel(sld.rId)
        liste.remove(sld)
    return prs


def lay(prs, name):
    return next(l for l in prs.slide_layouts if l.name == name)


def kapitel(prs, nummer, titel, frage, notiz):
    s = prs.slides.add_slide(lay(prs, "Chapter"))
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.text = f"{nummer} · {titel}"
    leitfrage(s, frage)
    notizen(s, notiz)
    return s


def folie(prs, kicker, titel, intro=None, quelle=Q):
    s = prs.slides.add_slide(lay(prs, "Slide"))
    s._intro_unten = kopf(s, kicker, titel, quelle=quelle, intro=intro)
    return s


def unter_intro(s, abstand=16):
    """Erste freie Zeile unterhalb der Einleitung."""
    return getattr(s, '_intro_unten', 110) + abstand


def baue() -> Presentation:
    prs = leere_praesentation()

    # ═══════════════════════════════════════════════════════ Titel
    s = prs.slides.add_slide(lay(prs, "Frontpage_Digital"))
    for ph in s.placeholders:
        i = ph.placeholder_format.idx
        if i == 10:
            ph.text_frame.text = "Datenbankentwurf am Fallbeispiel VeloCity"
        elif i == 11:
            ph.text_frame.text = ("Eine Fahrt, 61 Minuten, 4,96 Euro — und alles, was eine "
                                  "Datenbank dafür wissen und garantieren muss")
    notizen(s, "Wir entwerfen in dieser Einheit eine Datenhaltung von Grund auf. Kein "
               "vorhandenes Schema, kein Umbau: ein Unternehmen kommt mit einem Auftrag. "
               "Damit das nicht abstrakt bleibt, hängen wir alles an EINEN Fall — an eine "
               "einzige Fahrt. Die begleitet uns bis zur letzten Folie.")

    # ═══════════════════════════════════════════ Der Fall als roter Faden
    s = folie(prs, "Der Fall", "Anna fährt 61 Minuten. Das ist alles, was wir wissen.",
              "Aus diesem einen Satz entsteht das gesamte Datenmodell. Jede Frage, die "
              "diese Fahrt aufwirft, beantwortet ein Kapitel dieser Einheit — und am Ende "
              "steht eine Datenbank, die Annas Rechnung selbst berechnet und ihre Regeln "
              "erzwingt.")
    diagramm(s, bild("faden-annas-fahrt"), y=unter_intro(s), hoehe=190)
    sandband(s, "Prüfen Sie im Verlauf mit: Woher weiß die Datenbank am Ende, dass Anna "
                "genau 4,96 Euro zahlt — und nicht 6,20 oder 4,00?", y=396)
    notizen(s, "Die Zahlen sind echt: 61 Minuten, E-Bike, Studententarif mit Rabatt. Am "
               "Ende der Einheit können wir jeden Cent begründen. Die Zwischenbeträge 6,20 "
               "und 4,00 sind die beiden naheliegenden Fehler — der eine ignoriert den "
               "Tarif, der andere dreht die Reihenfolge um.")

    s = folie(prs, "Orientierung", "Zwölf Kapitel, zwölf Fragen an Annas Fahrt",
              "Die Reihenfolge ist zwingend. Anforderungen vor Modell, Modell vor "
              "Relationen, Relationen vor DDL. Wer sie umdreht, bindet sich an eine "
              "Umsetzung, bevor er weiß, was umzusetzen ist.")
    tabelle(s, ["Kapitel", "Die Frage, die Annas Fahrt aufwirft"],
            [["1 Fallstudie",     "Was muss die Datenbank über diese Fahrt überhaupt wissen?"],
             ["2 Konzeptionell",  "Welche Dinge sind das — und wie hängen sie zusammen?"],
             ["3 Normalisierung", "Warum reicht dafür nicht eine einzige Tabelle?"],
             ["4 Logisch",        "Wie wird aus dem Bild ein Satz von Relationen?"],
             ["5 Physisch",       "Wie erzwingt die Datenbank, dass Annas Rechnung stimmt?"],
             ["6 Implementierung","Was kostet die Fahrt genau — und warum?"],
             ["7 Sicherheit",     "Wer außer Anna darf ihre Daten sehen?"],
             ["8 Anbindung",      "Wie kommt das alles in die Web-Anwendung?"],
             ["9 Warenwirtschaft", "Wer arbeitet mit der Fahrt, wenn sie bezahlt ist?"],
             ["10 Analytik",      "Welche Frage stellt die Leitung — und welche Zeile antwortet?"],
             ["11 Oberflächen",   "Wie kommt Annas Rad auf zwei verschiedene Bildschirme?"],
             ["12 Abschluss",     "Was bleibt — und was ist entworfen, aber nicht gebaut?"]],
            y=180, spalten_b=[220, 683.5], zeilen_h=23)
    notizen(s, "Diese Tabelle ist die Landkarte. Ich komme nach jedem Kapitel kurz darauf "
               "zurück und hake ab.")

    # ═══════════════════════════════════════════════════ 1 Fallstudie
    kapitel(prs, 1, "Die Fallstudie",
            "Was muss die Datenbank über Annas Fahrt überhaupt wissen?",
            "Bevor irgendetwas modelliert wird, muss das Geschäft verstanden sein. "
            "Anforderungen werden in Substantiven und Regeln notiert, nicht in Tabellen.")

    s = folie(prs, "1 · Die Fallstudie", "VeloCity vermietet Räder minutengenau",
              "Ein Bike-Sharing-System in Würzburg. Kunden finden freie Räder auf einer "
              "Karte, entleihen sie, fahren und stellen sie wieder ab. Vielfahrer schließen "
              "einen Tarif ab, der Freiminuten und Rabatt mitbringt.")
    kachelreihe(s, [
        ("Drei Fahrzeugklassen",
         ["City-Bike, Stadtrad ohne Motor",
          "E-Bike Sport, Pedelec bis 25 km/h",
          "E-Cargo Loader, bis 75 kg Zuladung"]),
        ("Zwei Abstellarten",
         ["An einer festen Station, kostenfrei",
          "Frei im Stadtgebiet, ohne Zuschlag",
          "Beides muss das Modell abbilden"]),
        ("Vier Tarife",
         ["Basis ohne Freiminuten",
          "Student und ÖPNV mit Kontingent",
          "Premium mit Rabatt und Monatspreis"]),
    ], y=182, hoehe=150)
    faden(s, "Anna fährt ein E-Bike und hat den Studententarif — 300 Freiminuten im Monat.")
    notizen(s, "Wichtig ist die zweite Kachel: Free-Floating heißt, dass ein Rad auch ohne "
               "Station eine Position haben muss. Das entscheidet später über die Struktur "
               "der Positionsdaten.")

    s = folie(prs, "1 · Die Fallstudie", "Zehn Geschäftsregeln — jede muss erzwungen werden",
              "Diese Regeln sind die fachliche Substanz. Entscheidend ist nicht, sie zu "
              "notieren, sondern dass jede einzelne am Ende erzwungen wird — von der "
              "Datenbank, nicht von der Anwendung.")
    regel_streifen(s, [
        ("GR1", "Ein Rad ist höchstens einmal gleichzeitig ausgeliehen", "partieller UNIQUE-Index"),
        ("GR2", "Ein Kunde hat höchstens vier aktive Ausleihen", "Prüfung in der Funktion"),
        ("GR3", "Ein Kunde hat zu einem Zeitpunkt einen Tarif", "EXCLUDE-Constraint"),
        ("GR4", "Preise und Konditionen überlappen sich nie", "EXCLUDE-Constraint"),
        ("GR5", "Bepreist wird mit dem Preis zur Startzeit", "in fn_ausleihe_abrechnen"),
        ("GR6", "Angefangene Minuten werden aufgerundet", "GENERATED-Spalte"),
    ], y=174, hoehe=42, luecke=5)
    faden(s, "GR5 entscheidet über Annas Rechnung: es gilt der Preis um 10:00 Uhr, nicht der von heute.")
    notizen(s, "Sechs von zehn. Rechts steht schon, wo die Regel später landet — das ist "
               "der rote Faden durch die Umsetzung. GR7 bis GR15 folgen im physischen Entwurf.")

    # ═══════════════════════════════════════════ 2 Konzeptioneller Entwurf
    kapitel(prs, 2, "Konzeptioneller Entwurf",
            "Welche Dinge sind das — und wie hängen sie zusammen?",
            "Entitäten, Beziehungen, Kardinalitäten. Noch keine Datentypen, keine "
            "Schlüsselstrategie, keine Datenbank.")

    s = folie(prs, "2 · Konzeptioneller Entwurf", "Substantive werden Entitäten, Verben werden Beziehungen",
              "Eine Heuristik zum Einstieg, kein Automatismus — die eigentliche Arbeit steckt in "
              "den Kardinalitäten.")
    oben = unter_intro(s)
    diagramm(s, bild("erm-konzeptionell"), y=oben, hoehe=ZONE_UNTEN - 12 - oben)
    notizen(s, "Bewusst ohne Attribute. Auf dieser Stufe interessiert nur, welche Dinge es "
               "gibt und wie viele davon auf jeder Seite stehen. „Ein Kunde hat höchstens "
               "einen gültigen Tarif“ ist eine Entscheidung des Fachbereichs, keine "
               "technische Festlegung.")

    s = folie(prs, "2 · Konzeptioneller Entwurf", "Kardinalitäten lesen und prüfen",
              "Die Notation ist schnell gelernt. Die Frage dahinter ist die schwierige: "
              "Stimmt die Aussage über die Wirklichkeit?")
    tabelle(s, ["Zeichen", "Bedeutung", "Beispiel im Modell", "Ist das wirklich so?"],
            [["||--o{", "eins zu null oder viele", "KUNDE tätigt AUSLEIHE", "Neukunde ohne Fahrt: ja"],
             ["||--|{", "eins zu ein oder viele",  "RECHNUNG hat POSITION",  "Rechnung ohne Posten wäre sinnlos"],
             ["||--||", "eins zu eins",            "FAHRRAD hat POSITION",   "genau eine Position je Rad"],
             ["||--o|", "eins zu null oder eins",  "AUSLEIHE zu POSITION",   "noch nicht fakturiert: ja"]],
            y=192, spalten_b=[110, 230, 260, 303.5], zeilen_h=34)
    faden(s, "Annas Fahrt ist eine AUSLEIHE — sie verbindet Kunde, Rad, Station und Tarif in einem Vorgang.")
    notizen(s, "Die vierte Spalte ist die wichtige. Jede Kardinalität muss man gegen die "
               "Wirklichkeit prüfen, nicht aus dem Bauch setzen.")

    # Relationale Modelle je Bereich.
    #
    # Diese Folien tragen bewusst nur eine einzeilige Einleitung: das
    # Diagramm ist die Aussage, und es braucht die volle Hoehe, um im
    # Hoersaal lesbar zu bleiben.
    for kicker, titel, intro, bilddatei, fadentext, notiz in [
        ("2 · Bereich A", "Geschäftspartner: Adresse als eigene Entität",
         "Kunde, Station, Lieferant und Lager brauchen dieselbe Adressstruktur — vier Mal "
         "dieselben fünf Spalten wären Redundanz auf Schemaebene.",
         "rel-a-geschaeftspartner",
         "Anna ist ein Satz in kunde — mit Kundennummer, aber ohne Passwort.",
         "Auffällig: keine Passwortspalte. Die Anmeldung liegt vollständig bei Supabase "
         "Auth, verbunden über auth_uid. Und hausnummer ist NOT NULL mit leerem "
         "Vorgabewert — sonst greift der fachliche Schlüssel nicht, weil zwei NULL-Werte "
         "in einem UNIQUE-Index als verschieden gelten."),
        ("2 · Bereich B", "Netz und Flotte: Stammdaten und Bewegungsdaten getrennt",
         "Ein Rad ändert seine Stammdaten fast nie und seine Position ständig — die "
         "1:1-Trennung hält die Stammdatentabelle ruhig.",
         "rel-b-netz-und-flotte",
         "Annas E-Bike steht vor der Fahrt an einer Station, danach vielleicht frei im Stadtgebiet.",
         "station_id IS NULL bedeutet eindeutig: frei abgestellt. Zwischen Typ und Rad "
         "steht das Modell — das wirkt zunächst überflüssig, ist aber die Brücke zur "
         "Warenwirtschaft: Ersatzteile hängen am Modell, nicht am Einzelrad."),
        ("2 · Bereich C", "Tarif und Preis: alles ist zeitabhängig",
         "Ein Preis gilt für einen Zeitraum, nicht für immer — der EXCLUDE-Constraint "
         "verhindert Überschneidungen in der Datenbank, nicht in der Anwendung.",
         "rel-c-tarif-und-preis",
         "Annas Studententarif bringt 300 Freiminuten im Monat und wird beim Start ihrer Fahrt fixiert.",
         "freiminuten_periode ist der Kern: statt eines Zählers, der heruntergezählt wird, "
         "stehen Kontingent und Verbrauch je Monat nebeneinander. Bestand und Bewegung "
         "bleiben unterscheidbar, jeder Stand ist rekonstruierbar."),
        ("2 · Bereich D", "Nutzung: die Ausleihe und ihre Abrechnung",
         "Ein einzelner Kostenbetrag verrät nicht, wie er zustande kam — jede Zeile der "
         "Abrechnung bleibt sichtbar und trägt ihren Beleg.",
         "rel-d-nutzung",
         "Annas Fahrt erzeugt fünf Entgeltpositionen — jede einzeln nachvollziehbar.",
         "Der Fremdschlüssel von entgeltposition auf nutzungspreis ist der wichtigste Pfeil "
         "im ganzen Modell: er dokumentiert, welcher Preissatz tatsächlich angewandt wurde. "
         "Ohne ihn wäre eine Altrechnung nach einer Preisänderung nicht mehr erklärbar."),
        ("2 · Bereich E", "Abrechnung: Beleg, Positionen, Zahlung",
         "Das Kopf-Positionen-Muster, ergänzt um die Zahlung — gespeichert wird nur das "
         "Token des Dienstleisters, weder IBAN noch Kartennummer.",
         "rel-e-abrechnung",
         "Annas Fahrt erscheint am Monatsende als eine Zeile auf einer Rechnung.",
         "Was nicht gespeichert wird, kann nicht abfließen. Das ist keine Bequemlichkeit, "
         "sondern die wirksamste Schutzmaßnahme überhaupt."),
        ("2 · Bereich F", "Redaktionsinhalte: warum vier Tabellen und nicht eine",
         "FAQ, Nutzungsschritte, Kennzahlen und Höhenmarken könnten in einer "
         "generischen Schlüssel-Wert-Tabelle liegen — das wäre flexibler und "
         "schlechter.",
         "rel-f-inhalte",
         "Der Preis, den Anna auf der Tarifkarte sieht, wird aus nutzungspreis gerechnet — nicht getippt.",
         "Ein Entity-Attribute-Value-Modell wäre in dritter Normalform und trotzdem falsch: "
         "keine Typsicherheit, keine Fremdschlüssel, keine Constraints, unlesbare Abfragen. "
         "Normalform ist notwendig, nicht hinreichend."),
    ]:
        # Auf Diagrammfolien steht die Anna-Zeile OBEN statt unten und
        # ersetzt die Einleitung. Der Titel trägt die Aussage, die
        # Vertiefung steht in den Notizen — und das Diagramm bekommt die
        # ganze Fläche, sonst ist es im Hörsaal nicht lesbar.
        s = folie(prs, kicker, titel)
        faden(s, fadentext, y=88)
        diagramm(s, bild(bilddatei), y=134, hoehe=ZONE_UNTEN - 134 - 8)
        notizen(s, intro + "\n\n" + notiz)

    # ═══════════════════════════════════════════════════ 3 Normalisierung
    kapitel(prs, 3, "Normalisierung",
            "Warum reicht für Annas Fahrt nicht eine einzige Tabelle?",
            "Ein Prüfverfahren, kein Entwurfsverfahren. Es findet Redundanz — es sagt "
            "nicht, ob das Modell die Wirklichkeit trifft.")

    s = folie(prs, "3 · Normalisierung", "Der naive Versuch: alles in eine Tabelle",
              "So würde jemand Annas Fahrt aufschreiben, der noch nicht normalisiert hat. "
              "Die Tabelle funktioniert — bis zur ersten Änderung.")
    diagramm(s, bild("norm-vorher"), y=176, x0=FLUCHT_L, breite=390, hoehe=280)
    sandkarte(s, "Drei Anomalien, die daraus folgen",
              ["Einfügen: ein neuer Fahrradtyp lässt sich ohne Ausleihe gar nicht anlegen.",
               "Ändern: eine Preisanpassung müsste in jeder einzelnen Ausleihzeile nachgezogen werden.",
               "Löschen: mit der letzten Ausleihe eines Typs verschwindet auch dessen Preis."],
              y=200, x0=FLUCHT_L + 420, breite=BREITE - 420)
    notizen(s, "Erst die Anomalien zeigen, dann zerlegen. Wer sofort zerlegt, weiß "
               "hinterher nicht, welches Problem er gelöst hat.")

    s = folie(prs, "3 · Normalisierung", "1NF und 2NF: atomar, und alles hängt am ganzen Schlüssel",
              "Die 1NF verlangt atomare Attribute. Der lehrreichere Fall ist die "
              "wiederholende Gruppe: sie wird zu einer eigenen Relation, nicht zu mehreren "
              "Spalten.")
    vorher_nachher(s,
        ("verletzt", "1NF · nicht atomar",
         ["kunde_name  = „Anna Beispiel“",
          "positionen  = „Start 0,10; Zeit 6,10; …“",
          "",
          "Falsche Reparatur:",
          "position1, position2, position3",
          "→ legt die Anzahl im Schema fest"], True),
        ("erfüllt", "1NF · zerlegt",
         ["vorname     = „Anna“",
          "nachname    = „Beispiel“",
          "",
          "entgeltposition als eigene Relation,",
          "eine Zeile je Preisbestandteil",
          "→ die Anzahl ist offen"], True),
        y=176, hoehe=190)
    sandband(s, "2NF: bei zusammengesetztem Schlüssel (ausleihe_id, position_nr) hängt "
                "kunde_email nur an ausleihe_id — daraus folgt die Trennung in Kopf und "
                "Positionen.", y=380)
    notizen(s, "Das Kopf-Positionen-Muster kennen die meisten aus der Praxis. Hier sehen "
               "sie zum ersten Mal, dass es formal begründbar ist und nicht bloß Konvention.")

    s = folie(prs, "3 · Normalisierung", "3NF — und warum sie hier nicht genügt",
              "Die transitive Kette vom Rad über den Typ zum Preis ist der klassische "
              "3NF-Verstoß. Sie aufzulösen ist notwendig, reicht aber fachlich nicht aus.")
    vorher_nachher(s,
        ("3NF verletzt", "Die transitive Kette",
         ["ausleihe_id",
          "  → rahmennummer",
          "     → typ_bezeichnung",
          "        → preis_pro_minute",
          "",
          "Eine Preisänderung müsste in jeder",
          "Ausleihzeile nachgezogen werden."], True),
        ("3NF erfüllt — reicht trotzdem nicht", "Zerlegt, aber immer noch falsch",
         ["fahrradtyp und nutzungspreis",
          "werden eigene Relationen.",
          "",
          "ABER: liegt der Preis am Typ, ändert",
          "eine Anpassung weiterhin rückwirkend",
          "ALLE Altrechnungen — auch Annas."], False),
        y=176, hoehe=200)
    sandband(s, "Erst der Gültigkeitszeitraum in nutzungspreis macht das Modell fachlich "
                "richtig. Normalisierung beseitigt Redundanz — sie ersetzt keine fachliche "
                "Analyse.", y=390)
    notizen(s, "Das ist die zentrale Botschaft des Blocks und eine gute Prüfungsfrage: "
               "Nennen Sie eine Stelle, an der ein Modell in dritter Normalform trotzdem "
               "fachlich falsch ist.")

    s = folie(prs, "3 · Normalisierung", "Das Ergebnis der Zerlegung",
              "Aus einer Tabelle sind acht geworden — jede mit genau einer Aufgabe.")
    oben = unter_intro(s)
    diagramm(s, bild("norm-nachher"), y=oben, hoehe=ZONE_UNTEN - 40 - oben)
    faden(s, "Annas Fahrt liegt jetzt in ausleihe — ihr Preis in nutzungspreis, ihre Abrechnung in entgeltposition.")
    notizen(s, "Acht Relationen statt einer klingt nach mehr Aufwand. Der Aufwand entsteht "
               "einmal beim Entwurf und spart sich bei jeder Änderung wieder ein.")

    s = folie(prs, "3 · Normalisierung", "Exkurs: aus der Postleitzahl folgt nicht der Ort",
              "Der Lehrbuchklassiker lautet, plz bestimme ort, also gehöre ort in eine "
              "eigene Relation. In Deutschland stimmt das nicht — und der scheinbar saubere "
              "Zerlegungsschritt wäre fachlich falsch.")
    ampel_matrix(s, ["gilt"], [
        ("Eine PLZ hat genau einen Ort", [False],
         "Ländliche Sammel-PLZ umfassen mehrere Orte"),
        ("Ein Ort hat genau eine PLZ", [False],
         "Großstädte haben Dutzende"),
        ("plz → ort ist funktional", [False],
         "Also keine transitive Abhängigkeit, kein Split"),
    ], y=206, zeilen_h=52, chip_b=80, label_b=340)
    sandband(s, "Eine funktionale Abhängigkeit ist eine Behauptung über die Wirklichkeit "
                "und muss geprüft werden — nicht aus Namensähnlichkeit geschlossen.", y=396)
    notizen(s, "Hier lohnt der Blick ins Data Dictionary: die Begründung steht als "
               "Kommentar direkt an der Spalte adresse.ort und veraltet nicht.")

    # ═══════════════════════════════════════════════════ 4 Logischer Entwurf
    kapitel(prs, 4, "Logischer Entwurf",
            "Wie wird aus dem Bild ein Satz von Relationen?",
            "Regelbasiert und damit prüfbar. Wer eine 1:N-Beziehung auf der Eins-Seite "
            "verankert, hat nicht anders entworfen, sondern falsch abgebildet.")

    s = folie(prs, "4 · Logischer Entwurf", "Die Abbildung folgt festen Regeln",
              "Dieser Schritt ist mechanisch — und genau deshalb korrigierbar. Ein Fehler "
              "hier ist kein Geschmacksurteil, sondern ein Regelverstoß.")
    tabelle(s, ["ERM-Konstrukt", "Abbildung ins Relationenmodell", "Beispiel im Modell"],
            [["Entitätstyp", "eigene Relation", "kunde, station, ausleihe"],
             ["Attribut", "Spalte", "vorname, kapazitaet"],
             ["1:N-Beziehung", "Fremdschlüssel auf der N-Seite", "ausleihe.kunde_id"],
             ["1:1-Beziehung", "FK auf der optionalen Seite, dort zugleich PK", "fahrrad_position.fahrrad_id"],
             ["M:N-Beziehung", "eigene Verknüpfungsrelation", "Phase 2: wartungsposition"],
             ["Mehrwertiges Attribut", "eigene Relation", "fahrradtyp_merkmal"]],
            y=186, spalten_b=[210, 400, 293.5], zeilen_h=32)
    faden(s, "Annas Fahrt bekommt vier Fremdschlüssel: Kunde, Rad, Mitgliedschaft, Startstation.")
    notizen(s, "Die 1:1-Abbildung ist die elegante: fahrrad_id ist in fahrrad_position "
               "zugleich Primär- und Fremdschlüssel. Damit ist „höchstens eine Position je "
               "Rad“ ohne zusätzlichen Constraint erzwungen.")

    s = folie(prs, "4 · Logischer Entwurf", "Jede Relation trägt zwei Schlüssel",
              "Surrogatschlüssel als Primärschlüssel, fachlicher Schlüssel als "
              "Eindeutigkeitsbedingung. Beide haben eine Aufgabe, keiner ersetzt den anderen.")
    ampel_matrix(s, ["stabil", "eindeutig", "sprechend"], [
        ("Surrogatschlüssel  kunde_id", [True, True, False], "Primärschlüssel"),
        ("Fachschlüssel  kundennummer", [False, True, True], "UNIQUE-Bedingung"),
        ("Fachschlüssel als Primärschlüssel", [False, True, True], "Änderung zieht durch alle Verweise"),
    ], y=206, zeilen_h=52, chip_b=86, label_b=300)
    sandband(s, "Eine Rahmennummer wird beim Rahmentausch neu vergeben. Ein "
                "Primärschlüssel, der sich ändert, zieht die Änderung durch jede "
                "verweisende Relation.", y=396)
    notizen(s, "Die dritte Zeile ist der häufigste Anfängerfehler: den sprechenden "
               "Schlüssel zum Primärschlüssel machen. Er ist eindeutig und sprechend — aber "
               "eben nicht stabil.")

    # ═══════════════════════════════════════════════════ 5 Physischer Entwurf
    kapitel(prs, 5, "Physischer Entwurf",
            "Wie erzwingt die Datenbank, dass Annas Rechnung stimmt?",
            "Hier wird eine Geschäftsregel entweder erzwungen oder nur gewünscht. Jede "
            "Regel, die nicht in einem Constraint steht, wird früher oder später verletzt.")

    s = folie(prs, "5 · Physischer Entwurf", "Datentypen sind fachliche Entscheidungen",
              "Jede Zeile hat einen Grund, der über Geschmack hinausgeht. Die beiden ersten "
              "sind die, an denen in der Praxis am meisten schiefgeht.")
    tabelle(s, ["Zweck", "Typ", "Warum genau dieser"],
            [["Zeitpunkt", "timestamptz", "timestamp ohne Zone verliert den Bezug; Sommerzeit wird doppeldeutig"],
             ["Geldbetrag", "numeric(10,2)", "exakte Dezimalarithmetik; float trifft 0,10 nicht genau"],
             ["Zeitraum", "daterange", "ein Zeitraum ist ein Wert, kein Spaltenpaar — erst so prüfbar"],
             ["Surrogatschlüssel", "bigint IDENTITY", "standardkonform; ALWAYS verhindert Setzen von außen"],
             ["Koordinate", "numeric(9,6)", "sechs Nachkommastellen entsprechen etwa elf Zentimetern"],
             ["Text", "text", "PostgreSQL speichert text und varchar(n) identisch"]],
            y=186, spalten_b=[180, 170, 553.5], zeilen_h=32)
    faden(s, "Annas Startzeit als timestamp ohne Zone wäre in der Nacht der Zeitumstellung nicht eindeutig.")
    notizen(s, "Die Zeitzonenfalle ist real und teuer: eine Fahrt in der Nacht der "
               "Umstellung lässt sich mit timestamp ohne Zone nicht eindeutig einordnen.")

    s = folie(prs, "5 · Physischer Entwurf", "Zehn von fünfzehn Regeln erzwingt die Datenbank",
              "Die oberen acht gelten immer, auch bei direktem SQL-Zugriff. Die unteren "
              "vier nur, wenn man den vorgesehenen Weg nimmt — sie brauchen Kontext, den ein "
              "Constraint nicht hat.")
    tabelle(s, ["Regel", "Umsetzung", "Wirkt"],
            [["GR1", "CREATE UNIQUE INDEX … WHERE status = 'aktiv'", "immer"],
             ["GR3, GR4", "EXCLUDE USING gist (… WITH =, … WITH &&)", "immer"],
             ["GR6", "GENERATED ALWAYS AS (ceil(…)) STORED", "immer"],
             ["GR7", "CHECK (verbraucht <= kontingent)", "immer"],
             ["GR10", "UNIQUE (kunde_id, periode_jahr, periode_monat)", "immer"],
             ["GR11", "CHECK: Station oder Koordinaten, nie beides", "immer"],
             ["GR13", "CHECK + Constraint-Trigger (braucht den Radstatus)", "immer"],
             ["GR15", "Constraint-Trigger: zählt Räder gegen Stellplätze", "immer"],
             ["GR2, GR5", "Prüfung in fn_ausleihe_starten und fn_ausleihe_abrechnen", "nur über die Funktion"],
             ["GR8, GR9", "Prüfung in der api_-Schicht", "nur über die Funktion"],
             ["GR12", "Prüfung in fn_ausleihe_starten", "nur über die Funktion"],
             ["GR14", "fn_im_geschaeftsgebiet in fn_ausleihe_beenden", "nur über die Funktion"]],
            y=176, spalten_b=[130, 520, 253.5], zeilen_h=24)
    notizen(s, "Der Unterschied in der dritten Spalte ist der eigentliche Inhalt dieser "
               "Folie. Constraints sind unbestechlich, Funktionen kann man umgehen — "
               "deshalb muss die Fachlogik von außen unerreichbar sein.")

    s = folie(prs, "5 · Physischer Entwurf", "EXCLUDE verhindert, was UNIQUE nicht kann",
              "Ein UNIQUE kennt nur Gleichheit. Für „diese Zeiträume dürfen sich nicht "
              "überschneiden“ braucht es einen Operator, der Überlappung prüfen kann.")
    code_kacheln(s,
        ("Der Constraint",
         ["constraint nutzungspreis_ueberschneidung_ex",
          "  exclude using gist (",
          "    typ_id      with =,",
          "    gueltigkeit with &&",
          "  )",
          "",
          "Keine zwei Zeilen mit gleichem typ_id",
          "UND überlappenden Zeiträumen."], BLAU),
        ("Zwei Fallstricke",
         ["1  Voraussetzung btree_gist",
          "   Ohne die Erweiterung fehlt bigint",
          "   die Operatorklasse für gist:",
          "   „data type bigint has no default",
          "    operator class for access method",
          "    gist“",
          "",
          "2  Halboffene Grenzen '[)'",
          "   Sonst ist der Wechseltag doppelt",
          "   belegt."], ORANGE),
        y=176, hoehe=240)
    faden(s, "Ohne GR4 könnte es für Annas Fahrtzeitpunkt zwei gültige Preise geben — die Rechnung wäre nicht bestimmt.")
    notizen(s, "Beide Fallstricke sind beim Bauen tatsächlich aufgetreten. Der erste hat "
               "die Anlage abgebrochen, der zweite wäre bei einem nahtlosen Tarifwechsel "
               "aufgefallen.")

    s = folie(prs, "5 · Physischer Entwurf", "Warum das Mindestalter kein CHECK sein darf",
              "Naheliegend wäre eine Bedingung mit current_date. PostgreSQL akzeptiert sie "
              "sogar. Trotzdem ist sie falsch — und der Grund zeigt, was ein CHECK "
              "eigentlich ist.")
    vorher_nachher(s,
        ("verlockend und falsch", "CHECK mit current_date",
         ["check (geburtsdatum <=",
          "  current_date - interval '16 years')",
          "",
          "current_date ist nicht IMMUTABLE.",
          "Ein CHECK wird beim Schreiben geprüft",
          "UND beim Wiedereinspielen eines Dumps.",
          "",
          "Ein Kunde, der bei der Anmeldung 16 war,",
          "bleibt es — der Restore bricht trotzdem ab."], False),
        ("richtig", "Immutable auf der Tabelle, Regel in der Funktion",
         ["check (geburtsdatum is null",
          "   or geburtsdatum between",
          "      date '1900-01-01'",
          "  and date '2100-01-01')",
          "",
          "-- die Altersregel in api_profil_aktualisieren:",
          "if p_geburtsdatum >",
          "   current_date - interval '16 years'",
          "then return 'Mindestalter nicht erreicht';"], False),
        y=176, hoehe=230)
    sandband(s, "Merksatz: Ein CHECK darf nur von der Zeile selbst abhängen — und von "
                "nichts, was sich mit der Zeit ändert.", y=420)
    notizen(s, "Gute Prüfungsfrage: Warum ist eine CHECK-Bedingung mit current_date "
               "gefährlich, obwohl die Datenbank sie annimmt?")

    # ═══════════════════════════════════════════════════ 6 Implementierung
    kapitel(prs, 6, "Implementierung",
            "Was kostet Annas Fahrt genau — und warum?",
            "Jetzt lösen wir die Frage vom Anfang auf. Jeder Cent muss begründbar sein.")

    s = folie(prs, "6 · Implementierung", "21 Aufbauschritte, jeder für sich lauffähig",
              "Jede Datei ist idempotent: sie läuft zweimal hintereinander fehlerfrei. Das "
              "ist die Voraussetzung dafür, dass man einen Aufbau gefahrlos wiederholen kann.")
    schichtenstapel(s, [
        ("0001 Schema, Erweiterungen, Aufzählungstypen, Audit-Mechanik", True),
        ("0002 bis 0007 · die sechs Fachbereiche A bis F, 28 Tabellen", False),
        ("0008 Referenzdaten: Entgeltarten, Preise, Tarife, Inhalte", False),
        ("0009 Geschäftslogik: fn_-Fachlogik und api_-Zugriffsschicht", True),
        ("0010 Sichten · 0011 Zugriffsschutz · 0012 Data Dictionary", True),
        ("0014 bis 0020 · sieben Schritte mehr für die Warenwirtschaft", False),
    ], y=182, hoehe=44, luecke=8)
    notizen(s, "Die Reihenfolge ist die Reihenfolge dieser Vorlesung. Wer 0004 vor 0003 "
               "ausführt, bekommt einen Fremdschlüsselfehler — die Abhängigkeiten sind im "
               "Schema selbst dokumentiert. Zur Nummernlücke: 0013 gab es, sie hieß "
               "altsystem_abloesen und zog nach db/betrieb/, weil sie gegen eine leere "
               "Datenbank nicht durchlief. Lücken in Migrationsnummern sind normal — "
               "Nummern werden nie neu vergeben, sonst stimmt keine Historie mehr.")

    s = folie(prs, "6 · Implementierung", "Annas Rechnung: 4,96 Euro, Zeile für Zeile",
              "Das Zeitentgelt wird über ALLE Minuten gebildet und die Freiminuten als "
              "eigene Gutschrift abgezogen. So ist auf der Rechnung ablesbar, was der "
              "Tarifvorteil wert war.")
    tabelle(s, ["Position", "Menge", "Einzelbetrag", "Betrag", "Woher der Wert kommt"],
            [["STARTGEBUEHR", "1", "0,10", "+ 0,10", "nutzungspreis, gültig um 10:00 Uhr"],
             ["ZEITENTGELT", "61", "0,10", "+ 6,10", "dauer_minuten, aufgerundet"],
             ["FREIMINUTEN", "0", "0,10", "- 0,00", "Kontingent im Monat schon verbraucht"],
             ["TARIFRABATT", "1", "20 %", "- 1,24", "tarif_kondition zum Startzeitpunkt"],
             ["HOECHSTPREIS", "1", "—", "- 0,00", "6,20 liegt unter der Obergrenze 50,00"]],
            y=178, spalten_b=[190, 90, 130, 120, 373.5], zeilen_h=30)
    sandband(s, "6,20 - 1,24 = 4,96 Euro. Reihenfolge: Rabatt VOR der Kappung — umgekehrt "
                "käme 4,00 heraus, weil der Rabatt den bereits gedeckelten Betrag ein "
                "zweites Mal senken würde.", y=364)
    faden(s, "Damit ist die Frage vom Anfang beantwortet: 4,96 Euro, und jede Zeile ist belegt.")
    notizen(s, "Das ist die Auflösung des roten Fadens. Neun Testfälle sichern die "
               "Berechnung ab; einer prüft ausdrücklich die Reihenfolge, weil genau dort "
               "der teure Fehler sitzt.")

    s = folie(prs, "6 · Implementierung", "Dokumentation, die nicht veralten kann",
              "Ein Data Dictionary, das von Hand gepflegt wird, ist nach dem zweiten "
              "Release falsch. Die Lösung ist, es aus dem Systemkatalog zu erzeugen — und "
              "die Vollständigkeit maschinell zu erzwingen.")
    code_kacheln(s,
        ("Beschreibung an das Objekt",
         ["comment on column",
          "  velocity.fahrrad_position.station_id is",
          "  'NULL bedeutet: das Rad steht frei",
          "   abgestellt, nicht an einer Station';",
          "",
          "Der Kommentar erklärt, was NULL",
          "bedeutet — nicht, wie die Spalte heißt."], BLAU),
        ("Vollständigkeit erzwingen",
         ["-- Test schlägt fehl, solange eine",
          "-- Fachspalte ohne Kommentar ist",
          "where col_description(c.oid, a.attnum)",
          "      is null",
          "  and a.attname not in",
          "      ('erstellt_am','geaendert_am');",
          "",
          "338 Spalten im Aufbau, 264 beschrieben,",
          "74 technische Audit-Spalten ausgenommen."], GRUEN_D),
        y=176, hoehe=240)
    notizen(s, "Der Trick ist der Test, nicht der Kommentar. Ohne ihn schreibt man die "
               "ersten zwanzig Kommentare und vergisst die restlichen zweihundert.")

    # ═══════════════════════════════════════════════════ 7 Sicherheit
    kapitel(prs, 7, "Zugriffsschutz",
            "Wer außer Anna darf Annas Daten sehen?",
            "Der Schutz liegt in der Datenbank. Auf den Browser ist grundsätzlich kein "
            "Verlass.")

    s = folie(prs, "7 · Zugriffsschutz", "Der Schlüssel im Browser ist kein Geheimnis",
              "Supabase liefert den anon-Key an jeden Besucher aus. Er steht im Quelltext "
              "der Seite. Jeder kann damit beliebige Anfragen stellen — mit curl, ohne die "
              "Website je zu öffnen.")
    vorher_nachher(s,
        ("das Antipattern", "Der häufigste Supabase-Anfängerfehler",
         ["create policy \"alles fuer alle\"",
          "  on kunde for all",
          "  to anon, authenticated",
          "  using (true) with check (true);",
          "",
          "Folge: jeder mit dem öffentlichen",
          "Schlüssel liest UND ändert sämtliche",
          "Kundendaten — auch Annas."], False),
        ("default deny", "Was nicht erlaubt ist, ist verboten",
         ["-- RLS auf jeder Basistabelle",
          "alter table … enable row level security;",
          "",
          "-- keine Policy für anon",
          "-- Lesen nur über v_-Sichten",
          "-- Schreiben nur über api_-Funktionen",
          "",
          "Anna sieht ihre Fahrt, sonst niemand."], False),
        y=176, hoehe=220)
    sandband(s, "Jede Zugriffsbeschränkung, die im JavaScript steht, ist wirkungslos.", y=410)
    notizen(s, "Im Schema velocity_demo liegt eine Tabelle mit erfundenen Daten, an der die "
               "Studierenden das Antipattern live ausprobieren dürfen — ohne den echten "
               "Bestand zu öffnen.")

    s = folie(prs, "7 · Zugriffsschutz", "Was der Browser erreicht — und was nicht",
              "Zwei Schichten trennen die Anwendung von den Tabellen — die Fachlogik dazwischen "
              "prüft selbst nicht auf den angemeldeten Nutzer.")
    oben = unter_intro(s)
    diagramm(s, bild("sicherheit-schichten"), y=oben, hoehe=ZONE_UNTEN - 12 - oben)
    notizen(s, "Der untere rote Kasten ist der Kern: fn_ausleihe_beenden bekommt die "
               "kunde_id als Parameter und glaubt sie. Wäre die Funktion von außen "
               "aufrufbar, könnte jeder fremde Ausleihen abrechnen.")

    s = folie(prs, "7 · Zugriffsschutz", "Die Falle, die fast jeder übersieht",
              "PostgreSQL vergibt das Ausführungsrecht auf jede neu angelegte Funktion "
              "automatisch an die Rolle PUBLIC. Ein Entzug gegen einzelne Rollen greift "
              "deshalb nicht.")
    vorher_nachher(s,
        ("wirkungslos", "REVOKE gegen die Rollen allein",
         ["revoke all on all functions",
          "  in schema velocity",
          "  from anon, authenticated;",
          "",
          "Beide erben das Recht weiterhin",
          "über PUBLIC. Die Fachlogik bleibt",
          "aufrufbar."], False),
        ("wirksam — der einzige Schutz", "PUBLIC ausdrücklich mit entziehen",
         ["revoke all on all functions",
          "  in schema velocity",
          "  from public, anon, authenticated;",
          "",
          "Muss nach jeder neu angelegten",
          "Funktion in diesem Schema",
          "erneut laufen."], False),
        y=176, hoehe=190)
    sandkarte(s, "Placebo gemessen statt angenommen",
              ["alter default privileges … revoke execute on functions from public "
               "sollte künftige Funktionen automatisch schützen. Gegenprobe auf dieser "
               "Instanz: kein Eintrag in pg_default_acl, eine danach angelegte "
               "Testfunktion bekam trotzdem EXECUTE für PUBLIC.",
               "In einer gewöhnlichen PostgreSQL-Installation wirkt die Anweisung — "
               "hier nachweislich nicht. Deshalb bleibt das REVOKE oben die einzige "
               "Absicherung, und es muss nach jedem CREATE FUNCTION erneut laufen."],
              y=376, hoehe=108, warnung=True)
    notizen(s, "Der interessante Lehrpunkt ist nicht die ALTER-DEFAULT-PRIVILEGES-Zeile "
               "selbst, sondern dass sie in einer gewöhnlichen PostgreSQL-Installation "
               "wirkt und auf dieser Instanz gemessen nicht: kein Eintrag in "
               "pg_default_acl, gegengeprüft per Abfrage direkt danach und durch eine "
               "Testfunktion in einem Scratch-Schema, die trotzdem EXECUTE für PUBLIC "
               "bekam. Ein Sicherheitskonzept, das eine Schutzmaßnahme behauptet, die es "
               "nicht gibt, ist gefährlicher als eine bekannte Lücke. Der einzige "
               "tatsächliche Schutz für künftig angelegte Funktionen ist das explizite "
               "REVOKE weiter oben — erneut ausgeführt nach jeder neu angelegten Funktion "
               "in diesem Schema. Aufgedeckt hat das ein Test, nicht Nachdenken: "
               "test_s_keine_oeffentliche_funktion in db/tests/t0011_sicherheit.sql und "
               "Abnahmeprüfung 25 fangen das Vergessen zweifach ab.")

    s = folie(prs, "7 · Zugriffsschutz", "Nachweis statt Behauptung — auf drei Wegen",
              "Ein Sicherheitskonzept, das nur beschrieben ist, ist wertlos. Geprüft wird in "
              "der Datenbank, über die Schnittstelle und im Browser.")
    regel_streifen(s, [
        ("In der Datenbank", "RLS überall aktiv, anon ohne Tabellenrechte, Rollenwechsel-Probe", "pgTAP, t0011"),
        ("Über die Schnittstelle", "13 gesperrte Ressourcen, 9 öffentliche Sichten", "HTTP 401 gegen 200"),
        ("Im Browser", "Aufrufe der abgemeldeten Seite in der Konsole", "permission denied"),
    ], y=188, hoehe=52, chip_b=200)
    sandkarte(s, "Eine Falle im Prüfwerkzeug selbst",
              ["Der erste Entwurf meldete alle 13 gesperrten Ressourcen als bestanden — "
               "obwohl er nichts geprüft hatte: das Schema war gar nicht exponiert, also war "
               "alles unerreichbar, Sicheres wie Unsicheres.",
               "Ein Test, der „abgesichert“ nicht von „gar nicht erreichbar“ unterscheidet, "
               "ist gefährlicher als kein Test."],
              y=368, warnung=True)
    notizen(s, "Diese Folie ist mir die wichtigste des Blocks. Der häufigste Fehler bei "
               "Sicherheitstests ist nicht der fehlende Test, sondern der Test, der aus dem "
               "falschen Grund grün wird.")

    # ═══════════════════════════════════════════════════ 8 Anbindung
    kapitel(prs, 8, "Anwendung anbinden",
            "Wie kommt Annas Fahrt in die Web-Anwendung?",
            "Die Sichten sind der Vertrag zwischen Datenbank und Oberfläche.")

    s = folie(prs, "8 · Anwendung anbinden", "Nur Sichten lesen, nur Funktionen schreiben",
              "Die Anwendung kennt keine einzige Basistabelle. Das ist keine Konvention, "
              "sondern erzwungen: sie käme gar nicht an sie heran.")
    code_kacheln(s,
        ("Lesen — ausschließlich Sichten",
         ["from('v_station')",
          "from('v_verfuegbares_fahrrad')",
          "from('v_tarifkarte')",
          "from('v_faq')",
          "from('v_kennzahl')",
          "",
          "-- angemeldet:",
          "from('v_meine_ausleihe')",
          "from('v_mein_profil')"], BLAU),
        ("Schreiben — ausschließlich Funktionen",
         ["rpc('api_kunde_sicherstellen')",
          "rpc('api_ausleihe_starten')",
          "rpc('api_ausleihe_beenden')",
          "rpc('api_profil_aktualisieren')",
          "",
          "-- nicht erreichbar:",
          "from('kunde')      permission denied",
          "from('ausleihe')   permission denied",
          "rpc('fn_…')        permission denied"], GRUEN_D),
        y=176, hoehe=240)
    faden(s, "Anna sieht ihre Fahrt über v_meine_ausleihe — begrenzt durch RLS auf ihre eigenen Zeilen.")
    notizen(s, "Auch die Inhalte der Seite kommen aus der Datenbank: Tarifkarten, FAQ, die "
               "Schritte der Anleitung und die Kennzahlen. Vorher standen sie fest im HTML "
               "und mussten bei jeder Preisänderung von Hand nachgezogen werden.")

    # ═══════════════════════════════════════════ 9 Warenwirtschaft
    kapitel(prs, 9, "Die Warenwirtschaft",
            "Wer arbeitet mit Annas Fahrt, wenn sie längst bezahlt ist?",
            "Bis hierher ging es um die Kundensicht. Jetzt kommt die zweite Anwendung "
            "auf dieselbe Datenbank — und die stellt ganz andere Fragen an dieselben Zeilen.")

    s = folie(prs, "9 · Die Warenwirtschaft", "Zweite Anwendung, dasselbe Datenmodell",
              "Die Warenwirtschaft bekam kein eigenes Schema und keine eigene Datenbank. "
              "Sie legt Sichten und Funktionen über genau die Tabellen, die schon da waren.")
    kachelreihe(s, [
        ("Kein zweites Schema",
         ["38 Basistabellen legt der Aufbau an — für beide dieselben",
          "Eine Kundenadresse ist eine Kundenadresse",
          "Keine Kopie, die veralten könnte"]),
        ("19 Sichten, 15 Funktionen",
         ["Gelesen wird über v_wawi_…",
          "Geschrieben über api_…",
          "Fünf weitere api_ dienen der Website"]),
        ("Fünf Rollen",
         ["leitung · disposition · werkstatt",
          "kundenservice · demo",
          "Die Rolle steht in der Datenbank"]),
    ], y=182, hoehe=150)
    faden(s, "Annas Fahrt steht in genau einer Zeile — die Werkstatt und die Leitung lesen dieselbe.")
    notizen(s, "Der wichtigste Satz dieser Folie ist der letzte der ersten Kachel. Ein "
               "getrenntes Auswertungssystem mit eigener Kopie ist der Normalfall in der "
               "Praxis — und die häufigste Quelle für Zahlen, die sich widersprechen. "
               "Hier gibt es diese Möglichkeit gar nicht.")

    s = folie(prs, "9 · Die Warenwirtschaft", "Vier Fachrollen, vier Ausschnitte derselben Daten",
              "Was jemand sieht, entscheidet die Datenbank anhand seiner Rolle — nicht die "
              "Oberfläche. Ein verstecktes Menü ist kein Schutz.")
    ampel_matrix(s, ["Leitung", "Dispo", "Werkstatt", "Service"], [
        ("Flotte, einzelne Räder", [True, True, True, False], "v_wawi_flotte"),
        ("Stationen und Belegung", [True, True, False, False], "v_wawi_station"),
        ("Kundschaft, Stammdaten", [True, False, False, True], "v_wawi_kunde, api_kunde_…"),
        ("Schäden und Aufträge", [True, False, True, False], "v_wawi_schaden, v_wawi_auftrag"),
        ("Umsatz und Kennzahlen", [True, False, False, False], "v_wawi_umsatz_…, v_wawi_km_co2"),
        ("Einzelfahrten mit Ort", [False, False, False, False], "v_wawi_fahrt_km: für niemanden freigegeben"),
    ], y=182, zeilen_h=40, luecke=6, chip_b=78, label_b=250)
    notizen(s, "Die letzte Zeile ist die wichtigste: v_wawi_fahrt_km führt Einzelfahrten mit "
               "Kundennummer — ein Bewegungsprofil. Die Sicht ist für keine Rolle freigegeben; "
               "sie dient nur als Zwischenschritt innerhalb anderer Sichten. Datensparsamkeit "
               "heißt hier: die aggregierte Sicht ist offen, die feine nicht. "
               "Zur fünften Rolle: demo steht bewusst nicht in der Matrix. Sie ist keine "
               "Fachrolle, sondern der öffentliche Vorführzugang — sie liest dieselben "
               "Sichten wie die Leitung, schreibt aber nichts und sieht Einzelfahrten "
               "ebenso wenig. Wer sie als Spalte führt, suggeriert eine Aufgabe, die es "
               "nicht gibt.")

    s = folie(prs, "9 · Die Warenwirtschaft", "Zwei Arten von Sicht — und woran man sie erkennt",
              "Das Korn einer Sicht ist die Antwort auf die Frage: Wofür steht eine Zeile? "
              "Es unterscheidet die beiden Arten — und wer es vermischt, baut Sichten, die "
              "für beides zu langsam und für beides zu ungenau sind.")
    code_kacheln(s,
        ("Arbeitssicht — ein Ding je Zeile",
         ["v_wawi_flotte      ein Rad",
          "v_wawi_kunde       ein Kunde",
          "v_wawi_station     eine Station",
          "v_wawi_schaden     eine Meldung",
          "",
          "Korn:   das Objekt selbst",
          "Zweck:  suchen, öffnen, ändern",
          "Frage:  Wo steht Rad CB-00035?"], BLAU),
        ("Auswertungssicht — eine Gruppe je Zeile",
         ["v_wawi_umsatz_radtyp     Monat × Typ",
          "v_wawi_fahrten_je_tag    Tag × Typ",
          "v_wawi_km_co2            Monat × Typ",
          "v_wawi_stationsauslastung  Station",
          "",
          "Korn:   die Gruppierung",
          "Zweck:  vergleichen, entscheiden",
          "Frage:  Welcher Typ trägt den Umsatz?"], GRUEN_D),
        y=182, hoehe=250)
    notizen(s, "Das Korn — die Frage „wofür steht eine Zeile?“ — ist der wichtigste Begriff "
               "der ganzen Auswertung. Wer es nicht benennen kann, wird die Summen falsch "
               "bilden. Ein häufiger Fehler: Umsatz über eine Sicht summieren, in der jede "
               "Fahrt mehrfach vorkommt, weil ein Join sie vervielfacht hat.")

    # ═══════════════════════════════════════ 10 Analytisches Modell
    kapitel(prs, 10, "Das analytische Datenmodell",
            "Welche Frage stellt die Leitung — und welche Zeile beantwortet sie?",
            "Jetzt kommt der Perspektivwechsel: vom Schreiben zum Lesen, vom Einzelfall "
            "zur Summe. Dasselbe Modell, eine andere Brille.")

    s = folie(prs, "10 · Analytisches Modell", "Fakten sind Ereignisse, Dimensionen sind Fragen",
              "Ein Fakt ist etwas, das passiert ist und sich zählen lässt. Eine Dimension "
              "ist die Richtung, aus der jemand darauf schaut.")
    oben = unter_intro(s)
    diagramm(s, bild("analytik-stern"), y=oben, hoehe=196)
    sandkarte(s, "Keine dieser Dimensionen ist eine eigene Tabelle",
              ["Zeit steckt in ausleihe.startzeit · Fahrradtyp in fahrrad → fahrradmodell · "
               "Station in start_station_id und end_station_id",
               "Tarifgruppe in kunde → mitgliedschaft → tarif · Wohnort in kunde → adresse. "
               "Die Warenwirtschaft benennt nur die Rollen und legt Sichten darüber."],
              y=oben + 206)
    faden(s, "Annas Fahrt ist eine Faktenzeile — und taucht in jeder der vier Auswertungen auf.")
    notizen(s, "Der Kasten unten rechts trägt die eigentliche Aussage: keine dieser "
               "Dimensionen ist eine eigene Tabelle. Wir haben kein Sternschema gebaut, "
               "sondern die Rollen benannt, die vorhandene Tabellen beim Auswerten einnehmen. "
               "Warum das hier richtig ist, steht zwei Folien weiter.")

    s = folie(prs, "10 · Analytisches Modell", "Vier Auswertungen — und was sie tatsächlich messen",
              "Jede Kennzahl braucht drei Angaben: woraus sie entsteht, worüber sie gruppiert "
              "und was sie NICHT sagt. Die dritte fehlt in der Praxis fast immer.")
    tabelle(s, ["Auswertung", "Maß", "Gruppiert über", "Sagt nichts über"],
            [["Umsatz nach Radtyp", "Σ betrag", "Monat × Fahrradtyp", "Kosten, also nicht Gewinn"],
             ["Umsatz nach Kundengruppe", "Σ betrag", "Monat × Tarif", "Kunden ohne Tarif"],
             ["Kilometer und CO2", "Σ kilometer", "Monat × Typ", "Strecken sind geschätzt"],
             ["Stationsauslastung", "Zu-/Abgang", "Station", "Warum umgesetzt wurde"]],
            y=186, spalten_b=[236, 118, 190, 359.5], zeilen_h=34)
    sandband(s, "Die vierte Spalte gehört auf jede Kennzahlenfolie. Eine Zahl ohne ihre "
                "Grenze wird als Antwort auf Fragen benutzt, die sie nicht beantwortet.",
             y=372)
    notizen(s, "Beispiel für die Tragweite: „Umsatz nach Radtyp“ zeigt das Lastenrad vorn. "
               "Daraus folgt NICHT, dass Lastenräder sich am meisten lohnen — Anschaffung, "
               "Wartung und Stellplatzbedarf stehen in keiner dieser Zeilen. Wer so "
               "entscheidet, verwechselt Umsatz mit Deckungsbeitrag.")

    s = folie(prs, "10 · Analytisches Modell", "Vom Monat bis zum einzelnen Rad in zwei Klicks",
              "Drei Sichten, drei Körner. Jede Ebene beantwortet dieselbe Frage eine Stufe "
              "genauer — und jede ist eine eigene Sicht, nicht ein Filter auf derselben.")
    oben = unter_intro(s)
    diagramm(s, bild("analytik-drilldown"), y=oben, hoehe=150)
    sandkarte(s, "Warum drei Sichten und nicht eine mit Filter",
              ["Eine Sicht hat genau ein Korn. Wer Monat, Tag und Rad in dieselbe Sicht legt, "
               "erzeugt Zeilen, die sich ohne doppeltes Zählen nicht mehr summieren lassen."],
              y=oben + 172)
    faden(s, "Zwei Klicks von „E-Bike, Oktober“ bis zu dem Rad, mit dem Anna gefahren ist.")
    notizen(s, "Drill-Down ist der Übergang von der Kennzahl zum Beleg. Er ist der Grund, "
               "warum Auswertungen überhaupt Vertrauen verdienen: eine Zahl, die man nicht "
               "bis zur einzelnen Zeile aufklappen kann, muss man glauben.")

    s = folie(prs, "10 · Analytisches Modell", "Sichten statt Sternschema — und was das kostet",
              "Ein eigenes Sternschema mit nächtlicher Beladung wäre der Lehrbuchweg. Für "
              "diesen Fall ist er falsch — aber nicht immer.")
    kachelreihe(s, [
        ("Sternschema mit Beladung",
         ["Eigene Fakten- und Dimensionstabellen",
          "",
          "✓ Schnell auch bei sehr vielen Zeilen",
          "✓ Historie bleibt erhalten, wenn sich",
          "   Stammdaten später ändern",
          "× Zweite Kopie der Wahrheit",
          "× Beladung bauen und überwachen",
          "× Zahlen so alt wie der letzte Lauf"]),
        ("Sichten auf dem operativen Modell",
         ["Was die Warenwirtschaft tut",
          "",
          "✓ Keine Kopie, keine Beladung, kein Versatz",
          "✓ Jede Zahl so aktuell wie die Fahrt selbst",
          "✓ Eine Wahrheit, ein Schutzkonzept",
          "× Jede Abfrage rechnet neu",
          "× Bei 12 269 Ausleihen unkritisch,",
          "   bei Millionen nicht mehr"]),
    ], y=180, hoehe=238, spalten=2)
    sandband(s, "Keine Spalte ist die richtige. Die Wahl hängt an Datenmenge und "
                "geforderter Aktualität — und dreht sich um, wenn eine davon sich ändert.",
             y=424)
    notizen(s, "Die Entscheidung hängt an zwei Zahlen: Datenmenge und geforderte Aktualität. "
               "Hier sind es gut zwölftausend Ausleihen und die Erwartung, dass eine gerade "
               "beendete Fahrt sofort in der Auswertung steht. Bei zehn Millionen Zeilen und "
               "einem Tagesbericht dreht sich das Urteil um. Wichtig ist, dass Studierende "
               "die Frage stellen — nicht, dass sie diese eine Antwort auswendig können.")

    s = folie(prs, "10 · Analytisches Modell", "Die eine Stelle, an der geschätzt wird",
              "Die Datenbank kennt Start und Ziel, aber nicht den gefahrenen Weg. Für "
              "Kilometer und CO2 muss geschätzt werden — und genau das steht in den Daten.")
    code_kacheln(s,
        ("Was gemessen ist",
         ["ausleihe.startzeit    Zeitstempel",
          "ausleihe.endzeit      Zeitstempel",
          "ausleihe.dauer_minuten  gerechnet",
          "entgeltposition.betrag  gebucht",
          "",
          "Diese Werte sind Belege.",
          "Sie stammen aus dem Vorgang selbst."], BLAU),
        ("Was geschätzt ist",
         ["v_wawi_fahrt_km",
          "  km              Schätzwert",
          "  ist_geschaetzt  true / false",
          "",
          "Luftlinie × Umwegfaktor,",
          "wenn keine Strecke aufgezeichnet ist.",
          "Die Spalte sagt, welcher Fall vorliegt."], ROT_A),
        y=182, hoehe=232)
    sandband(s, "Eine Schätzung darf in eine Auswertung — aber sie muss sich als solche zu "
                "erkennen geben. Genau dafür ist ist_geschaetzt da.", y=428)
    notizen(s, "Das ist der Punkt, an dem Datenmodellierung zur Redlichkeit wird. Die "
               "CO2-Zahl ist ein Werbeargument; sie beruht auf einer Schätzung. Wer die "
               "Schätzung nicht kennzeichnet, produziert eine Zahl, die niemand mehr "
               "einordnen kann — und im Zweifel eine, die er selbst irgendwann glaubt.")

    # ═══════════════════════════════════════════════ 11 Oberflächen
    kapitel(prs, 11, "Zwei Oberflächen, ein Weg",
            "Wie kommt Annas Rad auf zwei verschiedene Bildschirme?",
            "Kundenwebsite und Warenwirtschaft sind zwei Anwendungen — aber sie gehen "
            "denselben Weg in dieselbe Datenbank.")

    s = folie(prs, "11 · Zwei Oberflächen", "Ein Weg, zwei Anwendungen, dieselbe Regel",
              "Beide Oberflächen sprechen dieselbe Schnittstelle an. Was sie unterscheidet, "
              "ist nicht der Weg, sondern die Rolle des angemeldeten Kontos.")
    oben = unter_intro(s)
    diagramm(s, bild("web-architektur"), y=oben, hoehe=ZONE_UNTEN - 40 - oben)
    faden(s, "Anna und die Werkstatt fragen dieselbe Adresse — und bekommen verschiedene Zeilen.")
    notizen(s, "Der öffentliche Schlüssel ist in beiden Anwendungen derselbe und steht im "
               "Quelltext. Er ist kein Geheimnis, sondern nur die Angabe, welches Projekt "
               "gemeint ist. Wer etwas sehen darf, entscheidet die Rolle im angemeldeten "
               "Konto — das ist der Inhalt von Kapitel 7.")

    s = folie(prs, "11 · Zwei Oberflächen", "Was die Warenwirtschaft anders macht",
              "Dieselbe Datenbank, andere Aufgabe — und daraus folgen vier Unterschiede, "
              "die man am Modell ablesen kann.")
    regel_streifen(s, [
        ("Vorgangsklammer", "Jede Änderung gehört zu einem Vorgang und ist ihm zuzuordnen",
         "neuerVorgang()"),
        ("Änderungsprotokoll", "Wer hat wann was geändert — die Datenbank schreibt mit",
         "aenderungsprotokoll"),
        ("Rolle statt Besitz", "Nicht „meine Zeilen“, sondern „Zeilen meiner Aufgabe“",
         "hat_rolle('werkstatt')"),
        ("Kein Löschen", "Ausgemustert ist ein Zustand, keine gelöschte Zeile",
         "status = 'ausgemustert'"),
    ], y=190, hoehe=54, luecke=10, chip_b=210)
    notizen(s, "Der vierte Punkt ist der, den Studierende am häufigsten falsch machen. Ein "
               "ausgemustertes Rad zu löschen würde jede Fahrt entwerten, die damit gemacht "
               "wurde — und damit auch Annas Rechnung. Historie verträgt kein DELETE.")

    # ═══════════════════════════════════════════════════ 12 Abschluss
    kapitel(prs, 12, "Zusammenfassung und Ausblick",
            "Was bleibt aus zwölf Kapiteln — und was ist entworfen, aber nicht gebaut?",
            "Acht Sätze zum Mitnehmen — und der Rest des Modells, der auf dem Papier steht.")

    s = folie(prs, "12 · Zusammenfassung", "Acht Sätze, die diese Einheit tragen",
              "Wenn Sie nichts anderes mitnehmen, dann diese acht. Zu jedem sollten Sie ein "
              "Beispiel aus Annas Fahrt nennen können.")
    regel_streifen(s, [
        ("Reihenfolge", "Anforderungen vor Modell, Modell vor Relationen, Relationen vor DDL", ""),
        ("Kardinalität", "ist eine fachliche Aussage, keine technische Festlegung", ""),
        ("Normalform", "ist notwendig, aber nicht hinreichend", ""),
        ("Constraint", "was nicht erzwungen wird, wird verletzt", ""),
        ("Schutz", "liegt in der Datenbank, nie im Browser", ""),
        ("Nachweis", "Rechte werden geprüft, nicht angenommen", ""),
        ("Korn", "Eine Sicht hat genau ein Korn — wer es nicht benennt, summiert falsch", ""),
        ("Grenze", "Jede Kennzahl braucht die Angabe, was sie NICHT sagt", ""),
    ], y=180, hoehe=34, luecke=4, chip_b=0)
    notizen(s, "Zu jedem Satz gibt es eine Folie in dieser Einheit und eine Stelle im "
               "Modell. Das eignet sich als Prüfungsvorbereitung.")

    s = folie(prs, "12 · Ausblick", "Was entworfen ist, aber noch nicht gebaut",
              "Die Instandhaltung steht seit Kapitel 9. Beschaffung, Lager und die Logistik "
              "sind entworfen, aber nicht gebaut — grau hinterlegt im Diagramm.")
    oben = unter_intro(s)
    diagramm(s, bild("wawi-anschluss"), y=oben, hoehe=ZONE_UNTEN - 40 - oben)
    faden(s, "Meldet Anna einen Defekt, entsteht ein Wartungsauftrag — die Lagerbewegung dazu fehlt noch.")
    notizen(s, "Der Lehrpunkt beim Lager ist derselbe wie bei den Freiminuten: Bestand "
               "gegen Bewegung. Wer den Lagerbestand als Zahl pflegt, kann ihn nicht mehr "
               "erklären.")

    return prs


if __name__ == "__main__":
    prs = baue()
    prs.save(str(ZIEL))
    print(f"{len(prs.slides._sldIdLst)} Folien geschrieben nach {ZIEL}")
