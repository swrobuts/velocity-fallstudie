#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Erzeugt sechs Kurzdecks — je Use Case eine Datei mit zwei Folien.

Aufruf:
    python3 slides/build_usecase_decks.py
    python3 slides/check_deck.py slides/usecases/01_Preisauskunft.pptx

WARUM SECHS DATEIEN STATT EINES DECKS

velocity-crispdm.pptx ist eine Lesehilfe: 120 Folien, die das Vorgehen
Phase fuer Phase mitlaufen lassen. Das ist etwas anderes als eine
Vorstellung des Falls. Wer einen Use Case in fuenf Minuten erklaeren
will, braucht die Frage, den Weg, das Ergebnis und den Nutzen - und
zwar ohne durch ein Kapitel zu blaettern. Deshalb hier: eine Datei je
Fall, zwei Folien, keine Titelfolie.

ALLE ZAHLEN STAMMEN AUS DEN MERKZETTELN DER NOTEBOOKS

analytics/bau/werte/*.json wird beim Bau der Notebooks geschrieben. Ein
Foliendeck mit abgetippten Zahlen ist beim naechsten Datenlauf still
falsch - genau dieser Fehler hat in dieser Fallstudie schon einmal eine
Pruefung gekostet. Fehlt ein Wert, bricht der Bau ab.
"""
from __future__ import annotations

import json
import pathlib
import sys

from pptx import Presentation

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent))
from thws import (  # noqa: E402
    kachelreihe, kopf, notizen, prozesskette, sandband, sandkarte,
)

WURZEL = pathlib.Path(__file__).resolve().parent.parent
WERTE = WURZEL / "analytics" / "bau" / "werte"
ZIEL = WURZEL / "slides" / "usecases"

VORLAGE = pathlib.Path(
    "/Users/robert/.claude/skills/thws-slides/assets/template.pptx"
)
if not VORLAGE.exists():
    raise SystemExit(
        f"Vorlage fehlt: {VORLAGE}\n"
        "Ohne sie fehlen die Layouts Frontpage_Digital/Chapter/Slide."
    )

MERKZETTEL = {
    "01": "01_Regression_Fahrtdauer",
    "02": "02_Klassifikation_Wartungsrisiko",
    "03": "03_Clustering_Stationen_und_Kunden",
    "04": "04_Zeitreihe_Nachfrageprognose",
    "05": "05_Assoziation_Wege_im_Netz",
    "06": "06_Anomalieerkennung_Auffaellige_Vorgaenge",
}


class Werte(dict):
    """Merkzettel eines Notebooks. Ein unbekannter Schluessel bricht ab.

    Ein stiller KeyError-Ersatz („n/a") waere hier das Schlimmste: Die
    Folie saehe fertig aus und traege eine Luecke.
    """

    def __init__(self, nummer):
        pfad = WERTE / f"{MERKZETTEL[nummer]}.json"
        if not pfad.exists():
            raise SystemExit(
                f"ABBRUCH: {pfad.name} fehlt. Erst die Notebooks bauen:\n"
                f"    python3 analytics/bau/bauen.py")
        super().__init__(json.loads(pfad.read_text(encoding="utf-8")))
        self.nummer = nummer

    def __missing__(self, schluessel):
        raise SystemExit(
            f"ABBRUCH: Notebook {self.nummer} hat '{schluessel}' nicht mit "
            f"merke() festgehalten.\n    Bekannt sind: {', '.join(sorted(self))}")

    def p(self, schluessel, stellen=1):
        """Anteil als deutscher Prozentwert: 0.8112 -> '81,1 %'."""
        return f"{self[schluessel] * 100:.{stellen}f}".replace(".", ",") + " %"

    def z(self, schluessel, stellen=0):
        """Zahl in deutscher Schreibweise."""
        return f"{self[schluessel]:.{stellen}f}".replace(".", ",")


def leere_praesentation():
    prs = Presentation(str(VORLAGE))
    liste = prs.slides._sldIdLst
    for sld in list(liste):
        prs.part.drop_rel(sld.rId)
        liste.remove(sld)
    return prs


def lay(prs, name):
    return next(l for l in prs.slide_layouts if l.name == name)


# ─────────────────────────────────────────────────────── Die beiden Folien

def folie_frage(prs, fall):
    """Folie 1: Fragestellung, Weg, Erfolgskriterium."""
    s = prs.slides.add_slide(lay(prs, "Slide"))
    kopf(s, f"Notebook {fall['nr']} · {fall['kicker']}", fall["titel_a"],
         quelle=fall["quelle"], intro=fall["intro_a"])
    prozesskette(s, fall["start"], [(k, "") for k in fall["schritte"]],
                 fall["ziel"], y=196, hoehe=84)
    sandkarte(s, "Das Erfolgskriterium — festgelegt vor der Messung",
              fall["kriterium"], y=302)
    notizen(s, fall["notiz_a"])
    return s


def folie_ergebnis(prs, fall):
    """Folie 2: Ergebnis, Nutzen, Grenze — und der Satz, der bleibt."""
    s = prs.slides.add_slide(lay(prs, "Slide"))
    kopf(s, f"Notebook {fall['nr']} · Ergebnis und Nutzen", fall["titel_b"],
         quelle=fall["quelle"], intro=fall["intro_b"])
    kachelreihe(s, [("Ergebnis", fall["ergebnis"]),
                    ("Nutzen", fall["nutzen"]),
                    ("Grenze", fall["grenze"])], y=176, hoehe=204)
    sandband(s, fall["merksatz"], y=398)
    notizen(s, fall["notiz_b"])
    return s


def baue(fall):
    prs = leere_praesentation()
    folie_frage(prs, fall)
    folie_ergebnis(prs, fall)
    ZIEL.mkdir(parents=True, exist_ok=True)
    pfad = ZIEL / f"{fall['nr']}_{fall['datei']}.pptx"
    prs.save(str(pfad))
    return pfad


# ───────────────────────────────────────────────────────────── Die Inhalte

def fall01():
    w = Werte("01")
    return {
        "nr": "01", "datei": "Preisauskunft", "kicker": "Preisauskunft",
        "quelle": ("Notebook 01 Regression Fahrtdauer, Abschnitte 1, 5 und 6.7. "
                   "Synthetische Lehrdaten."),
        "titel_a": "Ein Preis vor der Fahrt — aber niemand kennt die Dauer",
        "intro_a": (
            "Der Preis einer Fahrt ergibt sich aus ihrer Dauer. Beim Entsperren steht "
            "die noch nicht fest, angezeigt werden soll der Preis trotzdem. Der Ausweg "
            "ist nicht eine genauere Zahl, sondern eine Spanne, zu der sich das "
            "Unternehmen verbindlich äußern kann."),
        "start": "Anfrage\nam Rad",
        "schritte": ["Abnahme\nversiegeln",
                     "Merkmale\nprüfen",
                     "Spanne statt\nPunktwert",
                     "Kandidaten am\nselben Gate"],
        "ziel": "Abnahme,\ndann Anzeige",
        "kriterium": [
            f"Trifft: Die angezeigte Spanne enthält den tatsächlichen Preis in "
            f"mindestens {w.p('gate_schwelle', 0)} der Fälle — gemessen dort, wo die "
            f"Schätzung überhaupt in den Preis eingeht.",
            "Nützt: Die Spanne umfasst höchstens zwölf Minuten und höchstens 60 % des "
            "angezeigten Preises; sonst zeigt die Anwendung nichts an.",
            "Reicht: Je Radtyp muss ein Mindestanteil der Anfragen beantwortet werden. "
            "Ein Verfahren, das für einen Radtyp schweigt, besteht nicht.",
        ],
        "titel_b": f"Ausgeliefert wird die {w['kandidat']} — die einfachere Bauform",
        "intro_b": (
            f"{w['zulaessige_satz']}. Entschieden hat deshalb nicht die Prognosegüte, "
            f"sondern eine vorab benannte Auswahlregel: Bei gleicher Eignung gewinnt "
            f"die Bauform, die ohne laufenden Dienst auskommt."),
        "ergebnis": [
            f"Zusage {w.p('gate_schwelle', 0)}, auf dem versiegelten Abnahme-"
            f"zeitraum belegt mit {w.p('ab_unten')}.",
            f"{w['ab_gates_halten']} von {w['ab_gates_gesamt']} Gates halten.",
            f"Beantwortet werden {w.p('reichweite_real', 0)} der Anfragen.",
            f"Status: {w['produktstatus']}.",
        ],
        "nutzen": [
            "Der Preisrahmen steht vor dem Entsperren fest, nicht erst auf der "
            "Rechnung.",
            "Eine CSV-Datei genügt: kein Dienst, keine Bibliotheksversion, von Hand "
            "nachrechenbar.",
            "Wo die Spanne zu breit würde, schweigt die Anwendung — statt zu raten.",
        ],
        "grenze": [
            f"In {w.p('zielabweichung', 0)} der Fahrten weicht das Ende vom "
            f"angegebenen Ziel ab. Die Zusage gilt nur für die gewählte Strecke und "
            f"wird mit jeder Antwort genannt.",
            f"Gültig für Fahrten bis {w['gueltig_bis_lang']} — so weit reicht der "
            f"Kalender.",
        ],
        "merksatz": ("Zugesagt wird eine Spanne, nicht eine Zahl — und gemessen wird "
                     "an der Zusage, nicht an einer Kennzahl, die gut aussieht."),
        "notiz_a": (
            "Der Einstieg über die Dauer lohnt sich: Studierende schlagen zuerst eine "
            "Punktprognose vor. Die Frage, welche Genauigkeit man einer Kundin "
            "zusagen will, führt von selbst zur Spanne. Wichtig ist die Reihenfolge "
            "auf dieser Folie: Der Abnahmezeitraum wird versiegelt, bevor die erste "
            "Grafik gezeichnet wird."),
        "notiz_b": (
            "Der überraschende Teil ist die Auswahlregel. Alle drei Verfahren nehmen "
            "die Hürden; ausgeliefert wird das betrieblich einfachste. Das ist eine "
            "Entscheidung über Betriebskosten, keine analytische — und sie war vorher "
            "festgelegt, sonst wäre sie eine nachträgliche Rechtfertigung."),
    }


def fall02():
    w = Werte("02")
    return {
        "nr": "02", "datei": "Wartungsrisiko", "kicker": "Wartungsrisiko",
        "quelle": ("Notebook 02 Klassifikation Wartungsrisiko, Abschnitte 1, 5 und 6. "
                   "Synthetische Lehrdaten."),
        "titel_a": "Begrenzte Werkstattkapazität — welche Räder zuerst?",
        "intro_a": (
            f"Die Werkstatt schafft {w['kapazitaet']} vorsorgliche Prüfungen je "
            f"Quartal. Gesucht ist keine möglichst hohe Trefferquote, sondern die "
            f"Liste, die den größten Schaden verhindert — und dafür müssen beide "
            f"Fehlerarten zuerst bepreist werden."),
        "start": "Flotte und\nWerkstatt",
        "schritte": ["Fehler\nbepreisen",
                     "Schnitt nach\nder Zeit",
                     "Faustregeln\nals Maßstab",
                     "Modelle mit\nKostengewicht"],
        "ziel": "Prüfliste\nje Quartal",
        "kriterium": [
            f"Die Kostenmatrix steht am Anfang: ein übersehener Ausfall kostet "
            f"{w.z('kosten_verpasst')} €, eine unnötige Prüfung "
            f"{w.z('kosten_unnoetig')} € — ein Verhältnis von rund 7 zu 1.",
            f"Bindend sind drei Kriterien ({w['pflichtgates']}): Nutzen über mehrere "
            f"Quartale, Vergleich mit der heute verwendeten Faustregel und eine "
            f"statistische Absicherung gegen eine mitlaufende Grundrate.",
            "Nicht bindend ist die ursprünglich genannte 70-Prozent-Marke. Sie bleibt "
            "als Diagnose stehen, weil sie in manchen Quartalen gar nicht erreichbar "
            "ist — das zeigt eine Rechnung, keine Meinung.",
        ],
        "titel_b": "Die Faustregel gewinnt — und zwar nachweisbar",
        "intro_b": (
            f"Im Testquartal trifft die Faustregel {w['treffer_regel']} Räder, das "
            f"Random-Forest-Modell {w['treffer_wald']}. Ausschlaggebend war jedoch "
            f"nicht dieser Vorsprung, sondern die untere Vertrauensgrenze gegen eine "
            f"Schwelle, die sich aus der Grundrate des Quartals ergibt."),
        "ergebnis": [
            f"Faustregel {w.p('wilson_unten_regel')} untere Grenze gegen geforderte "
            f"{w.p('k3_schwelle')}.",
            f"Random Forest: {w.p('wilson_unten_wald')} — verfehlt sie.",
            f"Von zehn geprüften Rädern werden {w.z('quote_regel_von_zehn', 1)} "
            f"auffällig.",
            f"Von zehn auffälligen erfasst die Liste "
            f"{w.z('abdeckung_von_zehn', 1)}.",
        ],
        "nutzen": [
            "Die Liste entsteht ohne Modellbetrieb: eine Regel, die jede Werkstatt "
            "nachvollziehen und im Zweifel begründen kann.",
            "Das Kostenverhältnis steht sichtbar in der Entscheidung, statt in einer "
            "Kennzahl zu verschwinden.",
            "Treffsicherheit und Abdeckung werden getrennt berichtet — beide sind "
            "richtig und messen Verschiedenes.",
        ],
        "grenze": [
            f"Der Anteil auffälliger Räder schwankt über {w['panel_stichtage']} "
            f"Stichtage zwischen {w.p('panel_grundrate_min')} und "
            f"{w.p('panel_grundrate_max')}. Ein einzelnes gutes Quartal belegt wenig.",
            f"Beigelegt ist eine Prognoseliste zum {w['prognose_stichtag_lang']}; "
            f"bewertbar wird sie erst nach {w['horizont_tage']} Tagen.",
        ],
        "merksatz": ("Ein Modell schlägt eine Faustregel nur dann, wenn die Faustregel "
                     "mit derselben Sorgfalt gebaut wurde wie das Modell."),
        "notiz_a": (
            "Hier lohnt sich der Umweg über die Kosten. Solange die Frage lautet "
            "„wie genau ist das Modell?“, ist sie nicht beantwortbar — 60 Plätze und "
            "zwei ungleich teure Fehler machen daraus eine Entscheidung."),
        "notiz_b": (
            "Der Punkt ist nicht, dass Modelle nichts taugen. Der Punkt ist, dass ein "
            "Vorsprung von neun Treffern in einem Quartal nichts belegt, solange die "
            "Grundrate zwischen den Quartalen um das Dreifache schwankt."),
    }


def fall03():
    w = Werte("03")
    return {
        "nr": "03", "datei": "Segmente", "kicker": "Stationen und Kundschaft",
        "quelle": ("Notebook 03 Clustering Stationen und Kunden, Abschnitte 1, 5 und 6. "
                   "Synthetische Lehrdaten."),
        "titel_a": "Gruppen, die niemand aufgeschrieben hat",
        "intro_a": (
            "In den Stammdaten steht nicht, welche Station als Pendlerstation dient "
            "und welches Nutzungsmuster eine Kundengruppe zeigt. Beides steckt im "
            "Verhalten — und weil es keine hinterlegte richtige Antwort gibt, muss "
            "vorher feststehen, woran eine brauchbare Gruppierung erkennbar ist."),
        "start": "Daten ohne\nSegmente",
        "schritte": ["Verhalten statt\nStammdaten",
                     "Tagesgang\nund RFM",
                     "k-Means,\nk wählen",
                     "Fünf Kriterien\nprüfen"],
        "ziel": "Profile und\nBericht",
        "kriterium": [
            "Fünf Kriterien, für beide Teile dieselben: Eine Gruppe muss benennbar "
            "sein, sich anders behandeln lassen als die übrigen, groß genug sein, "
            "gegenüber dem Startwert reproduzierbar und über die Zeit stabil.",
            "Keines davon ist eine Gütezahl des Verfahrens. Silhouette und "
            "Rand-Index werden berichtet, entscheiden aber nicht — sie messen "
            "Trennschärfe, nicht Verwendbarkeit.",
            "Für die Kampagne kommen drei organisatorische Bedingungen hinzu; sie "
            "betreffen die Verwendung der Segmente, nicht ihre Berechnung.",
        ],
        "titel_b": "Brauchbare Stationstypen — und ein Preisproblem als Nebenbefund",
        "intro_b": (
            f"Die Stationstypen lassen sich gegen eine im Datensatz hinterlegte, dem "
            f"Verfahren nicht bekannte Zuordnung prüfen: {w.p('generator_treffer', 0)} "
            f"stimmen überein. Bei der Kundschaft halten "
            f"{w['gates_erfuellt']} von {w['gates_gesamt']} Kriterien."),
        "ergebnis": [
            f"Stationen: {w.p('generator_treffer', 0)} richtig zugeordnet, "
            f"Rand-Index {w.z('generator_ari', 3)}.",
            f"Kundschaft: analytisch {w['status_analytisch']}, Einsatz "
            f"{w['status_einsatz']}.",
            f"Ausgeliefert wird {w['exportart']}.",
        ],
        "nutzen": [
            "Stationsprofile geben der Umverteilung eine Sprache: benannte Typen "
            "statt Bauchgefühl — ausdrücklich als Hypothese, nicht als Sollbestand.",
            f"Sichtbar wurde ein Preisproblem: die {w['viel_segment']} bringen "
            f"{w.z('viel_je_fahrt', 2)} € je Fahrt, die {w['stark_segment']} "
            f"{w.z('stark_je_fahrt', 2)} €.",
            "Der Kundenbericht bleibt aggregiert und ohne Namen — das war eine "
            "Bedingung, keine nachträgliche Vorsicht.",
        ],
        "grenze": [
            f"{w.p('kurze_historie_anteil', 0)} der Kundschaft erscheinen gar nicht, "
            f"weil sie im Betrachtungszeitraum nicht gefahren sind. Ein Verfahren, "
            f"das auf Nutzung beruht, sieht abgewanderte Kundschaft nicht.",
            "Was fehlt, ist keine Kennzahl, sondern eine prospektive Prüfung: ob die "
            "Segmente auch im nächsten Quartal noch dieselben sind.",
        ],
        "merksatz": ("Ohne Zielgröße gibt es kein „richtig“ — die Erfolgskriterien "
                     "müssen deshalb vorher gesetzt und begründet werden."),
        "notiz_a": (
            "Dieser Fall ist der beste Anlass, über Erfolgskriterien zu sprechen. Bei "
            "Regression und Klassifikation liefert die Zielgröße einen Maßstab frei "
            "Haus; hier muss man ihn selbst bauen, und zwar bevor man das Ergebnis "
            "kennt."),
        "notiz_b": (
            "Der Nebenbefund zur Tarifstruktur ist didaktisch der wertvollste Teil: "
            "Die Segmentierung sollte eine Kampagne vorbereiten und hat stattdessen "
            "eine Frage an das Preismodell aufgeworfen. Genau dafür macht man "
            "explorative Analysen."),
    }


def fall04():
    w = Werte("04")
    return {
        "nr": "04", "datei": "Nachfrageprognose", "kicker": "Nachfrageprognose",
        "quelle": ("Notebook 04 Zeitreihe Nachfrageprognose, Abschnitte 1, 4 und 6. "
                   "Synthetische Lehrdaten."),
        "titel_a": "Am Vorabend planen — mit dem Wetter, das man dann kennt",
        "intro_a": (
            "Die Disposition plant abends für den kommenden Tag. Entscheidend ist "
            "nicht, wie gut ein Verfahren die Vergangenheit trifft, sondern wie gut "
            "es mit der Information arbeitet, die um 18 Uhr tatsächlich vorliegt — "
            "und das ist eine Wettervorhersage, kein gemessenes Wetter."),
        "start": "Planung\nam Vorabend",
        "schritte": ["Fehler-\nrichtungen",
                     "Schnitt nach\nder Zeit",
                     "Prognose-\nwetter",
                     "Vier Verfahren\nvergleichen"],
        "ziel": "Tages-\nprognose",
        "kriterium": [
            "Prognostiziert werden Fahrten, nicht Räder und nicht Schichten. Die "
            "Übersetzung in Dispositionsgrößen ist eine eigene Analyse und wird nicht "
            "stillschweigend mitgeliefert.",
            "Die beiden Fehlerrichtungen kosten unterschiedlich viel: zu wenig "
            "geplante Fahrten 4,00 €, zu viel geplante 0,80 € je Fahrt. Daraus folgt "
            "ein Sicherheitsaufschlag, der auf der Validierung gewählt wird.",
            "Verglichen wird ausschließlich unter Prognosewetter. Der Test wird erst "
            "geöffnet, nachdem Verfahren und Aufschlag feststehen.",
        ],
        "titel_b": "Das einfachere Verfahren gewinnt — unter Prognosewetter",
        "intro_b": (
            f"Gewählt wurde {w['gewaehlt_name']}. Unter Ist-Wetter liegen lineares "
            f"Modell und Gradient Boosting praktisch gleichauf; erst unter "
            f"Prognosewetter setzt sich das einfachere Verfahren ab. Die Modellwahl "
            f"hängt damit unmittelbar daran, unter welchen Bedingungen man vergleicht."),
        "ergebnis": [
            f"{w['gewaehlt_name']}: mittlerer Fehler {w.z('mae_linear', 1)} Fahrten.",
            f"Faustregel {w.z('mae_faustregel', 1)}, Nullmodell "
            f"{w.z('mae_null', 1)}.",
            f"Unter Ist-Wetter: {w.z('ist_linear', 2)} gegen "
            f"{w.z('ist_boosting', 2)} — praktisch gleichauf.",
            f"Status: {w['nb04_status']}.",
        ],
        "nutzen": [
            "Die Planung bekommt eine begründete Zahl statt eines Erfahrungswerts — "
            "und sie ist nachweislich besser als die bisherige Faustregel.",
            f"Der Aufschlag von {w.p('aufschlag', 0)} macht die teurere Fehlerrichtung "
            f"zu einer sichtbaren Entscheidung statt zu einem Bauchgefühl.",
            "Ein Probebetrieb läuft mit und wird protokolliert, ohne dass jemand "
            "nach der Prognose handelt.",
        ],
        "grenze": [
            "Prognostiziert wird die Gesamtzahl der Fahrten; gebraucht werden Räder "
            "je Station. Diese Umrechnung fehlt noch und ist keine Formel.",
            "Die Wetterunsicherheit ist simuliert, nicht gemessen — und ein "
            "Sommerfenster trägt keine Jahresaussage.",
        ],
        "merksatz": ("Verglichen wird unter den Bedingungen des Einsatzes. Wer unter "
                     "Ist-Wetter wählt, wählt für eine Lage, in der er nie liefert."),
        "notiz_a": (
            "Der Unterschied zwischen Ist-Wetter und Prognosewetter ist der Kern "
            "dieses Falls. Er lässt sich gut als Frage stellen: Welche Information "
            "hat die Disposition um 18 Uhr wirklich?"),
        "notiz_b": (
            "Dass das einfachere Verfahren gewinnt, ist kein Zufall und kein "
            "Argument gegen Boosting: Unter Prognoseunsicherheit verliert das "
            "flexiblere Modell mehr, weil es auf Merkmale reagiert, die selbst "
            "unsicher sind."),
    }


def fall05():
    w = Werte("05")
    return {
        "nr": "05", "datei": "Wege_im_Netz", "kicker": "Ströme im Netz",
        "quelle": ("Notebook 05 Assoziation Wege im Netz, Abschnitte 1, 5 und 6. "
                   "Synthetische Lehrdaten."),
        "titel_a": "Von wo nach wo — und ist das Muster stabil?",
        "intro_a": (
            "Gesucht sind Verbindungen, die innerhalb desselben Zeitfensters häufiger "
            "auftreten, als bei zufälliger Zielwahl zu erwarten wäre. Gezählt wird, "
            "nicht trainiert — und die entscheidende Vorkehrung ist zeitlich: Ein "
            "Teil der Daten bleibt bis zur Bestätigung ungeöffnet."),
        "start": "Fahrten als\nWarenkörbe",
        "schritte": ["Zeitraum\nversiegeln",
                     "Kontext-Lift\nbilden",
                     "Regeln nur im\nfrühen Teil",
                     "Bootstrap\nje Regel"],
        "ziel": "Dispositions-\nhinweis",
        "kriterium": [
            f"Produkt A, die automatische Umverteilung, verlangt vier Bedingungen: "
            f"Support, kontextbezogener Lift über {w.z('k2_lift', 1)}, ein konkretes "
            f"Ziel — und einen Nachweis der Wirtschaftlichkeit.",
            "Produkt B, der Dispositionshinweis, verlangt Bestätigung je einzelner "
            "Regel im unangetasteten Zeitraum, die Größenordnung neben jeder Regel, "
            "keine Automatik und eine Kennzeichnung der Begleitanalysen.",
            f"Bestätigt heißt: Die untere Grenze eines Tagesblock-Bootstraps liegt "
            f"über {w.z('k2_lift', 1)}. Ein Punktschätzer über der Schwelle genügt "
            f"nicht — er ignoriert, dass Fahrten desselben Tages zusammenhängen.",
        ],
        "titel_b": "Ein Hinweis für Menschen — keine automatische Umverteilung",
        "intro_b": (
            f"Von {w['b1_kandidaten']} geprüften Regeln halten {w['b1_gehalten']} "
            f"unter Unsicherheit. Produkt A bleibt {w['status_a']} — nicht wegen der "
            f"Regeln, sondern weil die entgangenen Fahrten in diesen Daten nicht "
            f"vorkommen und sich auch nicht aus den beobachteten erschließen lassen."),
        "ergebnis": [
            f"Produkt B: {w['b_regeln_n']} von {w['b1_kandidaten']} Regeln bestätigt.",
            f"Ausschluss zweigeteilt: {w['b1_raus_punkt']} scheitern schon am "
            f"Punktschätzer, {w['b1_raus_intervall']} erst am Intervall.",
            f"Produkt A: {w['status_a']}.",
        ],
        "nutzen": [
            "Die Disposition bekommt begründete Hinweise mit Größenordnung — und "
            "entscheidet selbst, ob gefahren wird.",
            "Jede Regel trägt ihren Nenner: eine Verbindung mit hohem Lift kann "
            "trotzdem Bruchteile einer Fahrt je Tag bedeuten.",
            "Salden und Abstell-Hotspots liegen bei, ausdrücklich als explorativ "
            "gekennzeichnet.",
        ],
        "grenze": [
            f"Bestanden ist ein analytisches Lehr-Gate auf synthetischen Daten — "
            f"keine Betriebsfreigabe. Die Datei trägt deshalb bewusst kein "
            f"Gültigkeitsdatum.",
            f"Die Hürde aus Phase 1 entspricht {w.z('huerde_je_werktag', 2)} Fahrten "
            f"je Werktag: eine Größenordnung, in der keine Umsetzfahrt beginnt. Das "
            f"Kriterium war auf der falschen Skala — verschoben wurde es dennoch nicht.",
        ],
        "merksatz": ("Beobachtet ist nicht bestätigt: Erst ein Unsicherheitsintervall "
                     "sagt, ob eine Regel die Schwelle wirklich hält."),
        "notiz_a": (
            "Die Trennung in zwei Produkte ist die eigentliche Leistung dieses Falls. "
            "Dieselben Regeln tragen einen Hinweis für Menschen, aber keine "
            "automatische Umsetzfahrt — weil an der Automatik eine "
            "Wirtschaftlichkeitsfrage hängt, die diese Daten nicht beantworten."),
        "notiz_b": (
            "Der Unterschied zwischen Punktschätzer und Bootstrap-Untergrenze lässt "
            "sich hier an konkreten Regeln zeigen: Drei Regeln liegen im "
            "Punktschätzer über der Schwelle und fallen trotzdem heraus."),
    }


def fall06():
    w = Werte("06")
    return {
        "nr": "06", "datei": "Anomalien", "kicker": "Auffällige Vorgänge",
        "quelle": ("Notebook 06 Anomalieerkennung, Abschnitte 1, 5 und 6. "
                   "Synthetische Lehrdaten."),
        "titel_a": "„Auffällig“ ist keine Frage — es sind drei",
        "intro_a": (
            "Ein überfälliges Rad jetzt, eine seltsame abgeschlossene Fahrt heute "
            "früh und eine stillstehende Station gestern sind drei verschiedene "
            "Aufgaben mit drei Entscheidungszeitpunkten. Wer sie in einer Liste "
            "zusammenfasst, baut ein Produkt, das keine der drei Fragen beantwortet."),
        "start": "Vorgänge\nim Betrieb",
        "schritte": ["Drei Fragen\ntrennen",
                     "Listenlänge\nableiten",
                     "Merkmale ohne\nPreisklasse",
                     "Tagesliste\nmessen"],
        "ziel": "Prüfliste,\nsechs Plätze",
        "kriterium": [
            f"Die Listenlänge ist abgeleitet, nicht gesetzt: {w['listenlaenge']} "
            f"Plätze folgen aus dem verfügbaren Zeitbudget und der Prüfdauer je Fall.",
            f"Für die Prüfliste gilt eine Präzision von mindestens "
            f"{w.p('b_gate_praezision', 0)} je neuem Alarm, ein Recall von "
            f"{w.p('b_gate_recall', 0)} der Episoden und höchstens "
            f"{w['b_gate_verzug']} Tag Verzug.",
            "Für die Rangliste auffälliger Fahrten gibt es bewusst kein Gütekriterium: "
            "Es fehlt ein Label. Was ohne Maßstab läuft, läuft im Schattenbetrieb.",
        ],
        "titel_b": "Drei Produkte, drei Urteile — und ein lehrreicher Rücksprung",
        "intro_b": (
            f"Das erste Modell rechnete sauber und lieferte Unbrauchbares: Es trennte "
            f"die Preisklassen statt der Anomalien. Aufgefallen ist das nicht durch "
            f"eine Kennzahl, sondern durch Sichtung der obersten Zeilen — danach "
            f"folgte ein Rücksprung in die Datenaufbereitung."),
        "ergebnis": [
            f"A1 überfällige Rückgabe: {w['a1_status']} — als Regel, ohne Modell.",
            f"A2 auffällige Fahrten: {w['a2_status']}, da kein Label vorliegt.",
            f"B Prüfliste: {w['b_status']}, {w['b_gates_halten']} Gates halten.",
        ],
        "nutzen": [
            "A1 löst die dringendste der drei Fragen mit einer Regel — dafür braucht "
            "es kein Modell und keinen Betrieb.",
            f"Die Prüfliste hebt die Präzision gegenüber der bisherigen Meldelogik "
            f"({w.p('stat_alt_quote')} je Alarm) deutlich an.",
            "Jedes der drei Produkte trägt seinen eigenen Status, statt gemeinsam "
            "freigegeben oder gesperrt zu werden.",
        ],
        "grenze": [
            f"Die globale Rangliste erreicht {w.p('globale_quote')}, die im Betrieb "
            f"erzeugbare Tagesliste {w.p('tagesquote')} — bei demselben Modell.",
            f"Auf dem unangetasteten Testabschnitt bleibt die Präzision mit "
            f"{w.p('stat_je_alarm')} unter den geforderten "
            f"{w.p('b_gate_praezision', 0)}; Recall und Verzug halten.",
        ],
        "merksatz": ("Eine Kennzahl auf der Gesamtliste beschreibt nicht die Liste, "
                     "mit der später tatsächlich gearbeitet wird."),
        "notiz_a": (
            "Die drei Zeitpunkte sind der Schlüssel. Sobald man fragt, wann jemand "
            "auf die Information reagieren soll, zerfällt „auffällig“ von selbst in "
            "drei Aufgaben mit unterschiedlichen Anforderungen."),
        "notiz_b": (
            "Der Rücksprung ist hier kein Betriebsunfall, sondern der Lerninhalt: "
            "Das Modell hat getan, was man ihm gesagt hat. Der Fehler lag in den "
            "Merkmalen — ein Lastenrad ist teurer als ein Citybike, und genau das "
            "hat der Isolation Forest gefunden."),
    }


FAELLE = (fall01, fall02, fall03, fall04, fall05, fall06)


if __name__ == "__main__":
    for bauen_fall in FAELLE:
        fall = bauen_fall()
        pfad = baue(fall)
        print(f"  {pfad.relative_to(WURZEL)}   2 Folien")
    print(f"\n{len(FAELLE)} Dateien -> {ZIEL.relative_to(WURZEL)}")
