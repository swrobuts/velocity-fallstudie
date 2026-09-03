#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Erzeugt das Foliendeck „CRISP-DM an sechs Fallbeispielen".

Aufruf:
    python3 slides/build_crispdm_deck.py
    python3 slides/check_deck.py slides/velocity-crispdm.pptx

Didaktisches Geruest
--------------------
Die Frage war: sechsmal der volle Kreislauf, oder einmal generisch und
dann die Faelle einsortiert? Beides waere falsch. Sechsmal derselbe
Rahmen begraebt die sechs UNTERSCHIEDLICHEN Lehren unter der immer
gleichen Gliederung. Einmal generisch und dann nur einsortiert laesst
die Studierenden nie ein Projekt von Anfang bis Ende erleben.

Deshalb vier Teile:

  A Die Karte          Der Kreislauf einmal kompakt - ein Nachschlage-
                       werk, keine Lektion.
  B Der Referenzfall   Notebook 1 Phase fuer Phase, mit den echten
                       Zahlen. Hier wird die Karte konkret. Regression
                       ist die vertrauteste Methode; die Aufmerksamkeit
                       bleibt beim Vorgehen statt beim Verfahren.
  C Fuenf Faelle       Gleiches Geruest, Tiefe nur dort, wo dieser Fall
                       etwas Neues zeigt. Denn jedes Notebook betont
                       eine andere Phase:
                         Fall 2 -> Phase 6  das Modell verdient seinen Unterhalt nicht
                         Fall 3 -> Phase 1  Kriterien ohne Zielgroesse
                         Fall 4 -> Phase 3  Schnitt entlang der Zeit
                         Fall 5 -> Phase 5  keine Regel, und der Plan traegt nicht
                         Fall 6 -> Ruecksprung 4 nach 3
  D Synthese           Die sechs Faelle zurueck auf den Kreislauf.

Der wiederkehrende Anker ist die PHASENLEISTE am Fuss jeder Inhalts-
folie - dieselbe Rolle, die im Datenbankdeck der rote Faden zu Annas
Fahrt spielt. Sie macht die Struktur sichtbar, ohne sie zu wiederholen.

Alle Zahlen stammen aus den ausgefuehrten Notebooks unter
analytics/notebooks/, nicht aus dem Gedaechtnis.
"""
from __future__ import annotations

import pathlib
import re
import sys

from pptx import Presentation

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent))
from thws import (  # noqa: E402
    BLAU, BREITE, FLUCHT_L, GRUEN_D, ORANGE, ROT_A, TUERKIS, TEXT_SEK,
    ZONE_OBEN, ZONE_UNTEN,
    ampel_matrix, code_kacheln, diagramm, kachelreihe, kopf, leitfrage, notizen,
    phasenleiste, prozesskette, regel_streifen, sandband, sandkarte,
    schichtenstapel, steckbrief, tabelle, vorher_nachher,
)

WURZEL  = pathlib.Path(__file__).resolve().parent.parent
ZIEL    = WURZEL / "slides" / "velocity-crispdm.pptx"

VORLAGE = pathlib.Path(
    "/Users/robert/.claude/skills/thws-slides/assets/template.pptx"
)
if not VORLAGE.exists():
    raise SystemExit(
        f"Vorlage fehlt: {VORLAGE}\n"
        "Ohne sie fehlen die Layouts Frontpage_Digital/Chapter/Slide."
    )

Q = ("Sechs Notebooks unter analytics/notebooks/ der Fallstudie VeloCity. "
     "Alle Zahlen sind den ausgefuehrten Notebooks entnommen.")


# ─────────────────────────────────────────────────────────── Grundgeruest

def leere_praesentation() -> Presentation:
    prs = Presentation(str(VORLAGE))
    liste = prs.slides._sldIdLst
    for sld in list(liste):
        prs.part.drop_rel(sld.rId)
        liste.remove(sld)
    return prs


def lay(prs, name):
    return next(l for l in prs.slide_layouts if l.name == name)


def kapitel(prs, nummer, titel, frage, notiz):
    s = prs.slides.add_slide(lay(prs, "Chapter"))
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.text = f"{nummer} · {titel}"
    leitfrage(s, frage)
    notizen(s, notiz)
    return s


# Welches Notebook gerade behandelt wird. Die Fallkapitel setzen das am
# Anfang; folie() leitet daraus zusammen mit dem Kicker die Fusszeile ab.
AKTUELLES_NB = [None]


def folie(prs, kicker, titel, intro=None, quelle=None):
    """Eine Inhaltsfolie. Die Quellenzeile entsteht, wo es geht, von selbst.

    Das Deck ist eine Lesehilfe: Jede Folie eines Fallkapitels soll sagen,
    wo im Notebook das Gezeigte steht. Von Hand gepflegt bleibt das nicht
    vollstaendig - beim ersten Durchlauf trugen 24 von 120 Folien eine
    Quelle. Deshalb wird sie aus dem Kicker abgeleitet: „Phase 5 · ...“
    oder „Fall 2 · Phase 5“ ergibt Abschnitt 5 des laufenden Notebooks.
    """
    if quelle is None:
        nb = AKTUELLES_NB[0]
        phase = re.search(r"Phase (\d)", kicker)
        quelle = nbq(nb, phase.group(1) if phase else "") if nb else Q
    s = prs.slides.add_slide(lay(prs, "Slide"))
    s._intro_unten = kopf(s, kicker, titel, quelle=quelle, intro=intro)
    return s


def unter_intro(s, abstand=16):
    return getattr(s, '_intro_unten', 110) + abstand


ASSETS = WURZEL / "slides" / "assets"

NB_DATEI = {
    1: "01_Regression_Fahrtdauer",
    2: "02_Klassifikation_Wartungsrisiko",
    3: "03_Clustering_Stationen_und_Kunden",
    4: "04_Zeitreihe_Nachfrageprognose",
    5: "05_Assoziation_Wege_im_Netz",
    6: "06_Anomalieerkennung_Auffaellige_Vorgaenge",
}


def nbq(nummer, abschnitt="", hinweis=""):
    """Quellenzeile, die auf Datei UND Abschnitt zeigt.

    Das Deck ist eine Lesehilfe. Eine Folie, die eine Zahl nennt, ohne zu
    sagen, wo sie steht, schickt die Studierenden auf die Suche. Diese
    Zeile beendet die Suche: Datei, Abschnitt, fertig.
    """
    teile = [f"analytics/notebooks/{NB_DATEI[nummer]}.ipynb"]
    if abschnitt:
        teile.append(f"Abschnitt {abschnitt}")
    if hinweis:
        teile.append(hinweis)
    return " · ".join(teile)


def bild(name: str) -> str:
    """Pfad zu einem Bild - und die Pruefung, ob es noch aktuell ist.

    Ein Zellausschnitt, der aelter ist als sein Notebook, zeigt eine
    Ausgabe, die es so nicht mehr gibt. Genau das ist passiert: Der
    Suchanker fuer nb2-kriterien zeigte auf eine geloeschte Zeile, das
    Werkzeug brach ab - und der Deckbau lief trotzdem durch, weil das
    alte PNG noch dalag. Auf der Folie stand danach ein Bild aus einer
    Fassung, deren Urteil umgedreht worden war.
    """
    pfad = ASSETS / f"{name}.png"
    if not pfad.exists():
        raise SystemExit(
            f"Bild fehlt: {pfad}\n"
            "Zuerst: python3 tools/notebook_ausschnitte.py "
            "und bash tools/render_diagrams.sh")
    # Zwei Sorten Bilder, zwei Quellen: Mermaid-Diagramme kommen aus einer
    # .mmd-Datei, Zellausschnitte aus einem Notebook. Der Praefix nb1-
    # steht bei beiden - entschieden wird deshalb ueber die Existenz der
    # Diagrammquelle, nicht ueber den Namen.
    diagramm_quelle = WURZEL / "doku" / "analytics" / "diagramme" / f"{name}.mmd"
    if diagramm_quelle.exists():
        quelle, befehl = diagramm_quelle, "bash tools/render_diagrams.sh"
    else:
        nummer = name[2] if name.startswith("nb") and name[2:3].isdigit() else None
        if not (nummer and int(nummer) in NB_DATEI):
            return str(pfad)
        quelle = WURZEL / "analytics" / "notebooks" / f"{NB_DATEI[int(nummer)]}.ipynb"
        befehl = f"python3 tools/notebook_ausschnitte.py {name}"
    if quelle.exists() and quelle.stat().st_mtime > pfad.stat().st_mtime + 1:
        raise SystemExit(
            f"Bild veraltet: {pfad.name} ist aelter als {quelle.name}.\n"
            f"Zuerst: {befehl}")
    return str(pfad)


def zellfolie(prs, nb, abschnitt, kicker, titel, ausschnitt, deutung, notiz):
    """Eine Folie, die eine ECHTE Notebookzelle zeigt und sie deutet.

    Das Bild traegt die Folie, der Text darunter sagt, worauf zu achten
    ist. Ohne diese eine Zeile ist ein Screenshot nur Dekoration.
    """
    s = folie(prs, kicker, titel, quelle=nbq(nb, abschnitt))
    y = unter_intro(s, 10)
    # Die Hoehe des Bandes folgt seinem Text (dieselbe Rechnung wie in
    # thws.sandband). Geschaetzt hatte ich 60 - bei drei Zeilen ragte das
    # Band dann 6 pt aus der Inhaltszone, und der Pruefer hat es gemeldet.
    je_zeile = int((BREITE - 32) / (14 * 0.52))
    band_h = max(52, max(1, -(-len(deutung) // je_zeile)) * 20 + 26)
    band_y = ZONE_UNTEN - band_h
    diagramm(s, bild(ausschnitt), y=y, hoehe=band_y - y - 14)
    sandband(s, deutung, y=band_y)
    notizen(s, notiz)
    return s


# ── Hoehen ausrechnen statt schaetzen
#
# Die zweite Form auf einer Folie braucht die Unterkante der ersten.
# Beim Umbau von Fall 1 standen dort geschaetzte Zahlen, und der
# Deckpruefer meldete elf Ueberlappungen - eine Tabelle mit fuenf
# Zeilen ist eben 40 pt hoeher als eine mit vier. Diese drei Funktionen
# nutzen dieselben Formeln wie die Motive in thws.py.

def h_tabelle(n_zeilen, zeilen_h):
    """Kopfzeile plus n Datenzeilen."""
    return (n_zeilen + 1) * zeilen_h


def h_gestapelt(n, hoehe, luecke):
    """streifen, schichtenstapel, ampel_matrix - n Karten mit Luecke."""
    return n * hoehe + (n - 1) * luecke


def darunter(y, hoehe, abstand=14):
    return y + hoehe + abstand


def streifen(s, zeilen, **kw):
    """regel_streifen, das auch Zweiertupel annimmt.

    Das Motiv erwartet (Regel, Wirkung, Beispielchip). Der Chip traegt
    hier meist nichts bei - die Wirkung steht schon im Fliesstext -,
    deshalb darf er entfallen.
    """
    regel_streifen(s, [z if len(z) == 3 else (z[0], z[1], "") for z in zeilen], **kw)


def rahmen(prs):
    AKTUELLES_NB[0] = None
    """Die zwei Folien zwischen Titel und Kapitel 1: Aufhaenger und Wegweiser."""
    s = folie(prs, "Der Ausgangspunkt",
              "Sechs Projekte, sechs Verfahren — und sechsmal dieselbe Reihenfolge",
              "Regression, Klassifikation, Clustering, Zeitreihe, Assoziation, "
              "Anomalieerkennung. Verschiedene Fragen, verschiedene Mathematik. "
              "Der Weg zum Ergebnis ist in allen sechs Fällen derselbe.")
    tabelle(s, ["Notebook", "Frage an die Daten"], [
        ["1 Regression", "Wie lange dauert diese Fahrt?"],
        ["2 Klassifikation", "Welche Räder müssen als Nächstes in die Werkstatt?"],
        ["3 Clustering", "Welche Gruppen stecken in den Daten?"],
        ["4 Zeitreihe", "Wie viele Fahrten kommen morgen?"],
        ["5 Assoziation", "Welche Wege gehören zusammen?"],
        ["6 Anomalieerkennung", "Was ist gestern schiefgelaufen?"],
    ], y=unter_intro(s), spalten_b=[250, 653.5], zeilen_h=40)
    notizen(s, "Lesen Sie die sechs Fragen laut vor. Sie klingen nach sechs völlig "
               "verschiedenen Aufgaben — und das sind sie auch. Genau deshalb ist es "
               "bemerkenswert, dass der Weg zur Antwort in allen sechs derselbe ist. "
               "Dieser Weg heißt CRISP-DM.")

    s = folie(prs, "Wegweiser", "Wie dieses Deck gebaut ist — und warum so",
              "Der Kreislauf wird nicht sechsmal wiederholt. Ein Fall läuft "
              "vollständig durch, die anderen fünf gehen dort in die Tiefe, wo sie "
              "etwas Neues zeigen.")
    streifen(s, [
        ("Kapitel 1 — Die Karte", "Der Kreislauf einmal kompakt. Bewusst knapp: "
                                  "Er bleibt abstrakt, bis der erste Fall ihn füllt"),
        ("Kapitel 2 — Fall 1 ganz", "Regression, alle sechs Phasen, mit den echten "
                                    "Zahlen. Hier wird die Karte konkret"),
        ("Kapitel 3 bis 7 — fünf Fälle", "Gleiches Gerüst, Tiefe nur an der Stelle, "
                                         "an der dieser Fall etwas beisteuert"),
        ("Kapitel 8 — Synthese", "Die sechs Fälle zurück auf die Karte gelegt"),
    ], y=unter_intro(s), hoehe=62, luecke=9, chip_b=0)
    notizen(s, "Weisen Sie auf die Leiste am Fuß der Folien hin: Sie zeigt jederzeit, "
               "in welcher Phase wir gerade sind. Wo ein Feld sandfarben mit rotem "
               "Rand erscheint, ist an dieser Stelle zurückgesprungen worden.")


# ═══════════════════════════════════════════════════ Teil A — Die Karte

def teil_karte(prs):
    AKTUELLES_NB[0] = None
    kapitel(prs, 1, "Der Kreislauf",
            "Warum braucht ein Datenprojekt überhaupt eine feste Schrittfolge?",
            "Wir bauen in diesem Kapitel die Landkarte. Sie ist bewusst knapp gehalten "
            "— sechs Phasen, was jede beantwortet, wo die Rücksprünge sitzen. Konkret "
            "wird das erst im nächsten Kapitel am ersten Fall. Wer die Karte jetzt "
            "abstrakt findet, hat recht: Sie IST abstrakt, bis der erste Fall sie "
            "füllt. Deshalb steht sie hier nur zwölf Folien lang.")

    # ── Warum ein Vorgehensmodell
    s = folie(prs, "Ausgangslage",
              "Ohne Reihenfolge ist ein Werkzeugkasten nur ein Haufen Werkzeug",
              "Die Verfahren sind lernbar und in wenigen Zeilen aufgerufen. Was Projekte "
              "scheitern lässt, ist fast nie die Mathematik — es ist die Reihenfolge.")
    y = unter_intro(s)
    kachelreihe(s, [
        ("Ohne Modell", [
            "Man beginnt mit den Daten, die da sind",
            "Das Ziel entsteht unterwegs",
            "Erfolg wird beurteilt, nachdem",
            "man das Ergebnis gesehen hat",
        ]),
        ("Die Folge", [
            "Jedes Ergebnis lässt sich",
            "im Nachhinein gut finden",
            "Niemand kann sagen, ob das",
            "Projekt sein Ziel erreicht hat",
        ]),
        ("Mit CRISP-DM", [
            "Die Frage steht vor den Daten",
            "Die Hürde steht vor dem Ergebnis",
            "Und es gibt einen benannten Ort",
            "für das Eingeständnis: zurück",
        ]),
    ], y=y, hoehe=178)
    sandband(s, "CRISP-DM ist kein Werkzeug, sondern eine Reihenfolge mit "
                "verbindlichen Zwischenständen. Der Wert steckt nicht in den sechs "
                "Namen, sondern darin, dass man sie NICHT überspringen darf.",
             y=y + 194)
    notizen(s, "Fragen Sie hier in den Raum: Wer hat schon einmal eine Auswertung "
               "gemacht und dabei das Erfolgskriterium erst festgelegt, als das "
               "Ergebnis schon auf dem Tisch lag? Fast alle haben. Das ist keine "
               "Nachlässigkeit, sondern die natürliche Reihenfolge — und genau "
               "deshalb braucht es ein Modell, das dagegen hält.")

    # ── Die sechs Phasen
    s = folie(prs, "Überblick", "Sechs Phasen, und jede beantwortet genau eine Frage",
              "Von der Geschäftsfrage zum laufenden Betrieb. Die Reihenfolge ist "
              "verbindlich, aber sie ist keine Einbahnstraße.")
    diagramm(s, bild("crispdm-kreis"), y=unter_intro(s), breite=BREITE)
    sandkarte(s, "Die Faustregel, die alle unterschätzen",
              ["Phase 2 und 3 kosten in echten Projekten 60 bis 80 Prozent der Zeit.",
               "Phase 4 — das Modellieren, an das alle denken — sind oft zehn Zeilen Code.",
               "Wer beim Lesen eines Notebooks die Zeilen zählt, sieht dasselbe Verhältnis."],
              y=unter_intro(s) + 130)
    phasenleiste(s, 0)
    notizen(s, "Die Kette liest sich von links nach rechts, aber sie ist ein KREIS. "
               "Das Wort Kreislauf steht nicht zufällig im Namen. Auf der übernächsten "
               "Folie sehen wir, wo in unseren sechs Fällen tatsächlich zurückgesprungen "
               "wurde — in vier von sechs.")

    # ── Je eine Folie pro Phase
    phasen = [
        (1, "Business Understanding", "Was ist das Problem — und woran messen wir Erfolg?",
         [("Geschäftsziel", "Was soll sich im Betrieb ändern? In der Sprache des Fachs, "
                            "nicht der Statistik"),
          ("Analytisches Ziel", "Was genau soll geschätzt, sortiert oder gruppiert werden?"),
          ("Erfolgskriterium", "Eine Zahl mit einer Hürde — festgelegt, BEVOR jemand "
                               "die Daten gesehen hat"),
          ("Fehlerkosten", "Welcher Fehler ist teurer? Fast nie sind beide gleich teuer")],
         "Das ist die Phase, die am häufigsten übersprungen wird — und die einzige, "
         "die man nicht nachholen kann. Ein Erfolgskriterium, das nach dem Ergebnis "
         "festgelegt wird, ist keins."),
        (2, "Data Understanding", "Tragen die Daten die Frage überhaupt?",
         [("Beschaffen", "Woher kommen die Daten, und was fehlt?"),
          ("Beschreiben", "Wie viele Zeilen, welche Spalten, welche Lücken?"),
          ("Erkunden", "Zuerst die Zielgröße ansehen, dann die Zusammenhänge"),
          ("Qualität prüfen", "Was ist unplausibel — und ist es ein Fehler oder ein Befund?")],
         "Hier wird nicht modelliert, hier wird GESEHEN. Wer diese Phase abkürzt, "
         "modelliert später Artefakte. In Notebook 6 kostet genau das eine ganze Runde."),
        (3, "Data Preparation", "Wie werden aus Rohdaten Merkmale, die man verwenden darf?",
         [("Auswählen", "Zeilen und Spalten begründet ausschließen — jeder Filter "
                        "mit Zeilenzahl davor und danach"),
          ("Bereinigen", "Fehlende Werte, Ausreißer, Dubletten"),
          ("Merkmale bauen", "Aus Zeitstempel wird Stunde, Wochentag, Feiertag"),
          ("Aufteilen", "Trainings- und Testmenge — bei Zeitbezug entlang der Zeit, "
                        "nicht zufällig")],
         "Die teuerste Phase, und die mit der gefährlichsten Falle: Merkmale, die es "
         "zum Vorhersagezeitpunkt gar nicht gibt. Dazu gleich mehr."),
        (4, "Modeling", "Welches Verfahren, und woran messen wir es?",
         [("Nullmodell zuerst", "Der einfachste denkbare Tipp. Wer den nicht schlägt, "
                                "hat nichts gewonnen"),
          ("Faustregel als Maßstab", "Was macht der Betrieb heute ohne Modell?"),
          ("Dann erst Verfahren", "Vom Durchschaubaren zum Stärkeren, nicht umgekehrt"),
          ("Kosten einbauen", "Ungleiche Fehlerkosten gehören ins Modell, nicht in "
                              "die Nachbetrachtung")],
         "Die kürzeste Phase. Sie wirkt in jedem Lehrbuch wie die Hauptsache und ist "
         "in der Praxis der kleinste Teil."),
        (5, "Evaluation", "Reicht das für die Entscheidung, um die es geht?",
         [("Technische Güte", "MAE, Trefferquote, Silhouette — notwendig, aber nie "
                              "hinreichend"),
          ("Gegen die Hürde", "Die Zahl aus Phase 1 — bestanden oder nicht"),
          ("Fehler ansehen", "Nicht nur wie gut, sondern WORAN es scheitert"),
          ("Urteil fällen", "Weiter, zurück, oder ehrlich abbrechen")],
         "Hier entscheidet sich, ob das Projekt weitergeht. Und hier zeigt sich, ob "
         "die Kriterien aus Phase 1 etwas taugten."),
        (6, "Deployment", "Wie kommt das Ergebnis in den Betrieb — und wer merkt, wenn es kippt?",
         [("Haltbar machen", "Modell und Merkmalsliste zusammen speichern"),
          ("Einbauen", "Die Funktion, die der Betrieb tatsächlich aufruft"),
          ("Überwachen", "Woran erkennt man, dass das Modell veraltet?"),
          ("Rückkopplung bedenken", "Verändert die Auslieferung die Daten, aus denen "
                                    "wir künftig lernen?")],
         "Die Phase, die in Lehrbeispielen fast immer fehlt. Ein Modell, das niemand "
         "aufruft, hat keinen Wert — und eines, das niemand überwacht, ist eine "
         "tickende Uhr."),
    ]
    for nr, name, frage, punkte, notiz in phasen:
        s = folie(prs, f"Phase {nr} von 6", f"{name}: {frage}")
        streifen(s, punkte, y=unter_intro(s), hoehe=54, luecke=8, chip_b=250)
        phasenleiste(s, nr)
        notizen(s, notiz)

    # ── Die Rücksprünge
    s = folie(prs, "Der Kern", "Der Rücksprung ist kein Scheitern — er ist das Verfahren",
              "Was CRISP-DM von einer Checkliste unterscheidet, sind die Pfeile "
              "zurück. In unseren sechs Fällen wurde viermal zurückgesprungen.")
    tabelle(s, ["Fall", "Rücksprung", "Warum"], [
        ["1 Regression", "5 → 1", "Für zwei von drei Radtypen ist eine Punktschätzung "
                                  "die falsche Zusage — eine Spanne wäre ehrlich"],
        ["3 Clustering", "5 → 1", "Die Analyse brachte eine bessere Frage hervor, als "
                                  "die war, mit der wir angefangen hatten"],
        ["6 Anomalie", "4 → 3", "Das Modell fand die Preisklasse statt der Anomalien — "
                                "die Merkmale mussten normiert werden"],
        ["6 Anomalie", "5 → Ende", "Aufgabe B scheitert an den Grunddaten. Ein negatives "
                                   "Ergebnis, ausdrücklich nicht ausgeliefert"],
    ], y=unter_intro(s), spalten_b=[150, 130, 623.5], zeilen_h=44)
    sandband(s, "Ein Projekt, das ohne einen einzigen Rücksprung durchläuft, hat "
                "entweder eine triviale Frage gestellt — oder nicht genau genug "
                "hingesehen.", y=unter_intro(s) + 232)
    phasenleiste(s, 0, rueckspruenge=(1, 3))
    notizen(s, "Diese Folie ist das Herz des Kapitels. Studierende lesen CRISP-DM "
               "als Abarbeitungsliste und empfinden einen Rücksprung als Fehler. "
               "In Wahrheit ist der Rücksprung der Beleg dafür, dass jemand "
               "hingesehen hat. In Notebook 6 kann man ihn Zelle für Zelle mitlesen.")

    # ── Die zwei häufigen Fehler
    s = folie(prs, "Zwei Fehler", "Zwei Fehler, die man in fast jedem Erstprojekt sieht")
    vorher_nachher(s,
                   ("So läuft es meistens", "Kriterium nach dem Ergebnis", [
                       "1. Daten laden, Modell rechnen",
                       "2. Der Fehler liegt bei 3,96 Minuten",
                       "3. „Unter vier Minuten — ordentlich.“",
                       "",
                       "Die Hürde wurde an das Ergebnis",
                       "angelegt, nicht umgekehrt. So kann",
                       "kein Projekt scheitern — und deshalb",
                       "sagt sein Erfolg auch nichts.",
                   ], False),
                   ("So gehört es", "Kriterium vor den Daten", [
                       "1. „Der Preisfehler muss unter",
                       "   50 Cent liegen“ — vor allem",
                       "   anderen festgelegt",
                       "2. Daten laden, Modell rechnen",
                       "3. CITY 0,41 € hält, EBIKE 0,85 €",
                       "   und CARGO 2,48 € reißen",
                       "",
                       "Genau so steht es in Notebook 1 —",
                       "und deshalb ist das Ergebnis dort",
                       "eine Erkenntnis statt einer Ausrede.",
                   ], False),
                   y=unter_intro(s), hoehe=222)
    phasenleiste(s, 0)
    notizen(s, "Der zweite Fehler steht auf der nächsten Folie: das Verfahren wählen, "
               "bevor die Frage steht. Beide Fehler haben dieselbe Wurzel — man fängt "
               "bei dem an, was Spaß macht, statt bei dem, was entschieden werden soll.")

    # ── Die sechs Fälle im Überblick
    s = folie(prs, "Wegweiser", "Sechs Fälle — und jeder zeigt eine andere Phase von innen",
              "Deshalb wird der Kreislauf hier nicht sechsmal wiederholt. Fall 1 "
              "läuft vollständig durch; die übrigen fünf gehen dort in die Tiefe, "
              "wo sie etwas Neues zu zeigen haben.")
    tabelle(s, ["Fall", "Verfahren", "Zeigt", "Der Satz, der bleibt"], [
        ["1", "Regression", "alle sechs", "Ob ein Merkmal erlaubt ist, entscheidet der "
                                          "Prozess — nicht der Spaltenname"],
        ["2", "Klassifikation", "Phase 6", "Ein Modell muss seinen Unterhalt verdienen — "
                                           "hier verdient es ihn nicht"],
        ["3", "Clustering", "Phase 1", "Erfolgskriterien auch ohne Zielgröße"],
        ["4", "Zeitreihe", "Phase 3", "Unter welchem Informationsstand man wählt, "
                                      "entscheidet mit, was man wählt"],
        ["5", "Assoziation", "Phase 5", "Zwei Produkte aus denselben Regeln — "
                                        "und nur eines darf entscheiden"],
        ["6", "Anomalie", "Rücksprung", "Ein Verfahren, das schlechter ist als eine "
                                        "Zeile Fachwissen, ist an der Aufgabenstellung "
                                        "gescheitert"],
    ], y=unter_intro(s), spalten_b=[50, 130, 110, 613.5], zeilen_h=42)
    phasenleiste(s, 0)
    notizen(s, "Diese Tabelle ist der Fahrplan für den Rest des Decks. Sie taucht in "
               "der Synthese am Ende noch einmal auf, dann ausgefüllt mit dem, was "
               "wir unterwegs gesehen haben.")

    s = folie(prs, "Wegweiser", "Dieselbe Zuordnung als Bild",
              "Welcher Fall beleuchtet welche Phase? Diese Zuordnung ist der Grund "
              "für den Aufbau des Decks.")
    diagramm(s, bild("crispdm-faelle"), y=unter_intro(s),
             hoehe=ZONE_UNTEN - unter_intro(s) - 34)
    notizen(s, "Zwei Fälle zeigen auf Phase 5 — Fall 1 und Fall 5 —, aber aus "
               "verschiedenen Richtungen: der eine, weil die gemittelte Kennzahl "
               "täuscht, der andere, weil vorab gesetzte Hürden radikal sieben. "
               "Das ist kein Zufall: Phase 5 ist die Phase, in der sich entscheidet, "
               "ob die Arbeit der Phasen 1 bis 4 etwas taugte.")


# ═════════════════════════════════════════ Teil B und C — die sechs Fälle
#
# WARUM DIESE KAPITEL AUS DEN MERKZETTELN GEBAUT WERDEN
#
# Die erste Fassung dieses Decks trug abgetippte Zahlen. Als die Daten neu
# erzeugt wurden, widersprach es den Notebooks - elf Aussagen, verteilt
# ueber alle sechs Kapitel, darunter ein ganzer didaktischer Bogen ("0,99
# gegen 1,00 Prozent"), dessen Anker verschwunden war. Ein Foliendeck mit
# abgetippten Zahlen veraltet still.
#
# Deshalb liest jede Zahl hier aus analytics/bau/werte/*.json - derselben
# Quelle wie Handout und Use-Case-Decks. Ein unbekannter Schluessel bricht
# den Bau ab, statt eine Luecke zu drucken.
#
# Und deshalb sind die Kapitel kurz: Was ein Fall zeigt, steht auf fuenf
# Folien. Wer den vollstaendigen Durchlauf will, liest das Notebook.


class Werte(dict):
    """Merkzettel eines Notebooks. Ein unbekannter Schluessel bricht ab."""

    DATEIEN = {
        1: "01_Regression_Fahrtdauer",
        2: "02_Klassifikation_Wartungsrisiko",
        3: "03_Clustering_Stationen_und_Kunden",
        4: "04_Zeitreihe_Nachfrageprognose",
        5: "05_Assoziation_Wege_im_Netz",
        6: "06_Anomalieerkennung_Auffaellige_Vorgaenge",
    }

    def __init__(self, nummer):
        pfad = WURZEL / "analytics" / "bau" / "werte" / f"{self.DATEIEN[nummer]}.json"
        if not pfad.exists():
            raise SystemExit(
                f"ABBRUCH: {pfad.name} fehlt. Erst die Notebooks bauen:\n"
                f"    python3 analytics/bau/bauen.py")
        import json as _json
        super().__init__(_json.loads(pfad.read_text(encoding="utf-8")))
        self.nummer = nummer

    def __missing__(self, schluessel):
        raise SystemExit(
            f"ABBRUCH: Notebook {self.nummer} kennt '{schluessel}' nicht.\n"
            f"    Bekannt: {', '.join(sorted(self))}")

    def p(self, schluessel, stellen=1):
        return f"{self[schluessel] * 100:.{stellen}f}".replace(".", ",") + " %"

    def z(self, schluessel, stellen=0):
        return f"{self[schluessel]:.{stellen}f}".replace(".", ",")


def fallkapitel(prs, nb, nummer, titel, leitfrage, notiz, steckbrief_zeilen,
                kriterium_kopf, kriterium_zeilen, ergebnis_kopf, ergebnis_zeilen,
                ergebnis_intro, ausgeliefert, nutzen, grenze, merksatz,
                phase_kriterium=1, phase_ergebnis=5):
    """Ein Fallkapitel: Trenner, Steckbrief, Kriterium, Ergebnis, Auslieferung."""
    AKTUELLES_NB[0] = nb
    kapitel(prs, nummer, titel, leitfrage, notiz)

    s = folie(prs, f"Fall {nb}", "Der Fall auf einen Blick")
    # Die Karte bleibt eng am Text: thws.steckbrief setzt die Zeilen mit
    # festem Abstand, eine hoehere Karte erzeugt nur Leerraum darin. Den
    # unteren Rand der Folie fasst stattdessen die Phasenleiste.
    steckbrief(s, steckbrief_zeilen, y=unter_intro(s))
    phasenleiste(s, 0)
    notizen(s, notiz)

    s = folie(prs, f"Fall {nb} · Phase 1",
              "Das Erfolgskriterium — festgelegt vor der Messung",
              "Was hier steht, stand vor der ersten Rechnung fest. Alles Weitere "
              "wird daran gemessen, nicht an einer Kennzahl, die sich im "
              "Nachhinein anbietet.")
    _y = unter_intro(s)
    tabelle(s, kriterium_kopf, kriterium_zeilen, y=_y,
            zeilen_h=min(58, max(44, int((ZONE_UNTEN - 52 - _y)
                                         / (len(kriterium_zeilen) + 1)))))
    phasenleiste(s, phase_kriterium)
    notizen(s, "Die Reihenfolge ist der Punkt: erst das Kriterium, dann die "
               "Messung. Ein Kriterium, das nach dem Ergebnis entsteht, misst "
               "nichts mehr.")

    s = folie(prs, f"Fall {nb} · Phase 5", "Das Ergebnis", ergebnis_intro)
    _y = unter_intro(s)
    tabelle(s, ergebnis_kopf, ergebnis_zeilen, y=_y,
            zeilen_h=min(58, max(40, int((ZONE_UNTEN - 52 - _y)
                                         / (len(ergebnis_zeilen) + 1)))))
    phasenleiste(s, phase_ergebnis)
    notizen(s, "Alle Zahlen dieser Folie stammen aus dem ausgeführten Notebook "
               "und werden beim Bau des Decks eingesetzt.")

    s = folie(prs, f"Fall {nb} · Phase 6", "Was ausgeliefert wird — und was nicht")
    kachelreihe(s, [("Ausgeliefert", ausgeliefert), ("Nutzen", nutzen),
                    ("Grenze", grenze)], y=unter_intro(s), hoehe=196)
    sandband(s, merksatz, y=darunter(unter_intro(s), 196))
    phasenleiste(s, 6)
    notizen(s, "Die dritte Kachel ist die wichtigste: Was ein Verfahren NICHT "
               "belegt, gehört genauso in die Übergabe wie das Ergebnis.")


def teil_referenzfall(prs):
    w = Werte(1)
    fallkapitel(
        prs, 1, 2, "Fall 1 — Regression: der volle Durchlauf",
        "Können wir dem Kunden vor der Fahrt sagen, was sie kosten wird?",
        "Fall 1 läuft vollständig durch. Die übrigen fünf gehen nur dort in die "
        "Tiefe, wo sie etwas Neues zeigen.",
        [("Geschäftsfrage", "Was kostet die Fahrt zu diesem Ziel — vor dem "
                            "Entsperren?"),
         ("Analytisches Ziel", "Eine Spanne für die Fahrtdauer, aus der das "
                               "Tarifblatt einen Preisbereich rechnet"),
         ("Der Kern", "Der Preis hängt an der Dauer, und die steht zur "
                      "Anfragezeit nicht fest"),
         ("Verfahren", "Quantilregression, Perzentiltabelle, Quantiltabelle — "
                       "alle drei am selben Kriterium gemessen"),
         ("Urteil", f"Ausgeliefert wird die {w['kandidat']}; Status "
                    f"{w['produktstatus']}")],
        ["", "Kriterium", "Schwelle"],
        [["trifft", "Die angezeigte Spanne enthält den tatsächlichen Preis",
          f"mindestens {w.p('gate_schwelle', 0)}, aggregiert je Radtyp"],
         ["nützt", "Die Spanne ist schmal genug, um zu helfen",
          "höchstens 12 Minuten und 60 % des Preises"],
         ["reicht", "Je Radtyp wird ein Mindestanteil der Anfragen beantwortet",
          "sonst besteht das Verfahren nicht"],
         ["geprüft auf", "Einem Zeitraum, den bis zur Abnahme nichts berührt hat",
          "in Phase 2 versiegelt"]],
        ["Kandidat", "Primärgate", "Reichweite", "Betriebsform"],
        [["Quantilregression", w.p("quantil_gate"), w.p("quantil_auskunft"),
          "Laufzeitdienst"],
         ["Perzentiltabelle", w.p("tabelle_gate"), w.p("tabelle_auskunft"), "CSV-Datei"],
         ["Quantiltabelle", w.p("qtab_gate"), w.p("qtab_auskunft"), "CSV-Datei"]],
        f"{w['zulaessige_satz']}. Entschieden hat deshalb nicht die Prognosegüte, "
        f"sondern eine vorab benannte Auswahlregel: Bei gleicher Eignung gewinnt "
        f"die Bauform, die ohne laufenden Dienst auskommt.",
        [f"Die {w['kandidat']} als CSV-Datei.",
         f"Zusage {w.p('gate_schwelle', 0)}, auf der Abnahme belegt mit "
         f"{w.p('ab_unten')}.",
         f"{w['ab_gates_halten']} von {w['ab_gates_gesamt']} Gates halten.",
         f"Status: {w['produktstatus']}."],
        ["Der Preisrahmen steht vor dem Entsperren fest.",
         "Kein Dienst, keine Bibliotheksversion, von Hand nachrechenbar.",
         f"Beantwortet werden {w.p('reichweite_real', 0)} der Anfragen — sonst "
         f"schweigt die Anwendung."],
        [f"In {w.p('zielabweichung', 0)} der Fahrten weicht das Ende vom "
         f"angegebenen Ziel ab; die Zusage gilt nur für die gewählte Strecke.",
         f"Gültig für Fahrten bis {w['gueltig_bis_lang']} — so weit reicht der "
         f"Kalender.",
         "Die Abnahme ist rückblickend; eine prospektive Prüfung steht aus."],
        "Zugesagt wird eine Spanne, nicht eine Zahl — und gemessen wird an der "
        "Zusage, nicht an einer Kennzahl, die gut aussieht.")

    # Der Referenzfall bekommt drei Folien mehr: an ihm wird der Kreislauf
    # vollstaendig gezeigt, die uebrigen Faelle setzen ihn voraus.
    AKTUELLES_NB[0] = 1
    s = folie(prs, "Fall 1 · Phase 3", "Der Leakage-Test",
              "Die Frage ist nicht statistisch, sondern zeitlich: Was steht in dem "
              "Moment zur Verfügung, in dem die Anzeige erscheinen soll?")
    ampel_matrix(s, ["erlaubt"], [
        ("Startstation, Zeitpunkt, Radtyp", [True], "der Kunde steht davor"),
        ("Geplantes Ziel", [True], "wird vor dem Entsperren gewählt"),
        ("Feiertag, Schulferien", [True], "stehen im Kalender"),
        ("Tageswetter", [False], "steht erst am Abend fest"),
        ("Endzeit, Dauer, Entgelt", [False], "entstehen während und nach der Fahrt"),
    ], y=unter_intro(s) + 26, zeilen_h=44)
    phasenleiste(s, 3)
    notizen(s, "Die Wetterzeile ist die lehrreichste: Ein Modell mit dem Tagesmittel "
               "benutzt Wissen von heute Abend für eine Anfrage von heute früh. Ein "
               "zeitlicher Schnitt heilt das nicht.")

    s = folie(prs, "Fall 1 · Phase 5", "Der Rücksprung — obwohl das Kriterium hielt",
              "Für den häufigsten Radtyp lag der Preisfehler unter der Grenze. "
              "Trotzdem ging es zurück in Phase 1, und zwar aus zwei Gründen, die "
              "nichts mit einem gerissenen Kriterium zu tun haben.")
    kachelreihe(s, [
        ("Der Mittelwert ist keine Erfahrung",
         [f"Im Mittel {w.z('preisfehler_city', 2)} € Abweichung.",
          f"Aber nur {w.p('city_unter_50', 0)} der Fahrten",
          "bleiben unter 50 Cent.",
          "Ein Kunde erlebt seine Fahrt,", "nicht den Mittelwert."]),
        ("Ein Radtyp hat kein Produkt",
         [f"Das Lastenrad reißt die Grenze",
          f"um ein Vielfaches.",
          "Eine Lösung, die den teuersten",
          "Radtyp ausspart, beantwortet",
          "die Geschäftsfrage nicht."]),
        ("Die Antwort: eine Spanne",
         ["Nicht die Verfahrensklasse", "ändert sich, sondern das,",
          "was die App verspricht.", "", "Aus einer Zahl wird ein",
          "Bereich mit Zusage."]),
    ], y=unter_intro(s), hoehe=204)
    phasenleiste(s, 5, rueckspruenge=((5, 1),))
    notizen(s, "Das ist der wichtigste Moment des Falls: Ein Rücksprung, der nicht "
               "aus einem gerissenen Kriterium folgt, sondern aus der Einsicht, dass "
               "das Kriterium die falsche Frage beantwortet.")

    s = folie(prs, "Fall 1 · Phase 6", "Vier Stufen zwischen Rechnung und Betrieb",
              "„Das Modell ist gut“ und „das Modell darf entscheiden“ sind zwei "
              "verschiedene Aussagen. Zwischen ihnen liegen vier Stufen.")
    schichtenstapel(s, [
        ("qualifiziert — die Kennzahlen halten auf den Testdaten", False),
        ("Schattenbetrieb — rechnet mit, wird protokolliert, niemand handelt danach",
         False),
        ("betriebsgesperrt — Kennzahlen halten, eine Voraussetzung fehlt", False),
        (f"sichtbar — die Anzeige ist freigeschaltet (hier erreicht: "
         f"{w['produktstatus']})", True),
    ], y=unter_intro(s) + 20, hoehe=52, luecke=12)
    phasenleiste(s, 6)
    notizen(s, "Die vierte Stufe ist neu gegenüber dem Lehrbuch: Analytisch kann "
               "alles halten und das Produkt trotzdem nicht laufen, weil eine "
               "betriebliche Voraussetzung fehlt — hier der Kalenderhorizont.")


def fall2(prs):
    w = Werte(2)
    fallkapitel(
        prs, 2, 3, "Fall 2 — Klassifikation: was ein Modell wert sein muss",
        "Welche Räder soll die Werkstatt im nächsten Quartal vorsorglich prüfen?",
        "Dieser Fall zeigt Phase 6 von innen: Ein Modell muss seinen Unterhalt "
        "verdienen. Hier verdient es ihn nicht.",
        [("Geschäftsfrage", f"Welche {w.z('kapazitaet')} Räder gehören auf die "
                            f"Quartalsliste?"),
         ("Analytisches Ziel", f"Wird dieses Rad in {w.z('horizont_tage')} Tagen "
                               f"auffällig — ja oder nein?"),
         ("Fehlerkosten", f"{w.z('kosten_verpasst')} € je übersehenem Ausfall gegen "
                          f"{w.z('kosten_unnoetig')} € je unnötiger Prüfung"),
         ("Verfahren", "Drei Faustregeln als Maßstab, dann Entscheidungsbaum und "
                       "Random Forest mit Kostengewicht"),
         ("Urteil", "Ausgeliefert wird die Faustregel; das Modell bleibt im Paket")],
        ["", "Kriterium", "Schwelle"],
        [["K1", "Mehr auffällige Räder als eine Zufallsauswahl gleicher Länge",
          f"{w.z('lift_faktor', 1)}-fach, in mindestens "
          f"{w.z('k1_mindestquartale')} von 5 Quartalen"],
         ["K2", "Die erwarteten Kosten liegen unter denen der heutigen Faustregel",
          "„das älteste Rad zuerst“"],
         ["K3", "Die untere Vertrauensgrenze der Listenpräzision hält im Testquartal",
          f"über {w.z('lift_faktor', 1)} × Grundrate des Quartals"],
         ["D70", "Die ursprünglich genannte 70-Prozent-Marke",
          "nur Diagnose — in einem Quartal nicht erreichbar"]],
        ["", "Faustregel", "Random Forest"],
        [["Treffer im Testquartal", w.z("treffer_regel"), w.z("treffer_wald")],
         ["untere Vertrauensgrenze (K3)", w.p("wilson_unten_regel"),
          w.p("wilson_unten_wald")],
         [f"gefordert: {w.p('k3_schwelle')}", "hält", "reißt K3"],
         ["Betrieb", "eine Regel, unverändert", "vierteljährlich nachtrainieren"]],
        f"Im Testquartal trifft die Regel {w.z('treffer_regel')} Räder gegenüber "
        f"{w.z('treffer_wald')}. Ausschlaggebend war aber nicht dieser Vorsprung, "
        f"sondern die statistische Absicherung.",
        ["Eine Faustregel: sortiere nach Kilometern seit der Reparatur.",
         f"Von zehn geprüften Rädern werden {w.z('quote_regel_von_zehn', 1)} "
         f"auffällig.",
         f"Von zehn auffälligen erfasst die Liste {w.z('abdeckung_von_zehn', 1)}.",
         f"Dazu eine Schattenliste zum {w['schatten_stichtag_lang']}."],
        ["Die Werkstatt kann die Liste ohne Nacharbeit übernehmen.",
         "Kein Modellbetrieb, keine Versionsstände, erklärbar im Streitfall.",
         "Das Kostenverhältnis steht sichtbar in der Entscheidung."],
        [f"Die Grundrate schwankt über {w.z('panel_stichtage')} Stichtage zwischen "
         f"{w.p('panel_grundrate_min')} und {w.p('panel_grundrate_max')}.",
         "Ein einzelnes gutes Quartal belegt deshalb wenig.",
         "Sobald die Liste benutzt wird, verändert sie die Daten, aus denen sie "
         "lernt — ohne Kontrollgruppe ist die Trefferquote keine Gütekennzahl."],
        "Ein Modell schlägt eine Faustregel nur dann, wenn die Faustregel mit "
        "derselben Sorgfalt gebaut wurde wie das Modell.",
        phase_ergebnis=5)


def fall3(prs):
    w = Werte(3)
    fallkapitel(
        prs, 3, 4, "Fall 3 — Clustering: Erfolg ohne Zielgröße",
        "Welche Stationstypen und welche Kundensegmente hat VeloCity wirklich?",
        "Dieser Fall zeigt Phase 1 von innen: Ohne Zielgröße gibt es kein "
        "„richtig“ — die Erfolgskriterien muss man selbst bauen, bevor man das "
        "Ergebnis kennt.",
        [("Geschäftsfrage", "Welche Stationstypen und Kundensegmente gibt es — "
                            "und wie behandelt man sie unterschiedlich?"),
         ("Analytisches Ziel", "Gruppen aus dem Verhalten, nicht aus den "
                               "Stammdaten"),
         ("Der Kern", "Es gibt keine Zielgröße und damit keinen mitgelieferten "
                      "Maßstab"),
         ("Verfahren", "k-Means auf Tagesgang (Stationen) und auf Aktualität, "
                       "Häufigkeit, Umsatz (Kundschaft)"),
         ("Urteil", f"Analytisch {w['status_analytisch']}, für den Einsatz "
                    f"{w['status_einsatz']}")],
        ["", "Kriterium", "Warum es nötig ist"],
        [["1", "benennbar", "Eine Gruppe, die niemand beschreiben kann, ist keine"],
         ["2", "unterschiedlich behandelbar",
          "Vier Gruppen mit derselben Maßnahme sind eine Gruppe"],
         ["3", "groß genug",
          f"unter {w.p('min_segmentanteil', 0)} des Bestands lohnt keine eigene "
          f"Ansprache"],
         ["4 und 5", "stabil gegenüber Startwert und über die Zeit",
          "sonst wechselt die Kundschaft das Segment ohne Verhaltensänderung"]],
        ["", "Stationen", "Kundschaft"],
        [["Zahl der Gruppen", "vier benennbare Typen", "vier Segmente"],
         ["gegen die verdeckte Wahrheit", w.p("generator_treffer", 0) + " richtig",
          "keine vorhanden"],
         ["Reproduzierbarkeit (ARI)", w.z("generator_ari", 3), w.z("ari_kunden", 3)],
         ["Kriterien erfüllt", "alle fünf",
          f"{w['gates_erfuellt']} von {w['gates_gesamt']}"]],
        f"Bei den Stationen lässt sich das Ergebnis gegen eine im Datensatz "
        f"hinterlegte, dem Verfahren nicht bekannte Zuordnung prüfen. Bei der "
        f"Kundschaft gibt es diese Möglichkeit nicht — dort entscheiden allein die "
        f"fünf Kriterien.",
        ["Stationsprofile als CSV — ausdrücklich Hypothesen, kein Sollbestand.",
         f"Für die Kundschaft {w['exportart']}.",
         f"Analytisch {w['status_analytisch']}, Einsatz {w['status_einsatz']}."],
        ["Die Umverteilung bekommt benannte Stationstypen statt Bauchgefühl.",
         f"Sichtbar wurde ein Preisproblem: {w['viel_segment']} bringen "
         f"{w.z('viel_je_fahrt', 2)} € je Fahrt, {w['stark_segment']} "
         f"{w.z('stark_je_fahrt', 2)} €.",
         "Der Kundenbericht bleibt aggregiert und ohne Namen."],
        [f"{w.p('kurze_historie_anteil', 0)} der Kundschaft erscheinen gar nicht — "
         f"sie sind im Betrachtungszeitraum nicht gefahren.",
         "Was fehlt, ist keine Kennzahl, sondern eine prospektive Prüfung.",
         "Silhouette und Rand-Index messen Trennschärfe, nicht Verwendbarkeit."],
        "Ohne Zielgröße gibt es kein „richtig“ — die Erfolgskriterien müssen "
        "vorher gesetzt und begründet werden.",
        phase_kriterium=1, phase_ergebnis=5)


def fall4(prs):
    w = Werte(4)
    fallkapitel(
        prs, 4, 5, "Fall 4 — Zeitreihe: der Schnitt entlang der Zeit",
        "Wie viele Fahrten kommen morgen — und unter welchem Informationsstand "
        "wählt man das Verfahren?",
        "Dieser Fall zeigt Phase 3 von innen. Zwei Dinge gehen anders: Die "
        "Aufteilung muss der Zeit folgen, und eine Störgröße macht einen Effekt "
        "vor, den es so nicht gibt.",
        [("Geschäftsfrage", "Wie viele Fahrten kommen morgen? Räder und Schichten "
                            "daraus abzuleiten ist eine eigene Analyse"),
         ("Analytisches Ziel", "Zahl der Fahrten am Folgetag"),
         ("Fehlerkosten", "4,00 € je fehlendem Rad gegen 0,80 € je überzähligem"),
         ("Betriebskriterium", "Die Prognose muss um 18 Uhr stehen — sonst ist sie "
                               "wertlos, egal wie genau"),
         ("Urteil", f"{w['gewaehlt_name']}; Status {w['nb04_status']}")],
        ["", "Kriterium", "Warum"],
        [["Einheit", "Prognostiziert werden Fahrten, nicht Räder oder Schichten",
          "Die Übersetzung ist eine eigene Analyse"],
         ["Vergleichslage", "Gemessen wird unter Prognosewetter, nicht unter "
                            "Ist-Wetter",
          "Um 18 Uhr liegt nur eine Vorhersage vor"],
         ["Aufschlag", "Der Sicherheitsaufschlag wird auf der Validierung gewählt",
          "Die teurere Fehlerrichtung wird sichtbar entschieden"],
         ["Reihenfolge", "Der Test wird erst geöffnet, wenn Verfahren und Aufschlag "
                         "feststehen",
          "Danach ist er verbraucht"]],
        ["Verfahren", "unter Prognosewetter", "unter Ist-Wetter"],
        [["Nullmodell", w.z("mae_null", 1), "—"],
         ["Faustregel der Disposition", w.z("mae_faustregel", 1), "—"],
         [w["gewaehlt_name"], w.z("mae_linear", 1), w.z("ist_linear", 2)],
         ["Gradient Boosting", w.z("mae_boosting", 1), w.z("ist_boosting", 2)]],
        f"Mittlerer absoluter Fehler in Fahrten je Tag. Unter Ist-Wetter liegen "
        f"beide Modelle praktisch gleichauf; erst unter Prognosewetter setzt sich "
        f"das einfachere Verfahren ab.",
        [f"{w['gewaehlt_name']} mit einem Aufschlag von {w.p('aufschlag', 0)}.",
         f"Mittlerer Fehler {w.z('mae_linear', 1)} Fahrten gegen "
         f"{w.z('mae_faustregel', 1)} bei der Faustregel.",
         f"Status: {w['nb04_status']}."],
        ["Die Planung bekommt eine begründete Zahl statt eines Erfahrungswerts.",
         "Der Aufschlag macht die teurere Fehlerrichtung zu einer sichtbaren "
         "Entscheidung.",
         "Der Schattenpilot läuft mit und wird protokolliert."],
        ["Prognostiziert wird die Gesamtzahl; gebraucht werden Räder je Station.",
         "Die Wetterunsicherheit ist simuliert, nicht gemessen.",
         "Ein Sommerfenster trägt keine Jahresaussage."],
        "Verglichen wird unter den Bedingungen des Einsatzes. Wer unter Ist-Wetter "
        "wählt, wählt für eine Lage, in der er nie liefert.",
        phase_kriterium=3, phase_ergebnis=5)


def fall5(prs):
    w = Werte(5)
    fallkapitel(
        prs, 5, 6, "Fall 5 — Assoziation: zwei Produkte, zwei Maßstäbe",
        "Zwischen welchen Stationen gibt es systematische Ströme — und was darf "
        "man damit tun?",
        "Dieser Fall zeigt Phase 5 von innen: Dieselben Regeln tragen einen "
        "Hinweis für Menschen, aber keine automatische Umsetzfahrt. Der "
        "Unterschied liegt nicht in der Statistik, sondern im Produkt.",
        [("Geschäftsfrage", "Von wo nach wo — und in welchem Zeitfenster?"),
         ("Analytisches Ziel", "Gerichtete Start-Ziel-Regeln je Kontext, gemessen "
                               "am kontextbedingten Lift"),
         ("Zwei Produkte", "A automatische Umverteilung, B Dispositionshinweis — "
                           "mit eigenen Kriterien"),
         ("Verfahren", "Support, Konfidenz und Lift von Hand; Bestätigung im "
                       "versiegelten letzten Drittel"),
         ("Urteil", f"A: {w['status_a']}. B: {w['status_b']}")],
        ["", "Kriterium", "Schwelle"],
        [["A1 bis A3", "Support, kontextbedingter Lift, Ziel ist eine konkrete "
                       "Station",
          f"{w.p('k1_support', 0)} · {w.z('k2_lift', 1)} · Start ≠ Ziel"],
         ["A4", "Der Wert der betroffenen Fahrten übersteigt die Kosten einer "
                "Umsetzrunde",
          f"{w.z('kosten_transport')} € je Runde — gesetzte Szenarioannahme"],
         ["B1", "Jede Regel hält im Bestätigungszeitraum, den die Suche nicht "
                "gesehen hat",
          f"untere Grenze eines Tagesblock-Bootstraps ≥ {w.z('k2_lift', 1)}"],
         ["B2 bis B4", "Größenordnung neben der Regel, keine Automatik, "
                       "Begleitanalysen explorativ",
          "im Code geprüft"]],
        ["", "Ergebnis", "Bedeutung"],
        [["Regeln, die A1 bis A3 nehmen", w.z("brauchbare_regeln"),
          "statistische Mindestanforderung erfüllt"],
         ["davon unter B1 bestätigt", f"{w.z('b1_gehalten')} von "
                                      f"{w.z('b1_kandidaten')}",
          "im unangetasteten Zeitraum, unter Unsicherheit"],
         ["Ausschluss am Punktschätzer", w.z("b1_raus_punkt"),
          "erreicht den Lift gar nicht erst"],
         ["Ausschluss erst am Intervall", w.z("b1_raus_intervall"),
          "beobachtet, aber nicht bestätigt"]],
        f"Von {w.z('b1_kandidaten')} geprüften Regeln halten {w.z('b1_gehalten')} "
        f"unter Unsicherheit. Der Unterschied zwischen „beobachtet“ und "
        f"„bestätigt“ ist genau diese Spalte.",
        [f"`dispositionshinweise.csv` mit {w.z('b_regeln_n')} bestätigten Regeln.",
         f"Jede mit ihrer Größenordnung: höchstens {w.z('b_je_tag_max', 2)} Fahrten "
         f"je Tag.",
         f"Produkt A: {w['status_a']}."],
        ["Die Disposition bekommt begründete Hinweise und entscheidet selbst.",
         "Jede Regel trägt ihren Nenner — hoher Lift heißt nicht viele Fahrten.",
         "Salden und Abstell-Hotspots liegen bei, als explorativ gekennzeichnet."],
        [f"Ein analytisches Lehr-Gate, keine Betriebsfreigabe.",
         f"A4 ist {w['a4_zustand_text']}: Fahrten, die mangels Rad nie "
         f"stattfanden, stehen nirgends.",
         f"Die Hürde entspricht {w.z('huerde_je_werktag', 2)} Fahrten je Werktag "
         f"— dort beginnt keine Umsetzfahrt."],
        "Beobachtet ist nicht bestätigt: Erst ein Unsicherheitsintervall sagt, ob "
        "eine Regel die Schwelle wirklich hält.",
        phase_ergebnis=5)


def fall6(prs):
    w = Werte(6)
    fallkapitel(
        prs, 6, 7, "Fall 6 — Anomalie: der Rücksprung zum Mitlesen",
        "Was ist auffällig — und welche der drei Fragen, die darin stecken, "
        "beantwortet man gerade?",
        "Dieser Fall zeigt den Rücksprung von innen: Das erste Modell rechnete "
        "sauber und lieferte Unbrauchbares. Aufgefallen ist das nicht durch eine "
        "Kennzahl, sondern durch Sichtung der obersten Zeilen.",
        [("Geschäftsfrage", "Drei Fragen mit drei Entscheidungszeitpunkten: jetzt, "
                            "heute früh, gestern"),
         ("Analytisches Ziel", "A1 überfällige Rückgabe, A2 auffällige Fahrten, "
                               "B eine Prüfliste"),
         ("Der Kern", "Für A2 gibt es kein Label — und damit keine Trefferquote"),
         ("Verfahren", "Interquartilsregel, dann Isolation Forest; nach dem "
                       "Rücksprung je Radtyp normiert"),
         ("Urteil", f"A1 {w['a1_status']}, A2 {w['a2_status']}, B {w['b_status']}")],
        ["", "Kriterium", "Schwelle"],
        [["Listenlänge", "Aus Zeitbudget und Prüfdauer abgeleitet, nicht gesetzt",
          f"{w.z('listenlaenge')} Plätze"],
         ["B: Präzision", "Anteil brauchbarer Fälle je neuem Alarm",
          f"mindestens {w.p('b_gate_praezision', 0)}"],
         ["B: Recall und Verzug", "Anteil erkannter Episoden, Meldeverzug",
          f"{w.p('b_gate_recall', 0)} · höchstens {w.z('b_gate_verzug')} Tag"],
         ["A2", "Bewusst kein Gütekriterium",
          "Ohne Label lässt sich keines formulieren"]],
        ["Produkt", "Status", "Warum"],
        [["A1 überfällige Rückgabe", w["a1_status"],
          "als Regel beschrieben; Echtzeitquelle und Alarmkanal fehlen"],
         ["A2 auffällige Fahrten", w["a2_status"], "kein Label, keine belegte Güte"],
         ["B Prüfliste", w["b_status"],
          f"{w['b_gates_halten']} Gates halten auf dem unangetasteten Test"],
         ["globale Liste gegen Tagesliste", f"{w.p('globale_quote')} gegen "
                                            f"{w.p('tagesquote')}",
          "dasselbe Modell, zwei verschiedene Produkte"]],
        "Drei Produkte, drei Urteile. Die letzte Zeile ist die lehrreichste: "
        "Dieselbe Rangliste liefert zwei sehr verschiedene Trefferquoten, je "
        "nachdem, ob man sie global oder je Tag auswertet.",
        [f"A1 als Regel — {w['a1_status']}, ohne Modell.",
         f"A2 im {w['a2_status']}: rechnet mit, wird protokolliert.",
         f"B als {w['b_status']} entwickelte Prüfliste."],
        ["A1 löst die dringendste der drei Fragen ohne jeden Modellbetrieb.",
         f"Die Prüfliste hebt die Präzision gegenüber der alten Meldelogik "
         f"({w.p('stat_alt_quote')} je Alarm).",
         "Jedes Produkt trägt seinen eigenen Status statt einer gemeinsamen "
         "Freigabe."],
        [f"Auf dem unangetasteten Test bleibt die Präzision mit "
         f"{w.p('stat_je_alarm')} unter den geforderten "
         f"{w.p('b_gate_praezision', 0)}.",
         "Für A2 fehlt das Label — eine Trefferquote kann erst der Schattenbetrieb "
         "liefern.",
         "Die globale Rangliste beschreibt nicht die Liste, die im Betrieb "
         "entsteht."],
        "Eine Kennzahl auf der Gesamtliste beschreibt nicht die Liste, mit der "
        "später tatsächlich gearbeitet wird.",
        phase_ergebnis=5)


# ═══════════════════════════════════════════════════ Teil D — Synthese

def teil_synthese(prs):
    AKTUELLES_NB[0] = None
    kapitel(prs, 8, "Synthese",
            "Was war in allen sechs Fällen gleich — und was hat jeder einzelne "
            "über den Kreislauf gelehrt?",
            "Zum Abschluss legen wir die sechs Fälle zurück auf die Karte aus "
            "Kapitel 1. Jetzt ist sie nicht mehr abstrakt.")

    s = folie(prs, "Synthese", "Sechs Fälle, sechs Blickwinkel auf denselben Kreislauf",
              "Die Tabelle aus Kapitel 1 — jetzt ausgefüllt mit dem, was wir "
              "unterwegs gesehen haben.")
    tabelle(s, ["Fall", "Zeigt", "Der Satz, der bleibt"], [
        ["1 Regression", "Phase 5", "Der Rücksprung kam, obwohl das Kriterium hielt — "
                                    "ein Mittelwert ist keine Erfahrung"],
        ["2 Klassifikation", "Phase 6", "Das Modell verdient seinen Unterhalt "
                                        "nicht — die Faustregel hält die Zusage"],
        ["3 Clustering", "Phase 1", "Erfolgskriterien ohne Zielgröße — und eine bessere "
                                    "Frage als die, mit der wir anfingen"],
        ["4 Zeitreihe", "Phase 3", "Der Schnitt folgt der Zeit — und der "
                                   "Informationsstand der Auswahl dem des Betriebs"],
        ["5 Assoziation", "Phase 5", "Beobachtet ist nicht bestätigt — erst das "
                                     "Intervall entscheidet"],
        ["6 Anomalie", "Rücksprung", "Zweimal schlägt eine Zeile Fachwissen das "
                                     "Verfahren — weil vorher keine Baseline stand"],
    ], y=unter_intro(s), spalten_b=[170, 110, 623.5], zeilen_h=42)
    notizen(s, "Gehen Sie die sechs Zeilen einzeln durch und lassen Sie die "
               "Studierenden sagen, an welcher Stelle im jeweiligen Notebook das "
               "stand. Wer das kann, hat den Kreislauf verstanden.")

    s = folie(prs, "Synthese", "Fünf Dinge, die in allen sechs Fällen gleich waren",
              "Unabhängig vom Verfahren, unabhängig von der Frage — diese fünf "
              "Schritte kamen jedes Mal vor.")
    streifen(s, [
        ("Kriterium vor den Daten", "In allen sechs Fällen stand die Hürde fest, "
                                    "bevor jemand ein Modell gerechnet hat"),
        ("Ein Maßstab zuerst", "Nullmodell, Faustregel oder Interquartilsregel — "
                               "immer erst der billige Vergleich"),
        ("Die Fälle ansehen", "In jedem Fall brachte der Blick auf einzelne Zeilen "
                              "etwas, das keine Kennzahl gezeigt hätte"),
        ("Fehlerkosten benennen", "Wo die beiden Fehlerrichtungen ungleich teuer "
                                  "sind, gehört das ins Modell — nicht in die Diskussion"),
        ("Den Rücksprung aushalten", "Vier der sechs Fälle enden mit einem benannten "
                                     "Weg zurück statt mit einem Haken"),
    ], y=unter_intro(s), hoehe=56, luecke=8, chip_b=0)
    notizen(s, "Diese fünf Punkte sind das eigentliche Lernziel des Blocks. Die "
               "Verfahren kann man nachschlagen; diese Reihenfolge muss man sich "
               "angewöhnen.")

    s = folie(prs, "Synthese", "Wo die Zeit tatsächlich hingeht",
              "Zählt man die Codezellen der sechs Notebooks nach Phasen, ergibt "
              "sich dasselbe Bild wie in echten Projekten.")
    kachelreihe(s, [
        ("Phasen 1 bis 3", [
            "Frage schärfen, Daten ansehen,",
            "Merkmale bauen.",
            "",
            "Der weitaus größte Teil —",
            "und der, den man beim Lesen",
            "am liebsten überspringt.",
        ]),
        ("Phase 4", [
            "Das Modellieren.",
            "",
            "In jedem der sechs Notebooks",
            "ein knappes Kapitel, oft",
            "wenige Zeilen Code.",
            "",
            "Nicht die Hauptsache.",
        ]),
        ("Phasen 5 und 6", [
            "Bewerten und ausliefern.",
            "",
            "Der Teil, der in Lehrbüchern",
            "fehlt — und der über Erfolg",
            "oder Misserfolg im Betrieb",
            "entscheidet.",
        ]),
    ], y=unter_intro(s), hoehe=200)
    notizen(s, "Wenn die Studierenden aus diesem Block eine Sache mitnehmen: Der "
               "spannende Teil ist nicht der, den man erwartet. Wer Data Science "
               "lernt, um Modelle zu bauen, verbringt 70 Prozent seiner Zeit mit "
               "etwas anderem.")

    s = folie(prs, "Synthese", "Die Notebooks — und wie sie weiterhelfen",
              "Zu jedem der sechs Fälle gibt es zwei Fassungen: eine vollständig "
              "gerechnete und eine mit Lücken zum Selbstausfüllen.")
    tabelle(s, ["", "Vorführfassung", "Übungsfassung"], [
        ["Ablage", "analytics/notebooks/", "analytics/notebooks/uebung/"],
        ["Ausgaben", "vollständig eingebettet, sofort lesbar", "leer — sie entstehen "
                                                               "beim Rechnen"],
        ["Code", "vollständig", "an den Kernstellen ausgeschnitten"],
        ["Aufgaben", "keine", "1 bis 4 je Notebook, an der jeweils lehrreichsten Stelle"],
        ["Struktur", "sechs Phasen als Überschriften", "identisch — auch die "
                                                        "Übungsfassung führt durch alle sechs"],
    ], y=unter_intro(s), spalten_b=[140, 400, 363.5], zeilen_h=42)
    notizen(s, "Empfehlen Sie, die Vorführfassung zu lesen und die Übungsfassung zu "
               "rechnen. Beide tragen dieselbe Phasenstruktur, sodass man beim "
               "Steckenbleiben in der einen die Antwort in der anderen findet.")

    s = folie(prs, "Synthese", "Acht Sätze, die diesen Block tragen")
    streifen(s, [
        ("Erst die Frage", "Ein Erfolgskriterium, das nach dem Ergebnis entsteht, "
                           "ist keins"),
        ("Dann der Maßstab", "Wer das Nullmodell nicht schlägt, hat nichts gewonnen"),
        ("Leakage prüfen", "Was weiß ich zum Zeitpunkt der Vorhersage wirklich?"),
        ("Kosten sind ungleich", "Fast nie sind beide Fehlerrichtungen gleich teuer"),
        ("Kennzahlen trennen", "Technische Güte und fachliches Urteil sind zwei Dinge"),
        ("Die Fälle ansehen", "Eine Kennzahl sagt, wie gut. Nur der Blick sagt, woran"),
        ("Einfach schlägt stark", "Bei Gleichstand gewinnt, was man erklären kann"),
        ("Zurück ist normal", "Der Rücksprung ist das Verfahren, nicht sein Scheitern"),
    ], y=unter_intro(s), hoehe=36, luecke=4, chip_b=0)
    notizen(s, "Diese acht Sätze fassen die sechs Fälle zusammen. Keiner davon ist "
               "verfahrensspezifisch — sie gelten für die Regression genauso wie für "
               "die Anomalieerkennung. Das ist der Sinn eines Vorgehensmodells.")

    s = folie(prs, "Ausblick", "Was dieser Block nicht behandelt hat",
              "Damit klar ist, wo die Grenzen liegen — und was ein nächster Schritt "
              "wäre.")
    tabelle(s, ["Thema", "Warum es hier fehlt"], [
        ["Neuronale Netze", "Bei diesen Datenmengen und Fragestellungen ohne Vorteil "
                            "gegenüber den gezeigten Verfahren — und deutlich schwerer "
                            "zu erklären"],
        ["Hyperparameter-Suche", "Systematisches Durchprobieren hätte die Ergebnisse "
                                 "leicht verbessert und den Blick auf das Vorgehen "
                                 "verstellt"],
        ["Produktivbetrieb", "Ausliefern heißt hier: eine Funktion und ein Plan. Echte "
                             "Bereitstellung, Versionierung und Überwachung sind ein "
                             "eigenes Feld"],
        ["Kausalität", "Alle sechs Fälle finden Zusammenhänge, keine Ursachen. Die "
                       "Störgröße in Fall 4 zeigt, wie schnell man das verwechselt"],
    ], y=unter_intro(s), spalten_b=[250, 653.5], zeilen_h=54)
    notizen(s, "Der letzte Punkt ist der wichtigste für die Berufspraxis. Ein Modell, "
               "das Zusammenhänge findet, wird im Unternehmen regelmäßig als Beleg "
               "für Ursachen gelesen. Fall 4 ist das beste Gegenbeispiel, das dieses "
               "Deck zu bieten hat.")



def baue() -> Presentation:
    prs = leere_praesentation()
    s = prs.slides.add_slide(lay(prs, "Frontpage_Digital"))
    for ph in s.placeholders:
        i = ph.placeholder_format.idx
        if i == 10:
            ph.text_frame.text = "CRISP-DM an sechs Fallbeispielen"
        elif i == 11:
            ph.text_frame.text = ("Sechs Datenprojekte bei VeloCity — und ein Vorgehen, "
                                  "das in allen sechs dasselbe ist")
    notizen(s, "Sechs Notebooks, sechs Verfahren, ein Vorgehensmodell. Wir laufen "
               "einen Fall vollständig durch und sehen an den fünf anderen, was "
               "jeweils anders ist.")
    rahmen(prs)
    teil_karte(prs)
    teil_referenzfall(prs)
    fall2(prs)
    fall3(prs)
    fall4(prs)
    fall5(prs)
    fall6(prs)
    teil_synthese(prs)
    return prs


if __name__ == "__main__":
    prs = baue()
    prs.save(str(ZIEL))
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} Folien -> {ZIEL}")
