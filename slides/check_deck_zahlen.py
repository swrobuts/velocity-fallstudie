#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Haelt die Zahlen eines Foliendecks gegen die Merkzettel der Notebooks.

WARUM ES DIESE PRUEFUNG BRAUCHT

Handout und Use-Case-Decks werden aus analytics/bau/werte/*.json erzeugt;
eine Abweichung ist dort nicht moeglich. velocity-crispdm.pptx dagegen
traegt abgetippte Zahlen. Als die Daten neu erzeugt wurden, ist das Deck
still falsch geworden - und ein falsches Foliendeck faellt niemandem auf,
weil es plausibel aussieht.

Diese Pruefung vergleicht ausgewaehlte, von Hand gegengelesene
Behauptungen mit dem gemessenen Wert. Sie ist bewusst KEINE
Vollpruefung aller Zahlen: Ein Deck enthaelt auch Konstanten und
Beispielwerte, die in keinem Merkzettel stehen. Was hier steht, ist
gepruefte Substanz.

Aufruf: python3 slides/check_deck_zahlen.py [pfad.pptx]
"""
from __future__ import annotations

import json
import pathlib
import sys

from pptx import Presentation

WURZEL = pathlib.Path(__file__).resolve().parent.parent
WERTE = WURZEL / "analytics" / "bau" / "werte"
STANDARD = WURZEL / "slides" / "velocity-crispdm.pptx"


def merkzettel():
    return {f.stem[:2]: json.loads(f.read_text(encoding="utf-8"))
            for f in WERTE.glob("*.json")}


# (Folie, Textbaustein im Deck, Notebook, Beschreibung, Funktion -> Istwert)
# Der Textbaustein wird auf der Folie gesucht; fehlt er, gilt die Zeile
# als erledigt (die Folie wurde offenbar schon korrigiert).
BEHAUPTUNGEN = [
    (45, "Zwei Kandidaten", "01", "Zahl der Kandidaten",
     lambda w: f"drei: {w['01']['zulaessige']}"),
    (99, "20 %", "05", "Anteil Rundtouren",
     lambda w: f"{w['05']['anteil_rundtouren']:.1%}"),
    (102, "32", "05", "Zahl der gefundenen Regeln",
     lambda w: f"{w['05']['b1_kandidaten']} Kandidaten, "
               f"{w['05']['brauchbare_regeln']} nehmen A1 bis A3"),
    (102, "0,99", "05", "Support der staerksten Regel",
     lambda w: f"{w['05']['top_support']:.2%} - ueber der Huerde "
               f"{w['05']['k1_support']:.0%}"),
    (103, "Keine Regel wird freigegeben", "05", "Freigabeurteil Produkt B",
     lambda w: w["05"]["status_b"]),
    (105, "1,8 Räder je Werktag", "05", "groesster Stationssaldo",
     lambda w: f"{w['05']['saldo_max']:+.2f}"),
    (105, "20 bis 40", "05", "Stationskapazitaet",
     lambda w: f"{w['05']['kap_min']:.0f} bis {w['05']['kap_max']:.0f}"),
    (105, "elf Räder", "05", "frei abgestellte Raeder je Werktag",
     lambda w: f"{w['05']['frei_raeder_tag']:.1f}"),
    (106, "1,8 Räder je Werktag", "05", "groesster Stationssaldo (Abschluss)",
     lambda w: f"{w['05']['saldo_max']:+.2f}"),
    (15, "Keine Regel nimmt beide", "05", "Kurzfassung Fall 5",
     lambda w: f"{w['05']['brauchbare_regeln']} Regeln nehmen A1 bis A3, "
               f"{w['05']['b_regeln_n']} sind bestaetigt"),
    (67, "71,7 %", "02", "Trefferquote im Testquartal",
     lambda w: f"Regel {w['02']['quote_regel']:.1%}, Wald {w['02']['quote_wald']:.1%}"),
]


def main(argv):
    pfad = pathlib.Path(argv[0]) if argv else STANDARD
    if not pfad.exists():
        print(f"Deck fehlt: {pfad}")
        return 2
    w = merkzettel()
    folien = list(Presentation(str(pfad)).slides)

    def text(nr):
        if nr > len(folien):
            return ""
        return "\n".join(s.text_frame.text for s in folien[nr - 1].shapes
                         if s.has_text_frame)

    funde = []
    for nr, baustein, nb, was, ist in BEHAUPTUNGEN:
        if baustein in text(nr):
            funde.append((nr, was, baustein, ist(w)))

    print(f"{pfad.name}: {len(folien)} Folien, {len(BEHAUPTUNGEN)} gepruefte Aussagen")
    if not funde:
        print("  Keine der bekannten Falschaussagen mehr im Deck.")
        return 0
    print(f"  {len(funde)} veraltete Aussage(n):\n")
    print(f"  {'Folie':>5}  {'Gegenstand':<38} {'im Deck':<26} gemessen")
    for nr, was, baustein, ist in funde:
        print(f"  {nr:>5}  {was:<38} {baustein:<26} {ist}")
    print("\n  Ein Foliendeck mit abgetippten Zahlen veraltet still. "
          "Handout und\n  Use-Case-Decks lesen dieselben Werte aus "
          "analytics/bau/werte/*.json.")
    return 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
