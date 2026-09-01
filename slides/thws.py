"""Bausteine für THWS-Foliendecks nach dem Kachel-Raster.

Rasterwerte und Palette stammen aus dem Skill bint-folie. Alle Maße in
Punkt; die Foliengröße ist 1024 x 576.
"""
from __future__ import annotations

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

# ----------------------------------------------------------------- Palette
BLAU     = RGBColor(0x00, 0x3E, 0x6E)
GELB     = RGBColor(0xFF, 0xB4, 0x14)
GRUEN    = RGBColor(0x8A, 0xB8, 0x33)
GRUEN_D  = RGBColor(0x55, 0x80, 0x1C)   # abgedunkelt für weisse Schrift
ORANGE   = RGBColor(0xED, 0x70, 0x04)
ROT      = RGBColor(0xBE, 0x23, 0x44)
ROT_A    = RGBColor(0xA3, 0x26, 0x38)   # Warnleiste
TUERKIS  = RGBColor(0x4A, 0xB5, 0xC4)
SAND     = RGBColor(0xF7, 0xF5, 0xEF)
SAND_D   = RGBColor(0xE3, 0xDE, 0xD1)
HAAR     = RGBColor(0xC9, 0xC3, 0xB4)
WEISS    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT     = RGBColor(0x40, 0x40, 0x40)
TEXT_SEK = RGBColor(0x55, 0x51, 0x48)
GRAU     = RGBColor(0x8A, 0x83, 0x78)

# ------------------------------------------------------------------ Raster
FLUCHT_L   = 90.8
FLUCHT_R   = 994.3
BREITE     = 903.5
SP3        = (90.8, 402.6, 715.1)
SP3_B      = 279.2
SP2        = (90.8, 558.8)
SP2_B      = 435.0
KOPF_X     = 39.0
INTRO_X    = 84.0
INTRO_Y    = 90.0
INTRO_B    = 910.0
ZONE_OBEN  = 176.0
ZONE_UNTEN = 494.0
QUELLE_Y   = 506.0


def _tf(shape, *, rand=(10, 8, 10, 8)):
    tf = shape.text_frame
    tf.word_wrap = True
    l, t, r, b = rand
    tf.margin_left, tf.margin_top = Pt(l), Pt(t)
    tf.margin_right, tf.margin_bottom = Pt(r), Pt(b)
    return tf


def _absatz(tf, text, *, groesse=14, farbe=TEXT, fett=False, erste=False,
            ausricht=PP_ALIGN.LEFT, mono=False, abstand=0):
    p = tf.paragraphs[0] if erste else tf.add_paragraph()
    p.alignment = ausricht
    p.line_spacing = 0.95
    if abstand:
        p.space_before = Pt(abstand)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(groesse)
    r.font.bold = fett
    r.font.color.rgb = farbe
    if mono:
        r.font.name = 'Consolas'
    return p


def rechteck(slide, x, y, b, h, *, fuell=None, linie=None, linien_b=0.75,
             geom=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(geom, Pt(x), Pt(y), Pt(b), Pt(h))
    if fuell is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fuell
    if linie is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = linie
        sh.line.width = Pt(linien_b)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


def leiste(slide, x, y, h, farbe=BLAU, breite=4):
    """Die 4-pt-Akzentleiste am linken Rand einer Karte."""
    return rechteck(slide, x, y, breite, h, fuell=farbe)


# ------------------------------------------------------------- Kopfbereich
def kopf(slide, kicker, titel, *, quelle=None, intro=None):
    """Setzt Kicker, Titel, Einleitung und Quellenzeile.

    Schriftgrößen und Farben der drei Kopfplatzhalter werden bewusst
    NICHT gesetzt: das Vorlagendeck erbt sie vom Master. Wer sie hier
    ueberschreibt, lässt das Deck vom übrigen Kursdesign abweichen.
    """
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text_frame.text = titel
        elif idx == 12:
            ph.text_frame.text = kicker
        elif idx == 13:
            ph.text_frame.text = quelle or ''
    if not intro:
        return INTRO_Y
    # Hoehe aus der Textlaenge ableiten. Eine fest vorgegebene Hoehe laesst
    # kurze Einleitungen unnoetig viel Platz beanspruchen - Platz, der auf
    # Diagrammfolien fehlt.
    # Derselbe Faktor wie im Pruefer, plus Reserve: eine zu knapp
    # bemessene Einleitung schiebt sich sonst ins Motiv darunter.
    je_zeile = int(INTRO_B / (14 * 0.55))
    zeilen = max(1, -(-len(intro) // je_zeile))
    hoehe = zeilen * 18 + 6
    box = slide.shapes.add_textbox(Pt(INTRO_X), Pt(INTRO_Y), Pt(INTRO_B), Pt(hoehe))
    tf = _tf(box, rand=(0, 0, 0, 0))
    _absatz(tf, intro, groesse=14, farbe=TEXT, erste=True, ausricht=PP_ALIGN.JUSTIFY)
    return INTRO_Y + hoehe


def notizen(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ------------------------------------------------------------------ Motive
def kachelreihe(slide, karten, *, y=ZONE_OBEN, hoehe=170, spalten=3):
    """Motiv 1: weisse Kacheln mit Akzentleiste. karten = [(titel, [zeilen])]"""
    xs, b = (SP3, SP3_B) if spalten == 3 else (SP2, SP2_B)
    for i, (titel, zeilen) in enumerate(karten[:spalten]):
        x = xs[i]
        rechteck(slide, x, y, b, hoehe, fuell=WEISS, linie=HAAR)
        leiste(slide, x, y, hoehe)
        box = slide.shapes.add_textbox(Pt(x + 14), Pt(y + 10), Pt(b - 26), Pt(hoehe - 20))
        tf = _tf(box, rand=(0, 0, 0, 0))
        _absatz(tf, titel, groesse=15, farbe=BLAU, fett=True, erste=True)
        for z in zeilen:
            _absatz(tf, z, groesse=13, farbe=TEXT, abstand=5)


def regel_streifen(slide, zeilen, *, y=ZONE_OBEN, hoehe=48, luecke=8, chip_b=250):
    """Motiv 2: nummerierte Regelkarten mit Beispiel-Chip rechts."""
    for i, (regel, wirkung, beispiel) in enumerate(zeilen):
        yy = y + i * (hoehe + luecke)
        rechteck(slide, FLUCHT_L, yy, BREITE, hoehe, fuell=WEISS, linie=HAAR)
        leiste(slide, FLUCHT_L, yy, hoehe)
        badge = rechteck(slide, FLUCHT_L + 14, yy + (hoehe - 26) / 2, 26, 26,
                         fuell=BLAU, geom=MSO_SHAPE.OVAL)
        tfb = _tf(badge, rand=(0, 0, 0, 0))
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tfb, str(i + 1), groesse=13, farbe=WEISS, fett=True, erste=True,
                ausricht=PP_ALIGN.CENTER)
        txt_b = BREITE - 60 - chip_b - 20
        box = slide.shapes.add_textbox(Pt(FLUCHT_L + 50), Pt(yy + 6), Pt(txt_b), Pt(hoehe - 12))
        tf = _tf(box, rand=(0, 0, 0, 0))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.line_spacing = 0.95
        r1 = p.add_run(); r1.text = regel + '  '
        r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = BLAU
        r2 = p.add_run(); r2.text = wirkung
        r2.font.size = Pt(13); r2.font.color.rgb = TEXT
        if beispiel:
            chip = rechteck(slide, FLUCHT_L + BREITE - chip_b - 12, yy + 9, chip_b, hoehe - 18,
                            fuell=SAND_D, geom=MSO_SHAPE.ROUNDED_RECTANGLE)
            tfc = _tf(chip, rand=(8, 3, 8, 3))
            tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
            _absatz(tfc, beispiel, groesse=13, farbe=TEXT_SEK, erste=True, mono=True)


def prozesskette(slide, start, schritte, ziel, *, y=250, hoehe=88):
    """Motiv 6: Startknoten, Chevrons, Zielknoten."""
    n = len(schritte)
    knoten_b = 130
    luecke = 6
    chev_b = (BREITE - 2 * knoten_b - (n + 1) * luecke) / n
    x = FLUCHT_L
    k = rechteck(slide, x, y, knoten_b, hoehe, fuell=SAND_D,
                 geom=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _tf(k, rand=(8, 6, 8, 6)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _absatz(tf, start, groesse=13, farbe=TEXT, fett=True, erste=True, ausricht=PP_ALIGN.CENTER)
    x += knoten_b + luecke
    for i, (kopfz, sub) in enumerate(schritte):
        c = rechteck(slide, x, y, chev_b, hoehe, fuell=WEISS, linie=HAAR,
                     geom=MSO_SHAPE.PENTAGON)
        tf = _tf(c, rand=(10, 6, 26, 6)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, f'{i + 1} · {kopfz}', groesse=13, farbe=BLAU, fett=True, erste=True)
        # Bewusst keine Unterzeile: bei mindestens 13 pt bleibt im Chevron
        # kein Platz dafür, und Schrift zu verkleinern ist keine Option.
        x += chev_b + luecke
    z = rechteck(slide, x, y, knoten_b, hoehe, fuell=BLAU,
                 geom=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _tf(z, rand=(8, 6, 8, 6)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _absatz(tf, ziel, groesse=13, farbe=WEISS, fett=True, erste=True, ausricht=PP_ALIGN.CENTER)


def code_kacheln(slide, links, rechts, *, y=ZONE_OBEN, hoehe=250):
    """Motiv 9: zwei Kacheln mit Code. links/rechts = (titel, [zeilen], farbe)"""
    for (titel, zeilen, akzent), x in ((links, SP2[0]), (rechts, SP2[1])):
        rechteck(slide, x, y, SP2_B, hoehe, fuell=WEISS, linie=HAAR)
        leiste(slide, x, y, hoehe, akzent)
        box = slide.shapes.add_textbox(Pt(x + 14), Pt(y + 9), Pt(SP2_B - 26), Pt(hoehe - 18))
        tf = _tf(box, rand=(0, 0, 0, 0))
        _absatz(tf, titel, groesse=14, farbe=akzent, fett=True, erste=True)
        for z in zeilen:
            _absatz(tf, z, groesse=13, farbe=TEXT, mono=True, abstand=2)


def tabelle(slide, kopfzeile, zeilen, *, y=ZONE_OBEN, spalten_b=None, zeilen_h=26):
    """Motiv 10: Kopfzeile blau, Zebra, erste Spalte fett."""
    n = len(kopfzeile)
    if spalten_b is None:
        spalten_b = [BREITE / n] * n
    x = FLUCHT_L
    for j, (kopf_txt, b) in enumerate(zip(kopfzeile, spalten_b)):
        z = rechteck(slide, x, y, b, zeilen_h, fuell=BLAU)
        tf = _tf(z, rand=(8, 3, 8, 3)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, kopf_txt, groesse=13, farbe=WEISS, fett=True, erste=True)
        x += b
    for i, zeile in enumerate(zeilen):
        x = FLUCHT_L
        yy = y + zeilen_h + i * zeilen_h
        for j, (wert, b) in enumerate(zip(zeile, spalten_b)):
            z = rechteck(slide, x, yy, b, zeilen_h,
                         fuell=SAND_D if i % 2 else SAND, linie=HAAR, linien_b=0.5)
            tf = _tf(z, rand=(8, 3, 8, 3)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _absatz(tf, wert, groesse=13, farbe=BLAU if j == 0 else TEXT,
                    fett=(j == 0), erste=True,
                    mono=wert.startswith(('SELECT', 'select', 'CHECK', 'check', 'EXCLUDE')))
            x += b


def schichtenstapel(slide, schichten, *, y=ZONE_OBEN, hoehe=46, luecke=10, breite=None):
    """Motiv 4: gestapelte Balken, tragende Ebene blau."""
    b = breite or BREITE
    for i, (text, tragend) in enumerate(schichten):
        yy = y + i * (hoehe + luecke)
        s = rechteck(slide, FLUCHT_L, yy, b, hoehe,
                     fuell=BLAU if tragend else WEISS,
                     linie=None if tragend else HAAR)
        tf = _tf(s, rand=(16, 4, 16, 4)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, text, groesse=14, farbe=WEISS if tragend else TEXT,
                fett=tragend, erste=True)
        if i < len(schichten) - 1:
            rechteck(slide, FLUCHT_L + b / 2 - 7, yy + hoehe + 1, 14, luecke - 2,
                     fuell=HAAR, geom=MSO_SHAPE.ISOSCELES_TRIANGLE).rotation = 180


def ampel_matrix(slide, kopfzeile, zeilen, *, y=ZONE_OBEN, zeilen_h=46, luecke=8,
                 chip_b=92, label_b=300):
    """Motiv 8: Zeilenkarten mit Haken- und Kreuz-Chips."""
    n_chips = len(kopfzeile)
    rest_x = FLUCHT_L + label_b + 14
    for j, k in enumerate(kopfzeile):
        box = slide.shapes.add_textbox(Pt(rest_x + j * (chip_b + 10)), Pt(y - 26),
                                       Pt(chip_b), Pt(24))
        tf = _tf(box, rand=(0, 0, 0, 0))
        _absatz(tf, k, groesse=13, farbe=TEXT_SEK, fett=True, erste=True,
                ausricht=PP_ALIGN.CENTER)
    for i, (label, werte, umsetzung) in enumerate(zeilen):
        yy = y + i * (zeilen_h + luecke)
        rechteck(slide, FLUCHT_L, yy, BREITE, zeilen_h, fuell=WEISS, linie=HAAR)
        leiste(slide, FLUCHT_L, yy, zeilen_h)
        box = slide.shapes.add_textbox(Pt(FLUCHT_L + 16), Pt(yy + 5), Pt(label_b - 20),
                                       Pt(zeilen_h - 10))
        tf = _tf(box, rand=(0, 0, 0, 0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, label, groesse=13, farbe=BLAU, fett=True, erste=True)
        for j, w in enumerate(werte):
            chip = rechteck(slide, rest_x + j * (chip_b + 10), yy + (zeilen_h - 28) / 2,
                            chip_b, 28, fuell=GRUEN_D if w else ROT,
                            geom=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = _tf(chip, rand=(0, 0, 0, 0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            # U+2717 fehlt der Hausschrift und faellt auf ein Buchstaben-X
            # zurueck. U+00D7 ist in jeder Latin-Schrift vorhanden und wird
            # als Kreuz gelesen.
            _absatz(tf, '✓' if w else '×', groesse=16, farbe=WEISS, fett=True,
                    erste=True, ausricht=PP_ALIGN.CENTER)
        ux = rest_x + n_chips * (chip_b + 10) + 6
        box = slide.shapes.add_textbox(Pt(ux), Pt(yy + 5), Pt(FLUCHT_L + BREITE - ux - 10),
                                       Pt(zeilen_h - 10))
        tf = _tf(box, rand=(0, 0, 0, 0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, umsetzung, groesse=13, farbe=TEXT_SEK, erste=True)


def sandkarte(slide, titel, zeilen, *, y, hoehe=None, warnung=False, breite=None,
              x0=FLUCHT_L):
    b = breite or BREITE
    if hoehe is None:
        # Wie beim Band: die Höhe folgt dem Text, nicht umgekehrt.
        je_zeile = int((b - 30) / (13 * 0.52))
        n = max(1, -(-len(titel) // je_zeile))
        for z in zeilen:
            n += max(1, -(-len(z) // je_zeile))
        hoehe = max(80, n * 19 + 30)
    rechteck(slide, x0, y, b, hoehe, fuell=SAND, linie=HAAR)
    leiste(slide, x0, y, hoehe, ROT_A if warnung else BLAU)
    box = slide.shapes.add_textbox(Pt(x0 + 16), Pt(y + 10), Pt(b - 30), Pt(hoehe - 20))
    tf = _tf(box, rand=(0, 0, 0, 0))
    _absatz(tf, titel, groesse=15, farbe=ROT_A if warnung else BLAU, fett=True, erste=True)
    for z in zeilen:
        _absatz(tf, z, groesse=13, farbe=TEXT, abstand=5)


def sandband(slide, text, *, y, hoehe=None, mono=False):
    # Höhe aus der Textlänge ableiten: ein Band, das überläuft,
    # schneidet den Merksatz ab - und genau der soll hängenbleiben.
    if hoehe is None:
        je_zeile = int((BREITE - 32) / (14 * 0.52))
        zeilen = max(1, -(-len(text) // je_zeile))
        hoehe = max(52, zeilen * 20 + 26)
    rechteck(slide, FLUCHT_L, y, BREITE, hoehe, fuell=SAND_D)
    box = slide.shapes.add_textbox(Pt(FLUCHT_L + 16), Pt(y + 6), Pt(BREITE - 32), Pt(hoehe - 12))
    tf = _tf(box, rand=(0, 0, 0, 0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _absatz(tf, text, groesse=14, farbe=BLAU, fett=True, erste=True, mono=mono)


# --------------------------------------------------------------- Diagramme
def _png_masse(pfad: str) -> tuple[int, int]:
    """Liest Breite und Hoehe aus dem IHDR-Block einer PNG-Datei."""
    import struct
    with open(pfad, 'rb') as f:
        kopf = f.read(24)
    return struct.unpack('>II', kopf[16:24])


def diagramm(slide, pfad, *, y=ZONE_OBEN, hoehe=None, breite=None,
             x0=FLUCHT_L, rahmen=True, quelle=None):
    """Setzt ein gerendertes Mermaid-Diagramm mittig in die Inhaltszone.

    Das Bild wird proportional so skaliert, dass es in die verfuegbare
    Flaeche passt - Diagramme werden nie verzerrt und laufen nie ueber
    die Quellenzeile.
    """
    max_b = breite or BREITE
    max_h = hoehe or (ZONE_UNTEN - y)
    b_px, h_px = _png_masse(pfad)
    faktor = min(max_b / b_px, max_h / h_px)
    b, h = b_px * faktor, h_px * faktor
    x = x0 + (max_b - b) / 2
    yy = y + (max_h - h) / 2

    if rahmen:
        rechteck(slide, x - 8, yy - 8, b + 16, h + 16, fuell=WEISS, linie=HAAR)
    slide.shapes.add_picture(pfad, Pt(x), Pt(yy), Pt(b), Pt(h))

    if quelle:
        box = slide.shapes.add_textbox(Pt(x0), Pt(yy + h + 10), Pt(max_b), Pt(18))
        tf = _tf(box, rand=(0, 0, 0, 0))
        _absatz(tf, quelle, groesse=13, farbe=TEXT_SEK, erste=True)
    return yy + h


def leitfrage(slide, frage, *, y=ZONE_OBEN + 40):
    """Die Frage, die ein Kapitel beantwortet. Steht auf der Kapitelfolie."""
    hoehe = 92
    rechteck(slide, FLUCHT_L, y, BREITE, hoehe, fuell=SAND_D)
    leiste(slide, FLUCHT_L, y, hoehe, GELB, breite=6)
    box = slide.shapes.add_textbox(Pt(FLUCHT_L + 24), Pt(y + 14), Pt(BREITE - 48), Pt(hoehe - 28))
    tf = _tf(box, rand=(0, 0, 0, 0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _absatz(tf, 'Leitfrage', groesse=13, farbe=TEXT_SEK, fett=True, erste=True)
    _absatz(tf, frage, groesse=20, farbe=BLAU, fett=True, abstand=6)


def faden(slide, text, *, y=ZONE_UNTEN - 34):
    """Der rote Faden: der Bezug dieser Folie zu Annas Fahrt."""
    hoehe = 34
    rechteck(slide, FLUCHT_L, y, BREITE, hoehe, fuell=WEISS, linie=HAAR)
    leiste(slide, FLUCHT_L, y, hoehe, ROT_A)
    box = slide.shapes.add_textbox(Pt(FLUCHT_L + 16), Pt(y + 4), Pt(BREITE - 32), Pt(hoehe - 8))
    tf = _tf(box, rand=(0, 0, 0, 0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.line_spacing = 0.95
    r1 = p.add_run(); r1.text = 'Annas Fahrt:  '
    r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = ROT_A
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(13); r2.font.color.rgb = TEXT


def vorher_nachher(slide, links, rechts, *, y=ZONE_OBEN, hoehe=250):
    """Zwei Kacheln: der naive Versuch und die Loesung.

    links/rechts = (kopfzeile, titel, [zeilen], mono)
    """
    for (kopfzeile, titel, zeilen, mono), x, farbe in (
            (links,  SP2[0], ROT_A),
            (rechts, SP2[1], GRUEN_D)):
        rechteck(slide, x, y, SP2_B, hoehe, fuell=WEISS, linie=HAAR)
        kopf_h = 30
        rechteck(slide, x, y, SP2_B, kopf_h, fuell=farbe)
        kbox = slide.shapes.add_textbox(Pt(x + 12), Pt(y + 4), Pt(SP2_B - 24), Pt(kopf_h - 8))
        ktf = _tf(kbox, rand=(0, 0, 0, 0)); ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(ktf, kopfzeile, groesse=13, farbe=WEISS, fett=True, erste=True)

        box = slide.shapes.add_textbox(Pt(x + 14), Pt(y + kopf_h + 10),
                                       Pt(SP2_B - 26), Pt(hoehe - kopf_h - 20))
        tf = _tf(box, rand=(0, 0, 0, 0))
        _absatz(tf, titel, groesse=14, farbe=farbe, fett=True, erste=True)
        for z in zeilen:
            _absatz(tf, z, groesse=13, farbe=TEXT, mono=mono, abstand=3)


PHASEN = ('Business\nUnderstanding', 'Data\nUnderstanding', 'Data\nPreparation',
          'Modeling', 'Evaluation', 'Deployment')


def phasenleiste(slide, aktuell, *, y=ZONE_UNTEN - 40, rueckspruenge=()):
    """Motiv 11: Wo im CRISP-DM-Kreislauf stehen wir gerade?

    Der wiederkehrende Anker des Analytics-Decks - dieselbe Rolle, die im
    Datenbankdeck der rote Faden zu Annas Fahrt spielt. Sechs Felder,
    das laufende ausgefuellt.

    aktuell:       1..6, oder 0 fuer "keine Phase" (Uebersichtsfolien)
    rueckspruenge: Phasennummern, auf die dieser Fall zurueckgesprungen
                   ist. Sie bekommen den Warnton - der Rueckschritt ist
                   das, was CRISP-DM von einer Liste unterscheidet, und
                   er soll sichtbar sein, nicht nur erwaehnt.
    """
    hoehe = 40
    luecke = 5
    feld_b = (BREITE - 5 * luecke) / 6
    for i, name in enumerate(PHASEN):
        nr = i + 1
        x = FLUCHT_L + i * (feld_b + luecke)
        if nr == aktuell:
            fuell, schrift, rand = BLAU, WEISS, None
        elif nr in rueckspruenge:
            fuell, schrift, rand = SAND_D, ROT_A, ROT_A
        else:
            fuell, schrift, rand = WEISS, GRAU, HAAR
        rechteck(slide, x, y, feld_b, hoehe, fuell=fuell, linie=rand)
        box = slide.shapes.add_textbox(Pt(x + 6), Pt(y + 3), Pt(feld_b - 12), Pt(hoehe - 6))
        tf = _tf(box, rand=(0, 0, 0, 0))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Der Name steht zweizeilig; 13 pt ist die Untergrenze des Skills,
        # deshalb wird die Zeile nicht kleiner gesetzt, sondern umbrochen.
        _absatz(tf, f'{nr} {name.splitlines()[0]}', groesse=13, farbe=schrift,
                fett=(nr == aktuell), erste=True, ausricht=PP_ALIGN.CENTER)
        rest = name.splitlines()[1] if '\n' in name else ''
        if rest:
            _absatz(tf, rest, groesse=13, farbe=schrift, fett=(nr == aktuell),
                    ausricht=PP_ALIGN.CENTER)


def steckbrief(slide, zeilen, *, y=ZONE_OBEN, hoehe=None):
    """Motiv 12: Der Fall auf einen Blick - Frage, Daten, Verfahren, Urteil.

    Zwei Spalten Begriff/Wert. Steht am Anfang jedes Fallkapitels, damit
    die Studierenden vor dem Durchlaufen wissen, wohin die Reise geht.
    """
    zeilen_h = 34
    hoehe = hoehe or (len(zeilen) * zeilen_h + 16)
    rechteck(slide, FLUCHT_L, y, BREITE, hoehe, fuell=SAND, linie=HAAR)
    leiste(slide, FLUCHT_L, y, hoehe, BLAU, breite=6)
    begriff_b = 190
    for i, (begriff, wert) in enumerate(zeilen):
        zy = y + 8 + i * zeilen_h
        b = slide.shapes.add_textbox(Pt(FLUCHT_L + 20), Pt(zy), Pt(begriff_b), Pt(zeilen_h - 4))
        tf = _tf(b, rand=(0, 0, 0, 0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, begriff, groesse=13, farbe=BLAU, fett=True, erste=True)
        w = slide.shapes.add_textbox(Pt(FLUCHT_L + 20 + begriff_b), Pt(zy),
                                     Pt(BREITE - begriff_b - 44), Pt(zeilen_h - 4))
        tf = _tf(w, rand=(0, 0, 0, 0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _absatz(tf, wert, groesse=13, farbe=TEXT, erste=True)
