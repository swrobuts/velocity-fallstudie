#!/usr/bin/env python3
"""Prueft ein THWS-Deck auf Rasterverstöße und Überlauf.

Aufruf:
    python3 slides/check_deck.py slides/velocity-datenbankentwurf.pptx

Geprüft wird, was sich ohne Rendern feststellen lässt: Shapes ausserhalb
der Inhaltszone, Unterschreitung der Mindestschriftgröße, fehlende
Vortragsnotizen und geschätzte Textueberlaeufe.
"""
from __future__ import annotations

import sys
from pptx import Presentation
from pptx.util import Pt

ZONE_OBEN  = 176.0
ZONE_UNTEN = 494.0
QUELLE_Y   = 506.0
FOLIE_B    = 1024.0
MIN_PT     = 13.0
# grobe Schätzung: Zeichen je Zeile bei gegebener Breite und Größe
def zeilen_bedarf(text: str, breite_pt: float, groesse: float) -> int:
    if not text:
        return 0
    je_zeile = max(1, int(breite_pt / (groesse * 0.52)))
    zeilen = 0
    for absatz in text.split('\n'):
        zeilen += max(1, -(-len(absatz) // je_zeile))
    return zeilen


def main(pfad: str) -> int:
    prs = Presentation(pfad)
    befunde = []
    for i, s in enumerate(prs.slides, start=1):
        hat_notiz = bool(s.notes_slide.notes_text_frame.text.strip())
        if not hat_notiz:
            befunde.append((i, 'Notiz', 'keine Vortragsnotiz hinterlegt'))
        for sh in s.shapes:
            if sh.top is None or sh.height is None:
                continue
            unten = sh.top.pt + sh.height.pt
            rechts = sh.left.pt + sh.width.pt
            if unten > ZONE_UNTEN + 1 and sh.top.pt < ZONE_UNTEN:
                befunde.append((i, 'Zone',
                                f'Shape reicht bis y={unten:.0f}, Inhaltszone endet bei {ZONE_UNTEN:.0f}'))
            if rechts > FOLIE_B - 20:
                befunde.append((i, 'Rand', f'Shape reicht bis x={rechts:.0f}'))
            # Kopf- und Fussplatzhalter werden vom Master gestylt und in
            # PowerPoint automatisch angepasst. Sie hier zu schaetzen
            # erzeugt nur Rauschen.
            if sh.is_placeholder and sh.placeholder_format.idx in (0, 12, 13, 10, 11):
                continue
            if not sh.has_text_frame:
                continue
            groessen = [r.font.size.pt for p in sh.text_frame.paragraphs
                        for r in p.runs if r.font.size]
            if groessen and min(groessen) < MIN_PT:
                befunde.append((i, 'Schrift', f'{min(groessen):.0f} pt unter dem Minimum {MIN_PT:.0f} pt'))
            txt = sh.text_frame.text
            if txt and groessen:
                g = min(groessen)
                tf = sh.text_frame
                rand_x = (tf.margin_left.pt if tf.margin_left else 0) + \
                         (tf.margin_right.pt if tf.margin_right else 0)
                rand_y = (tf.margin_top.pt if tf.margin_top else 0) + \
                         (tf.margin_bottom.pt if tf.margin_bottom else 0)
                nötig = zeilen_bedarf(txt, sh.width.pt - rand_x, g) * g * 1.18 + rand_y
                if nötig > sh.height.pt + 6:
                    befunde.append((i, 'Überlauf',
                                    f'~{nötig:.0f} pt Text in {sh.height.pt:.0f} pt Shape: '
                                    f'„{txt[:44].strip()}…“'))
        # Ueberlappung zwischen Inhaltsformen. Genau diese Fehlerklasse
        # war der Grund, warum ein Sandband die letzte Tabellenzeile
        # verdeckte, ohne dass eine Hoehenpruefung angeschlagen haette.
        inhalt = [sh for sh in s.shapes
                  if not sh.is_placeholder and sh.top is not None
                  and sh.top.pt >= 80]
        gemeldet = set()
        for a_i, a in enumerate(inhalt):
            for b in inhalt[a_i + 1:]:
                ax1, ay1 = a.left.pt, a.top.pt
                ax2, ay2 = ax1 + a.width.pt, ay1 + a.height.pt
                bx1, by1 = b.left.pt, b.top.pt
                bx2, by2 = bx1 + b.width.pt, by1 + b.height.pt
                ux = min(ax2, bx2) - max(ax1, bx1)
                uy = min(ay2, by2) - max(ay1, by1)
                if ux <= 2 or uy <= 2:
                    continue
                # Textfelder liegen absichtlich in ihren Karten, Leisten am
                # Kartenrand, Badges und Chips ebenfalls. Gemeldet wird nur,
                # wenn KEINE der beiden Formen die andere umschliesst.
                a_in_b = ax1 >= bx1 - 1 and ax2 <= bx2 + 1 and ay1 >= by1 - 1 and ay2 <= by2 + 1
                b_in_a = bx1 >= ax1 - 1 and bx2 <= ax2 + 1 and by1 >= ay1 - 1 and by2 <= ay2 + 1
                if a_in_b or b_in_a:
                    continue
                key = (round(max(ay1, by1)), round(uy))
                if key in gemeldet:
                    continue
                gemeldet.add(key)
                befunde.append((i, 'Ueberlappung',
                                f'zwei Formen ueberschneiden sich um {uy:.0f} pt bei y={max(ay1, by1):.0f}'))

    for nr, art, text in befunde:
        print(f'Folie {nr:2d}  {art:<9} {text}')
    print(f'\n{len(prs.slides._sldIdLst)} Folien, {len(befunde)} Befund(e).')
    return 1 if befunde else 0


if __name__ == '__main__':
    raise SystemExit(main(sys.argv[1] if len(sys.argv) > 1 else 'slides/velocity-datenbankentwurf.pptx'))
