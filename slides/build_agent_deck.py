#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Drei Folien: KI-Agent und Warenwirtschaft.

WOZU EIN EIGENES DECK

velocity-datenbankentwurf.pptx traegt das Thema seit Kapitel 12 in fuenf
Folien, eingebettet in den Bogen von Annas Fahrt. Diese drei hier stehen
fuer sich: ein Schaubild, der technische Aufwand, der Nachweis. Sie
setzen den Rest der Einheit NICHT voraus und lassen sich einzeln
verwenden.

Gebaut mit slides/thws.py wie das grosse Deck - dieselbe Formsprache,
dieselben Bausteine, dieselbe Pruefung durch slides/check_deck.py.

    python3 slides/build_agent_deck.py
"""
from __future__ import annotations

import pathlib
import sys

from pptx import Presentation

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent))
from thws import (  # noqa: E402
    diagramm, kachelreihe, kopf, notizen, regel_streifen, sandband, sandkarte,
)

WURZEL = pathlib.Path(__file__).resolve().parent.parent


def _serverkennzahlen() -> dict[str, int]:
    """Zaehlt mcp/server.py aus, statt die Zahlen hier zu fuehren.

    Am 06.09.2026 standen auf dieser Folie "20 Werkzeuge: 4 lesend, 16
    aendernd" und "486 Zeilen Python". Beim Nachmessen war die Zeilenzahl
    laengst 603 und stieg im selben Lauf auf 690, und die Werkzeugzahl
    wurde durch ein neues Werkzeug falsch. Eine Zahl ueber fremden Code,
    die von Hand hier steht, veraltet - dieses Projekt hat das an einem
    Tag elfmal erlebt.

    Getrennt wird ueber den Syntaxbaum, nicht ueber Textsuche: ein
    Werkzeug gilt als aendernd, wenn es _rpc() tatsaechlich AUFRUFT. Beim
    ersten Versuch habe ich auf die Zeichenkette "_rpc(" geprueft und
    serverstand() faelschlich als aendernd gezaehlt, weil dessen Quelltext
    diese Zeichenkette in einem Vergleich fuehrt.
    """
    import ast
    quelle = (WURZEL / "mcp" / "server.py").read_text(encoding="utf-8")
    baum = ast.parse(quelle)
    lesend = aendernd = 0
    for knoten in baum.body:
        if not isinstance(knoten, ast.FunctionDef):
            continue
        if not any(isinstance(d, ast.Call) and getattr(d.func, "attr", "") == "tool"
                   for d in knoten.decorator_list):
            continue
        ruft_rpc = any(isinstance(n, ast.Call) and getattr(n.func, "id", "") == "_rpc"
                       for n in ast.walk(knoten))
        if ruft_rpc:
            aendernd += 1
        else:
            lesend += 1
    zeilen = quelle.count("\n") + 1
    begruendung = sum(1 for z in quelle.splitlines() if z.strip().startswith("#"))
    begruendung += sum(
        len(ast.get_docstring(n).splitlines())
        for n in ast.walk(baum)
        if isinstance(n, (ast.Module, ast.ClassDef, ast.FunctionDef))
        and ast.get_docstring(n))
    return {"lesend": lesend, "aendernd": aendernd, "werkzeuge": lesend + aendernd,
            "zeilen": zeilen, "anteil_begruendung": round(100 * begruendung / zeilen)}
ASSETS = WURZEL / "slides" / "assets"
# Dieselbe Vorlage wie das grosse Deck, aus build_deck.py bezogen statt
# ein zweites Mal eingetragen: Ein Pfad an zwei Stellen ist ein Pfad,
# der eines Tages an einer davon nicht mehr stimmt.
from build_deck import VORLAGE  # noqa: E402
ZIEL = WURZEL / "slides" / "velocity-agent-wawi.pptx"

Q = ("Fallstudie VeloCity, Schema velocity auf supabase.butscher.cloud. "
     "Quellen im Repository unter mcp/ und db/aufbau/.")


def bild(name: str) -> str:
    pfad = ASSETS / f"{name}.png"
    if not pfad.exists():
        raise SystemExit(f"Diagramm fehlt: {pfad}\nZuerst: bash tools/render_diagrams.sh")
    return str(pfad)


def leere_praesentation() -> Presentation:
    prs = Presentation(str(VORLAGE))
    liste = prs.slides._sldIdLst
    for sld in list(liste):
        prs.part.drop_rel(sld.rId)
        liste.remove(sld)
    return prs


def lay(prs, name):
    return next(l for l in prs.slide_layouts if l.name == name)


def folie(prs, kicker, titel, intro=None):
    s = prs.slides.add_slide(lay(prs, "Slide"))
    s._intro_unten = kopf(s, kicker, titel, quelle=Q, intro=intro)
    return s


def unter_intro(s, abstand=16):
    return getattr(s, "_intro_unten", 110) + abstand


def baue() -> Presentation:
    prs = leere_praesentation()

    # ═══════════════════════════════════════════════════ 1 Das Schaubild
    s = folie(prs, "KI und Warenwirtschaft",
              "Der Agent geht denselben Weg wie der Browser",
              "Beide Oberflächen sprechen die Datenbank ausschließlich über Sichten und "
              "api_-Funktionen an. Ein Programm kann denselben Weg gehen, ohne dass am "
              "Modell etwas zu ändern war. Der Schnitt entstand für zwei Browser mit einem "
              "Menschen davor — dass ein dritter Aufrufer ohne Umbau danebenpasst, war "
              "nicht geplant und ist deshalb der Beleg.")
    diagramm(s, bild("agent-wawi"), y=unter_intro(s), hoehe=250, rahmen=False)
    sandband(s, "Im Programm des Agenten steht keine einzige Rechteprüfung. Wer ihm "
                "weniger erlauben will, nimmt seinem Konto eine Rolle weg — nicht seiner "
                "Aufgabenbeschreibung.", y=420)
    notizen(s, "Der wichtigste Satz steht in der Mitte des Bildes: Die Entscheidung fällt "
               "in der Datenbank. Im Programm des Agenten gibt es keine Rechteprüfung — "
               "nicht, weil man sie weggelassen hätte, sondern weil dort keine hingehört. "
               "Wer dem Agenten weniger erlauben will, nimmt seinem Konto eine Rolle weg.")

    # ═══════════════════════════════════════════════════ 2 Der Aufwand
    s = folie(prs, "KI und Warenwirtschaft",
              "Am Datenmodell war nichts zu tun",
              "Die Anbindung kostete kein neues Schema, keine neue Tabelle und kein neues "
              "Recht. Der Aufwand lag daneben: ein Konto, ein Programm, ein Rückweg und "
              "eine Prüfung, die das Ganze grün hält.")
    _k = _serverkennzahlen()
    kachelreihe(s, [
        ("Am Modell: nichts", [
            "kein neues Schema",
            "keine neue Tabelle",
            "kein neues Recht",
            "keine Ausnahme für den Agenten",
            "keine Änderung an den Sichten",
        ]),
        ("Daneben: vier Dinge", [
            "ein Mitarbeiterkonto mit vier Fachrollen",
            "ein Programm, das anmeldet und weiterreicht",
            "Abzug und Rückspielung des Schemas",
            "eine Prüfung gegen den Systemkatalog",
        ]),
        ("Was messbar entstand", [
            f"{_k['werkzeuge']} Werkzeuge: {_k['lesend']} lesend, "
            f"{_k['aendernd']} ändernd",
            f"{_k['zeilen']} Zeilen Python, davon "
            f"{_k['anteil_begruendung']} % Begründung",
            "5 api_-Funktionen bewusst ausgelassen",
            "0 Zeilen Rechteprüfung im Programm",
        ]),
    ], y=unter_intro(s), hoehe=160)
    sandkarte(s, "Die fünf ausgelassenen Funktionen sind kein Sicherheitsargument", [
        "Sie handeln auf dem eigenen Kundensatz des Aufrufers, und ein Mitarbeiterkonto "
        "hat keinen. Sie wären nicht gefährlich, sondern wirkungslos — der Unterschied "
        "gehört benannt, sonst wird jedes Unbequeme zur Sicherheitsfrage.",
    ], y=372)
    notizen(s, "Die Zahl 0 in der dritten Kachel ist die eigentliche Aussage der Folie. "
               "Ein Programm ohne Rechteprüfung klingt nach einer Lücke und ist das "
               "Gegenteil: Es gibt dort nichts zu umgehen, weil dort nichts entschieden "
               "wird.")

    # ═══════════════════════════════════════════════════ 3 Nachweis
    s = folie(prs, "KI und Warenwirtschaft",
              "Wer ändern lässt, braucht ein Gedächtnis und einen Rückweg",
              "Ein Agent ist nur dann vorführbar, wenn hinterher nachvollziehbar ist, was "
              "er geändert hat — und wenn es zurückgeht. Das Modell führt dafür zwei "
              "getrennte Bücher; der Rückweg liegt außerhalb der Anwendung.")
    regel_streifen(s, [
        ("Änderungsprotokoll", "Kunde, Mitarbeiter, Station — wer, wann, welches Feld",
         "aenderungsprotokoll"),
        ("Lebenslaufakte", "Räder führen ihr Vorher und Nachher im Klartext mit",
         "fahrrad_ereignis"),
        ("Eigenes Konto", "Was ein Mensch tat und was ein Programm, bleibt trennbar",
         "personalnummer"),
        ("Weg zurück", "Ausmustern und Löschen sind endgültig — zurück nur über den Abzug",
         "pg_dump / pg_restore"),
    ], y=unter_intro(s), hoehe=54, luecke=10, chip_b=200)
    sandband(s, "Es gibt keine Rücknahmefunktion. Wer vorführen will, was ein Agent "
                "anrichten kann, sichert vorher das Schema.", y=428)
    notizen(s, "Die Aufteilung in zwei Bücher wirkt zufällig, ist es aber nicht: Den "
               "Protokolltrigger tragen drei Tabellen, weil dort Personendaten stehen, und "
               "die Sicht darauf lässt die Werte deshalb weg. Räder führen stattdessen "
               "eine Lebenslaufakte mit Vorher und Nachher im Klartext — ein Radstatus ist "
               "kein Personendatum. Der letzte Punkt ist der unbequeme: Es gibt keine "
               "Rücknahmefunktion. Wer vorführen will, was ein Agent anrichten kann, "
               "sichert vorher das Schema.")

    return prs


if __name__ == "__main__":
    prs = baue()
    prs.save(str(ZIEL))
    print(f"{len(prs.slides._sldIdLst)} Folien geschrieben nach {ZIEL}")
