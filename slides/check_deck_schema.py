#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Haelt die Objektzahlen in velocity-datenbankentwurf.pptx gegen db/aufbau.

WARUM ES DIESE PRUEFUNG BRAUCHT

Die Zahlen dieses Decks sind in slides/build_deck.py von Hand getippt.
Am 01.09.2026 kam die Preisschaetzung dazu - eine Tabelle, eine Sicht,
eine api_-Funktion - und vier Angaben im Deck wurden still falsch:
"27 Tabellen" (28), "Vier weitere api_" (fuenf), "16 Sichten" (17) und
die Spaltenzahl. Aufgefallen ist das erst zwei Tage spaeter bei einer
Pruefung von Hand.

Diese Pruefung liest die Zahlen AUS DEM DECK und zaehlt die Objekte im
Quelltext des Aufbaus nach. Sie prueft damit das ausgelieferte Ergebnis,
nicht die Absicht des Bauers.

ZUM BEZUGSRAHMEN

Gezaehlt wird, was db/aufbau/*.sql anlegt - nichts sonst. Die laufende
Datenbank enthaelt zusaetzlich, was die Skripte in db/betrieb/ erzeugen,
derzeit die Tabelle uebernahme_protokoll mit zehn Spalten. Eine Abfrage
gegen den Systemkatalog liefert deshalb hoehere Werte als diese Pruefung.
Der Bezug auf den Aufbau ist Absicht: er ist versioniert und damit
ueberhaupt pruefbar.

Der Spaltenzaehler ist gegen doku/datenmodell/06-data-dictionary.md
kalibriert - dieses wird aus dem Systemkatalog erzeugt. Fuer alle dort
erfassten Tabellen stimmen beide Zaehlungen ueberein.

Aufruf: python3 slides/check_deck_schema.py [pfad.pptx]
"""
from __future__ import annotations

import pathlib
import re
import sys

from pptx import Presentation

WURZEL = pathlib.Path(__file__).resolve().parent.parent
AUFBAU = WURZEL / "db" / "aufbau"
STANDARD = WURZEL / "slides" / "velocity-datenbankentwurf.pptx"

# Woerter, die am Anfang einer Zeile im create-table-Rumpf keine Spalte
# einleiten, sondern eine Tabellenbedingung.
KEINE_SPALTE = {"constraint", "primary", "foreign", "unique", "check",
                "exclude", "like", "references", "deferrable", "initially"}

ZAHLWORT = {"eine": 1, "zwei": 2, "drei": 3, "vier": 4, "fünf": 5, "sechs": 6,
            "sieben": 7, "acht": 8, "neun": 9, "zehn": 10, "elf": 11, "zwölf": 12}


# ------------------------------------------------------------ Quelltext lesen

def quelltext() -> str:
    """Alle Aufbaudateien, ohne Kommentare, mit neutralisierten Literalen.

    Die Literale muessen neutralisiert werden, weil eine Zeichenkette wie
    daterange(current_date, null, '[)') eine Klammer enthaelt - sie wuerde
    die Klammerzaehlung unten aus dem Tritt bringen.
    """
    q = "\n".join(p.read_text(encoding="utf-8") for p in sorted(AUFBAU.glob("*.sql")))
    q = re.sub(r"/\*.*?\*/", " ", q, flags=re.S)
    q = "\n".join(z.split("--")[0] for z in q.split("\n"))
    return re.sub(r"'(?:[^']|'')*'", lambda m: "'" + "_" * (len(m.group(0)) - 2) + "'", q)


def dateien(von: int, bis: int) -> str:
    q = "\n".join(p.read_text(encoding="utf-8") for p in sorted(AUFBAU.glob("*.sql"))
                  if von <= int(p.name[:4]) <= bis)
    return "\n".join(z.split("--")[0] for z in q.split("\n"))


def namen(quelle: str, muster: str) -> set[str]:
    return {m.group(1) for m in re.finditer(muster, quelle, re.I)}


TABELLE = r"create\s+table\s+(?:if\s+not\s+exists\s+)?velocity\.(\w+)"
SICHT = r"create\s+(?:or\s+replace\s+)?view\s+velocity\.(\w+)"
FUNKTION = r"create\s+(?:or\s+replace\s+)?function\s+velocity\.(\w+)"


def spalten_je_tabelle() -> dict[str, list[str]]:
    """{tabelle: [spalten]} aus create table und alter table ... add column."""
    q = quelltext()
    aus: dict[str, list[str]] = {}
    for m in re.finditer(TABELLE + r"\s*\(", q, re.I):
        i, tiefe = m.end(), 1
        while i < len(q) and tiefe:
            tiefe += (q[i] == "(") - (q[i] == ")")
            i += 1
        rumpf, tiefe, letzter, teile = q[m.end():i - 1], 0, 0, []
        for j, z in enumerate(rumpf):
            tiefe += (z == "(") - (z == ")")
            if z == "," and tiefe == 0:
                teile.append(rumpf[letzter:j])
                letzter = j + 1
        teile.append(rumpf[letzter:])
        aus[m.group(1)] = [w[0] for w in (t.split() for t in teile)
                           if len(w) >= 2 and re.fullmatch(r"[a-z_]\w*", w[0])
                           and w[0].lower() not in KEINE_SPALTE]
    for m in re.finditer(r"alter\s+table\s+velocity\.(\w+)(.*?);", q, re.I | re.S):
        for c in re.findall(r"add\s+column\s+(?:if\s+not\s+exists\s+)?(\w+)",
                            m.group(2), re.I):
            if c not in aus.setdefault(m.group(1), []):
                aus[m.group(1)].append(c)
    return aus


def gezaehlt() -> dict[str, int]:
    alles, sp = quelltext(), spalten_je_tabelle()
    spalten = sum(len(v) for v in sp.values())
    audit = sum(1 for v in sp.values() for c in v if c in ("erstellt_am", "geaendert_am"))
    return {
        "tabellen_a_f": len(namen(dateien(2, 7), TABELLE)),
        "tabellen": len(namen(alles, TABELLE)),
        "wawi_sichten": len([x for x in namen(alles, SICHT) if x.startswith("v_wawi_")]),
        "wawi_api": len([x for x in namen(dateien(17, 20), FUNKTION) if x.startswith("api_")]),
        "website_api": len([x for x in namen(dateien(1, 16), FUNKTION) if x.startswith("api_")]),
        "spalten": spalten,
        "audit": audit,
        "beschrieben": spalten - audit,
    }


# ------------------------------------------------------- Behauptungen im Deck
#
# (Beschreibung, Muster mit einer Gruppe je Zahl, Schluessel je Gruppe)
# Das Muster wird ueber alle Folien gesucht; findet es sich nirgends, gilt
# die Behauptung als geaendert und wird gemeldet - eine stillschweigend
# verschwundene Aussage soll auffallen.

BEHAUPTUNGEN = [
    ("Tabellen der Bereiche A bis F",
     r"Fachbereiche A bis F, (\d+) Tabellen", ["tabellen_a_f"]),
    ("Basistabellen des Aufbaus",
     r"(\d+) Basistabellen legt der Aufbau an", ["tabellen"]),
    ("Spalten und Audit-Spalten",
     r"(\d+) Spalten im Aufbau, (\d+) beschrieben", ["spalten", "beschrieben"]),
    ("technische Audit-Spalten",
     r"(\d+) technische Audit-Spalten", ["audit"]),
    ("Sichten und Funktionen der Warenwirtschaft",
     r"(\d+) Sichten, (\d+) Funktionen", ["wawi_sichten", "wawi_api"]),
    ("api_-Funktionen der Website",
     r"(\w+) weitere api_ dienen der Website", ["website_api"]),
]


def main(argv):
    pfad = pathlib.Path(argv[0]) if argv else STANDARD
    if not pfad.exists():
        print(f"Deck fehlt: {pfad}")
        return 2
    ist = gezaehlt()
    text = "\n".join(f.text_frame.text for folie in Presentation(str(pfad)).slides
                     for f in folie.shapes if f.has_text_frame)

    funde = []
    for was, muster, schluessel in BEHAUPTUNGEN:
        m = re.search(muster, text)
        if m is None:
            funde.append((was, "nicht gefunden", "-"))
            continue
        for wert, s in zip(m.groups(), schluessel):
            zahl = ZAHLWORT.get(wert.lower(), None) if not wert.isdigit() else int(wert)
            if zahl != ist[s]:
                funde.append((was, wert, ist[s]))

    print(f"{pfad.name}: {len(BEHAUPTUNGEN)} Aussagen gegen db/aufbau/*.sql")
    if not funde:
        print("  Alle Objektzahlen im Deck stimmen mit dem Aufbau überein.")
        return 0
    print(f"  {len(funde)} Abweichung(en):\n")
    print(f"  {'Gegenstand':<44} {'im Deck':<14} gezählt")
    for was, deck, zahl in funde:
        print(f"  {was:<44} {str(deck):<14} {zahl}")
    print("\n  Die Zahlen stehen von Hand in slides/build_deck.py. Nach jeder\n"
          "  Schemaänderung dort nachziehen und das Deck neu bauen.")
    return 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
