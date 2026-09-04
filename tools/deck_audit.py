#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Inhaltliche Pruefung des Foliendecks.

slides/check_deck.py prueft Geometrie: Zonen, Schriftgrade, Ueberlappung,
Ueberlauf, fehlende Notizen. Das sagt nichts darueber, ob auf den Folien
etwas Falsches steht. Dieser Pruefer sucht nach inhaltlichen Fehlern:

  UMLAUT     transliterierte Umlaute (fuer, waere, Massgeblich). Auf der
             Website steckten davon 268 Stueck; das Deck wird aus
             derselben Hand gespeist.
  VERBOTEN   Begriffe, die laut Hausregel nicht in ein Foliendeck
             gehoeren: Termine, Lehrformen, Werkzeugnamen ohne Stand,
             Aussagen ueber Pruefungsstoff.
  ABSOLUT    "objektiv", "neutral", "beweist" - Absolutheiten, die sich
             fachlich nicht halten lassen.
  ZEICHEN    Emoji und Zeichen ausserhalb der Hausschrift.
  ZAHL       Zahlenangaben, die von der Wirklichkeit im Repository
             abweichen (Testanzahl, Pruefungen, Tabellen, Hoehenmeter).
  TITEL      doppelte Folientitel.

Aufruf:  python3 tools/deck_audit.py [--text]
         --text gibt zusaetzlich den vollstaendigen Folientext aus.
"""
import re
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

WURZEL = Path(__file__).resolve().parent.parent
DECK = WURZEL / "slides" / "velocity-datenbankentwurf.pptx"

UMLAUT = re.compile(
    r'(?<![\w.])('
    r'[A-Za-zÄÖÜäöüß]*(?:ue|ae|oe|Ue|Ae|Oe)[A-Za-zÄÖÜäöüß]*'
    r')(?![\w])')
# Deutsche Woerter, in denen ue/ae/oe kein Umlaut ist.
UMLAUT_OK = {
    'quelle', 'quellen', 'quelltext', 'quellcode', 'aktuell', 'aktuelle',
    'aktuellen', 'aktueller', 'aktuellem', 'steuert', 'steuern', 'umsatzsteuer',
    'umsatzsteuersatz', 'dauer', 'zuerst', 'neue', 'neuen', 'neues', 'neuer',
    'unique', 'true', 'value', 'values', 'queue', 'sequence', 'due',
    'konsequenz', 'frequenz', 'sequenz', 'adaequat', 'niveau', 'museum',
    # Steigerungsformen von "genau" - die Folge aue steht hier ueber eine
    # Silbengrenze hinweg und ist kein transliteriertes ae.
    'genauer', 'genaue', 'genauen', 'genauere', 'genaueren', 'genauigkeit',
    # Verben auf -auen: die Folge aue laeuft ebenfalls ueber eine Silbengrenze.
    'bauen', 'gebauen', 'brauen', 'trauen', 'schauen', 'blauen', 'grauen',
}


def bezeichner():
    """Tabellen-, Spalten- und Funktionsnamen aus dem Systemkatalog.

    Die heissen bewusst ASCII (kapazitaet, verfuegbare_raeder, gueltigkeit).
    Sie zu raten geht schief, also werden sie gefragt.
    """
    sys.path.insert(0, str(WURZEL / 'db'))
    from run import verbinde
    cur = verbinde().cursor()
    cur.execute("""
        select table_name from information_schema.tables where table_schema='velocity'
        union select column_name from information_schema.columns where table_schema='velocity'
        union select routine_name from information_schema.routines where routine_schema='velocity'""")
    namen = set()
    for (n,) in cur.fetchall():
        namen.add(n.lower())
        namen.update(teil for teil in n.lower().split('_') if len(teil) > 2)

    # AUCH DIE BEZEICHNER DES QUELLTEXTS (ergaenzt 28.08.2026).
    # Der Katalog kennt nur Tabellen, Spalten und Funktionen der
    # Datenbank. Das Deck zitiert aber auch Bausteine der Oberflaeche -
    # etwa neuerVorgang(), die Klammer um jede Aenderung in der
    # Warenwirtschaft. Ohne diese Ergaenzung meldet der Umlautpruefer
    # solche Namen als transliterierten Umlaut ("neuerVorgang" enthaelt
    # die Folge eue) und verlangt eine Schreibweise, die den Code
    # zerlegen wuerde. Bezeichner bleiben ASCII - das ist die Hausregel,
    # nicht ihre Verletzung.
    for js in sorted(WURZEL.glob("wawi/*.js")) + sorted(WURZEL.glob("src/*.js")):
        for m in re.finditer(r"\bfunction\s+([A-Za-z_$][\w$]*)", js.read_text(encoding="utf-8")):
            namen.add(m.group(1).lower())
    return namen

VERBOTEN = [
    (r'\bTermin\s*\d', 'Bindung an einen Termin'),
    (r'\bPräsenz\b|\bonline\b(?!\w)', 'Bindung an eine Lehrform'),
    (r'\bZoom\b|\bBreakout', 'Bindung an ein Werkzeug der Lehre'),
    (r'Prüfungsstoff|klausurrelevant|\bKlausur\b|Gewichtung der Note',
     'Aussage ueber Pruefungsstoff'),
]
ABSOLUT = [
    (r'\bobjektiv\b', 'objektiv'), (r'\bneutral\b', 'neutral'),
    (r'\bbeweist\b', 'beweist'), (r'\bimmer sicher\b', 'immer sicher'),
    (r'\bgarantiert\b(?! nicht)', 'garantiert'),
]
# Codezeilen sind ausgenommen: SQL-Bezeichner heissen im ganzen Schema
# bewusst ASCII (verfuegbare_raeder, "alles fuer alle").
CODEZEILE = re.compile(
    r'\b(create|select|insert|update|alter|drop|grant|revoke|check|using|'
    r'policy|constraint|returns|language|begin|declare)\b', re.I)
EMOJI = re.compile('[\U0001F000-\U0001FAFF☀-➿️]')
ZEICHEN_OK = set('✓×ΣΨ▸■·–—…„“”‚‘’€•→')


def zahl_im_repo():
    """Was im Repository wirklich gilt - gegen die Folien gehalten."""
    def sh(cmd):
        return subprocess.run(cmd, shell=True, cwd=WURZEL,
                              capture_output=True, text=True).stdout.strip()
    return {
        'Testfunktionen': sh("grep -rhoE 'create or replace function velocity_test.test_[a-z0-9_]+' "
                             "db/tests/*.sql | sort -u | wc -l"),
        'Aufbaudateien':  sh("ls db/aufbau/*.sql | wc -l"),
        'Diagramme':      sh("ls doku/datenmodell/erd/*.mmd | wc -l"),
        'Abnahmepruefungen': sh("grep -c '^schritt ' tools/abnahme.sh"),
        'Basistabellen': sh("grep -rhoiE 'create table (if not exists )?velocity\\.' "
                            "db/aufbau/*.sql | wc -l"),
        # Seit Kapitel 12 nennt eine Folie die Zahl der Agentenwerkzeuge.
        # Sie steht im Server von Hand, eine je Funktion - genau die Sorte
        # Zahl, die beim naechsten hinzugefuegten Werkzeug veraltet.
        'Werkzeuge': sh("grep -c '^@server.tool()' mcp/server.py"),
    }


# Ausgeschriebene Zahlwoerter. Ohne sie bleibt "Zwoelf Aufbauschritte" unsichtbar,
# weil die Zahlpruefung nur nach Ziffern sucht - genau so ist eine veraltete
# Angabe ueber Monate im Deck stehen geblieben.
ZAHLWORT = {
    'zwei': '2', 'drei': '3', 'vier': '4', 'fuenf': '5', 'fünf': '5',
    'sechs': '6', 'sieben': '7', 'acht': '8', 'neun': '9', 'zehn': '10',
    'elf': '11', 'zwoelf': '12', 'zwölf': '12', 'dreizehn': '13',
    'vierzehn': '14', 'fuenfzehn': '15', 'fünfzehn': '15', 'sechzehn': '16',
    'siebzehn': '17', 'achtzehn': '18', 'neunzehn': '19', 'zwanzig': '20',
}


def main():
    if not DECK.exists():
        print(f"FEHLT  {DECK}"); return 1
    prs = Presentation(DECK)
    namen = bezeichner()
    befunde, titel_gesehen, volltext = [], {}, []

    for nr, folie in enumerate(prs.slides, 1):
        texte = []
        for sh in folie.shapes:
            if sh.has_text_frame:
                texte.append(sh.text_frame.text)
            if sh.has_table:
                for zeile in sh.table.rows:
                    texte.extend(z.text for z in zeile.cells)
        text = "\n".join(t for t in texte if t)
        volltext.append((nr, text))

        for zeile in text.split('\n'):
          if CODEZEILE.search(zeile):
            continue
          for w in UMLAUT.findall(zeile):
            k = w.lower()
            if k in UMLAUT_OK or k in namen or len(w) < 4:
                continue
            befunde.append((nr, 'UMLAUT', w))
        for muster, was in VERBOTEN:
            if re.search(muster, text, re.I):
                befunde.append((nr, 'VERBOTEN', was))
        for muster, was in ABSOLUT:
            if re.search(muster, text, re.I):
                befunde.append((nr, 'ABSOLUT', was))
        # Die Hausschrift kennt ✓ ✗ Σ ▸ ■ · und die Anfuehrungszeichen.
        # Alles andere jenseits von U+2000 ist erst einmal verdaechtig.
        for z in set(text):
            if ord(z) > 0x2000 and z not in ZEICHEN_OK:
                art = 'Emoji' if EMOJI.match(z) else 'Sonderzeichen'
                befunde.append((nr, 'ZEICHEN', f'{art} {z!r} U+{ord(z):04X}'))

        erste = text.split("\n")[0].strip() if text else ''
        if len(erste) > 8:
            if erste in titel_gesehen:
                befunde.append((nr, 'TITEL', f'wie Folie {titel_gesehen[erste]}: {erste[:50]}'))
            titel_gesehen[erste] = nr

    # Zahlen gegen das Repository
    ganz = "\n".join(t for _, t in volltext)
    ist = zahl_im_repo()
    wortzahl = '|'.join(ZAHLWORT)
    for name, wert in ist.items():
        # \b VOR die Alternative: ohne sie fand 'zwanzig' sich mitten in
        # 'Einundzwanzig' wieder und meldete 20 statt 21. Zahlwoerter
        # ueber zwanzig stehen nicht in ZAHLWORT - fuer die gehoert die
        # Ziffer auf die Folie, sonst prueft hier gar nichts mehr.
        muster = re.compile(r'\b(\d+|' + wortzahl + r')\s+' + name[:6], re.I)
        for m in muster.finditer(ganz):
            gesagt = ZAHLWORT.get(m.group(1).lower(), m.group(1))
            if gesagt != wert:
                befunde.append((0, 'ZAHL', f'{name}: Folie sagt {m.group(1)} '
                                           f'(= {gesagt}), im Repository sind es {wert}'))

    if '--text' in sys.argv:
        for nr, t in volltext:
            print(f'\n----- Folie {nr} -----\n{t}')

    for nr, art, was in befunde:
        ort = f'Folie {nr}' if nr else 'Deck'
        print(f'{art:9} {ort:10} {was}')
    print(f'\n{len(prs.slides)} Folien geprueft, {len(befunde)} Befund(e).')
    print('Repository-Istwerte: ' + ', '.join(f'{k}={v}' for k, v in ist.items()))
    return 1 if befunde else 0


if __name__ == '__main__':
    sys.exit(main())
