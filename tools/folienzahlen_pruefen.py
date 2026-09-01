#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Haelt jede Zahl im CRISP-DM-Deck gegen das Notebook, das die Folie zitiert.

Aufruf:
    python3 tools/folienzahlen_pruefen.py
    python3 tools/folienzahlen_pruefen.py slides/velocity-crispdm.pptx

Wozu
----
Das Deck ist eine Lesehilfe. Eine Folie, die eine Zahl nennt, die im
Notebook nicht mehr steht, schickt die Studierenden auf eine Suche, die
nicht enden kann - und beschaedigt genau das Vertrauen, das eine
Lesehilfe aufbauen soll.

Der Fall, der zu diesem Werkzeug gefuehrt hat: Nachdem der Lehrdatensatz
neu erzeugt worden war (die Stationen wurden an die Datenbank
angeglichen), aenderten sich Zahlen in vier Notebooks. Die Folien
behielten die alten. Von Hand faellt das nicht auf - es sind 120 Folien
und sechs Notebooks.

Wie geprueft wird
-----------------
Jede Inhaltsfolie traegt in der Fusszeile ihre Quelle, zum Beispiel
"analytics/notebooks/04_Zeitreihe_Nachfrageprognose.ipynb · Abschnitt 6.1".
Daraus liest dieses Werkzeug das zustaendige Notebook und sucht jede Zahl
der Folie in dessen Quelltext UND Ausgaben.

Zahlen ohne Quellenzeile werden nicht geprueft: Kapitelfolien und die
Karte in Teil A beziehen sich auf kein einzelnes Notebook.

Grenzen - bewusst benannt
-------------------------
Das Werkzeug findet Zahlen, die NIRGENDS im Notebook vorkommen. Es findet
nicht, wenn eine Folie eine Zahl aus dem falschen Zusammenhang nennt, und
es findet keine falschen Behauptungen ohne Zahl ("das Modell gewinnt").
Dafuer gibt es keinen Ersatz fuer das Lesen.
"""
from __future__ import annotations

import json
import pathlib
import re
import sys

from pptx import Presentation

WURZEL = pathlib.Path(__file__).resolve().parent.parent
NB = WURZEL / "analytics" / "notebooks"

# Zahlen, die zur Sprache gehoeren und nicht zu den Daten: Phasennummern,
# Prozentmarken aus Erfolgskriterien, Jahreszahlen, Tarifkonstanten.
HARMLOS = {
    "0", "1", "2", "3", "4", "5", "6", "7", "8", "9", "10", "12", "13", "15", "16",
    "20", "24", "25", "26", "30", "35", "40", "42", "45", "50", "60", "70", "80",
    "90", "95", "100", "180", "288", "1999", "2026",
    "0,5", "1,0", "1,3", "1,5", "2,0", "2,5", "5,0", "0,50", "1,00",
}

# Schwellen und Konstanten, die in JEDEM Notebook vorkommen duerfen, ohne
# dass sie dort ausgedruckt werden - sie stehen in Kriterientabellen.
EINHEIT_HARMLOS = {"80 %", "50 %", "100 %", "1,00 €", "0,50 €", "1 %", "20 %"}

MUSTER = re.compile(r"\d{1,3}(?:[.\s]\d{3})*(?:,\d{1,2})?")

# Zahlen MIT Einheit werden immer geprueft, auch kleine. Ohne diese
# Sonderbehandlung rutschen genau die Aussagen durch, auf die es ankommt:
# „36 %" und „14 %" sind Trefferquoten, „40" und „14" waeren als blosse
# Zahlen harmlos. So blieb eine Folie mit 40 % statt 36 % unbemerkt.
MIT_EINHEIT = re.compile(r"(\d{1,3}(?:[.\s]\d{3})*(?:,\d{1,2})?)\s?(%|€)")


def notizraum(stamm: str) -> str:
    """Ausgaben und Fliesstext eines Notebooks als ein Suchtext.

    Der QUELLTEXT bleibt bewusst draussen. Er enthaelt Zahlen, die mit
    dem Ergebnis nichts zu tun haben - Bildgroessen, Farbwerte, Startwerte
    fuer Zufallszahlen -, und die decken falsche Folienzahlen zu. Genau
    das ist passiert: Eine Folie nannte 9,5 % Trefferquote, das Notebook
    sagte 14 %, und der Befund blieb aus, weil irgendwo figsize=(9.5, 5.5)
    stand. Geprueft wird deshalb gegen das, was ein Leser tatsaechlich
    sieht: Ausgaben und Fliesstext.
    """
    nb = json.loads((NB / f"{stamm}.ipynb").read_text(encoding="utf-8"))
    teile = []
    for c in nb["cells"]:
        if c["cell_type"] == "markdown":
            teile.append("".join(c["source"]))
        for o in c.get("outputs", []):
            teile.append("".join(o.get("text", [])))
    return " ".join(teile)


def einheit_drin(zahl: str, einheit: str, raum: str) -> bool:
    """Prueft eine Angabe MIT ihrer Einheit.

    Die nackte Zahl genuegt hier ausdruecklich nicht: "14" steht in jedem
    Notebook irgendwo, "14 %" als Trefferquote aber nur, wenn es stimmt.
    Genau daran ist die erste Fassung dieses Werkzeugs gescheitert - eine
    Folie nannte 14 % Aufschlag, das Notebook rechnete mit 16 %, und der
    Befund blieb aus, weil "14" als Uhrzeit vorkam.

    Prozentwerte gelten auch als Bruchteil geschrieben als gefunden: Die
    Notebooks drucken teils "36.0%", teils "0.360".
    """
    try:
        wert = float(zahl.replace(".", "").replace(",", "."))
    except ValueError:
        return False
    formen = {f"{wert:g}{einheit}", f"{wert:g} {einheit}",
              f"{wert:.1f}{einheit}", f"{wert:.1f} {einheit}",
              f"{wert:.2f}{einheit}", f"{wert:.2f} {einheit}"}
    if einheit == "%":
        formen |= {f"{wert / 100:.3f}", f"{wert / 100:.2f}", f"{wert / 100:g}"}
    else:
        # Euro drucken die Notebooks oft ohne Zeichen: "12250.0" oder
        # "1.67 €". Beide Schreibweisen zaehlen.
        formen |= {f"{wert:.2f}", f"{wert:.1f}", f"{wert:g}",
                   f"{wert:,.1f}".replace(",", ".")}
    # Deutsche Schreibweise mitnehmen: Die Fliesstexte der Notebooks
    # schreiben "0,80 €", die Ausgaben drucken "0.80".
    formen |= {f.replace(".", ",") for f in formen}
    return any(f in raum for f in formen)


def steht_drin(zahl: str, raum: str) -> bool:
    """Deutsche und englische Schreibweise gelten als dieselbe Zahl.

    Die Notebooks drucken "60,425" und "3.96", die Folien schreiben
    "60.425" und "3,96". Ohne diese Umschrift meldet das Werkzeug jede
    zweite Zahl - und wird nach dem dritten Lauf ignoriert.
    """
    formen = {
        zahl,
        zahl.replace(",", "."),
        zahl.replace(".", ""),
        zahl.replace(".", ","),
        zahl.replace(".", "").replace(",", "."),
        zahl.replace(",", ".").rstrip("0").rstrip("."),
    }
    return any(f and f in raum for f in formen)


def main() -> int:
    ziel = pathlib.Path(sys.argv[1]) if len(sys.argv) > 1 else \
        WURZEL / "slides" / "velocity-crispdm.pptx"
    if not ziel.exists():
        raise SystemExit(f"Deck fehlt: {ziel}\n"
                         "Zuerst: python3 slides/build_crispdm_deck.py")
    raeume = {p.stem: notizraum(p.stem) for p in sorted(NB.glob("*.ipynb"))}
    prs = Presentation(str(ziel))
    geprueft = befunde = 0
    for nr, folie in enumerate(prs.slides, 1):
        quelle = next((sh.text_frame.text for sh in folie.shapes
                       if sh.has_text_frame and "analytics/notebooks/" in sh.text_frame.text),
                      None)
        if not quelle:
            continue
        treffer = re.search(r"notebooks/(\S+?)\.ipynb", quelle)
        if not treffer or treffer.group(1) not in raeume:
            continue
        raum = raeume[treffer.group(1)]
        geprueft += 1
        fremd = set()
        for sh in folie.shapes:
            if not sh.has_text_frame:
                continue
            text = sh.text_frame.text
            for zahl, einheit in MIT_EINHEIT.findall(text):
                if f"{zahl} {einheit}" in EINHEIT_HARMLOS:
                    continue
                if not einheit_drin(zahl, einheit, raum):
                    fremd.add(f"{zahl} {einheit}")
            for zahl in MUSTER.findall(text):
                if len(zahl) < 2 or zahl in HARMLOS:
                    continue
                if not steht_drin(zahl, raum):
                    fremd.add(zahl)
        if fremd:
            befunde += 1
            titel = next((sh.text_frame.text.split("\n")[0] for sh in folie.shapes
                          if sh.has_text_frame and sh.text_frame.text.strip()), "")
            print(f"Folie {nr:3}  {treffer.group(1)[:2]}  {titel[:56]}")
            print(f"           steht nicht im Notebook: {', '.join(sorted(fremd))}")
    print(f"\n{geprueft} Folien mit Quellenangabe geprueft, {befunde} Befund(e).")
    return 1 if befunde else 0


if __name__ == "__main__":
    raise SystemExit(main())
